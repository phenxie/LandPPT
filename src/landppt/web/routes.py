"""
Web interface routes for LandPPT
"""

from fastapi import APIRouter, Request, Form, UploadFile, File, HTTPException, Depends
from fastapi.responses import HTMLResponse, StreamingResponse, FileResponse, JSONResponse
from fastapi.templating import Jinja2Templates
from pydantic import BaseModel
import json
import uuid
import asyncio
import time
import os
import zipfile
import tempfile
import shutil
from pathlib import Path
from datetime import datetime
import urllib.parse
import subprocess
import logging
import time
from typing import Optional, Dict, Any, List

from ..api.models import PPTGenerationRequest, PPTProject, TodoBoard, FileOutlineGenerationRequest
from ..services.enhanced_ppt_service import EnhancedPPTService
from ..services.pdf_to_pptx_converter import get_pdf_to_pptx_converter
from ..services.pyppeteer_pdf_converter import get_pdf_converter
from ..core.config import ai_config
from ..ai import get_ai_provider, get_role_provider, AIMessage, MessageRole
from ..auth.middleware import get_current_user_required, get_current_user_optional
from ..database.models import User
from ..database.database import get_db
from sqlalchemy.orm import Session
from ..utils.thread_pool import run_blocking_io, to_thread
import re
from bs4 import BeautifulSoup

# Configure logger for this module
logger = logging.getLogger(__name__)

router = APIRouter()
templates = Jinja2Templates(directory="src/landppt/web/templates")

# Add custom filters
def timestamp_to_datetime(timestamp):
    """Convert timestamp to readable datetime string"""
    try:
        if isinstance(timestamp, (int, float)):
            return datetime.fromtimestamp(timestamp).strftime("%Y-%m-%d %H:%M:%S")
        return str(timestamp)
    except (ValueError, OSError):
        return "无效时间"

def strftime_filter(timestamp, format_string="%Y-%m-%d %H:%M"):
    """Jinja2 strftime filter"""
    try:
        if isinstance(timestamp, (int, float)):
            dt = datetime.fromtimestamp(timestamp)
            return dt.strftime(format_string)
        return str(timestamp)
    except (ValueError, OSError):
        return "无效时间"

# Register custom filters
templates.env.filters["timestamp_to_datetime"] = timestamp_to_datetime
templates.env.filters["strftime"] = strftime_filter

# Import shared service instances to ensure data consistency
from ..services.service_instances import ppt_service

# AI编辑请求数据模型
class AISlideEditRequest(BaseModel):
    slideIndex: int
    slideTitle: str
    slideContent: str
    userRequest: str
    projectInfo: Dict[str, Any]
    slideOutline: Optional[Dict[str, Any]] = None
    chatHistory: Optional[List[Dict[str, str]]] = None
    images: Optional[List[Dict[str, Any]]] = None  # 新增：图片信息列表（url/id/name/size 等）
    visionEnabled: Optional[bool] = False  # 新增：视觉模式启用状态
    slideScreenshot: Optional[str] = None  # 新增：幻灯片截图数据（data URL / base64）

# AI自由对话请求数据模型（不设系统提示词，仅当前页）
class AISlideNativeDialogRequest(BaseModel):
    slideIndex: int
    slideTitle: str
    slideContent: str
    userRequest: str
    chatHistory: Optional[List[Dict[str, str]]] = None
    images: Optional[List[Dict[str, str]]] = None  # 粘贴/上传图片信息列表（url/id/name/size）

# AI要点增强请求数据模型
class AIBulletPointEnhanceRequest(BaseModel):
    slideIndex: int
    slideTitle: str
    slideContent: str
    userRequest: str
    projectInfo: Dict[str, Any]
    slideOutline: Optional[Dict[str, Any]] = None
    contextInfo: Optional[Dict[str, Any]] = None  # 包含原始要点、其他要点等上下文信息

# 图像重新生成请求数据模型
class AIImageRegenerateRequest(BaseModel):
    slide_index: int
    image_info: Dict[str, Any]
    slide_content: Dict[str, Any]
    project_topic: str
    project_scenario: str
    regeneration_reason: Optional[str] = None

# 一键配图请求数据模型
class AIAutoImageGenerateRequest(BaseModel):
    slide_index: int
    slide_content: Dict[str, Any]
    project_topic: str
    project_scenario: str


class AutoLayoutRepairRequest(BaseModel):
    html_content: str
    slide_data: Dict[str, Any]


class SpeechScriptGenerationRequest(BaseModel):
    generation_type: str  # "single", "multi", "full"
    slide_indices: Optional[List[int]] = None  # For single and multi generation
    customization: Dict[str, Any] = {}  # Customization options

class SpeechScriptExportRequest(BaseModel):
    export_format: str  # "docx", "markdown"
    scripts_data: List[Dict[str, Any]]
    include_metadata: bool = True

# 图片导出PPTX请求数据模型
class ImagePPTXExportRequest(BaseModel):
    slides: Optional[List[Dict[str, Any]]] = None  # 包含index, html_content, title
    images: Optional[List[Dict[str, Any]]] = None  # 包含index, data(base64), width, height (向后兼容)

class SlideBatchRegenerateRequest(BaseModel):
    """Batch slide regeneration request (0-based indices)."""
    slide_indices: Optional[List[int]] = None
    regenerate_all: bool = False
    scenario: Optional[str] = None
    topic: Optional[str] = None
    requirements: Optional[str] = None
    language: str = "zh"

# Helper function to extract slides from HTML content
async def _extract_slides_from_html(slides_html: str, existing_slides_data: list) -> list:
    """
    Extract individual slides from combined HTML content and update slides_data
    """
    try:
        # Parse HTML content
        soup = BeautifulSoup(slides_html, 'html.parser')

        # Find all slide containers - look for common slide patterns
        slide_containers = []

        # Try different patterns to find slides
        patterns = [
            {'class': re.compile(r'slide')},
            {'class': re.compile(r'page')},
            {'style': re.compile(r'width:\s*1280px.*height:\s*720px', re.IGNORECASE)},
            {'style': re.compile(r'aspect-ratio:\s*16\s*/\s*9', re.IGNORECASE)}
        ]

        for pattern in patterns:
            containers = soup.find_all('div', pattern)
            if containers:
                slide_containers = containers
                break

        # If no specific slide containers found, try to split by common separators
        if not slide_containers:
            # Look for sections or divs that might represent slides
            all_divs = soup.find_all('div')
            # Filter divs that might be slides (have substantial content)
            slide_containers = [div for div in all_divs
                             if div.get_text(strip=True) and len(div.get_text(strip=True)) > 50]

        updated_slides_data = []

        # If we found slide containers, extract them
        if slide_containers:
            for i, container in enumerate(slide_containers):
                # Try to extract title from the slide
                title = f"第{i+1}页"
                title_elements = container.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])
                if title_elements:
                    title = title_elements[0].get_text(strip=True) or title

                # Get the HTML content of this slide
                slide_html = str(container)

                # Create slide data
                slide_data = {
                    "page_number": i + 1,
                    "title": title,
                    "html_content": slide_html,
                    "is_user_edited": True  # Mark as user edited since it came from editor
                }

                # If we have existing slide data, preserve some fields
                if i < len(existing_slides_data):
                    existing_slide = existing_slides_data[i]
                    # Preserve any additional fields from existing data
                    for key, value in existing_slide.items():
                        if key not in slide_data:
                            slide_data[key] = value

                updated_slides_data.append(slide_data)

        # If we couldn't extract individual slides, treat the entire content as slides
        if not updated_slides_data and existing_slides_data:
            # Fall back to using existing slides structure but mark as edited
            for i, existing_slide in enumerate(existing_slides_data):
                slide_data = existing_slide.copy()
                slide_data["is_user_edited"] = True
                updated_slides_data.append(slide_data)

        # If we still have no slides but have HTML content, create a single slide
        if not updated_slides_data and slides_html.strip():
            slide_data = {
                "page_number": 1,
                "title": "编辑后的PPT",
                "html_content": slides_html,
                "is_user_edited": True
            }
            updated_slides_data.append(slide_data)

        logger.info(f"Extracted {len(updated_slides_data)} slides from HTML content")
        return updated_slides_data

    except Exception as e:
        logger.error(f"Error extracting slides from HTML: {e}")
        # Fall back to marking existing slides as edited
        if existing_slides_data:
            updated_slides_data = []
            for slide in existing_slides_data:
                slide_copy = slide.copy()
                slide_copy["is_user_edited"] = True
                updated_slides_data.append(slide_copy)
            return updated_slides_data
        else:
            return []

@router.get("/home", response_class=HTMLResponse)
async def web_home(
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Main web interface home page - redirect to dashboard for existing users"""
    # Check if user has projects, if so redirect to dashboard
    try:
        projects_response = await ppt_service.project_manager.list_projects(page=1, page_size=1)
        if projects_response.total > 0:
            # User has projects, redirect to dashboard
            from fastapi.responses import RedirectResponse
            return RedirectResponse(url="/dashboard", status_code=302)
    except:
        pass  # If error, show index page

    # New user or error, show index page
    return templates.TemplateResponse("index.html", {
        "request": request,
        "ai_provider": ai_config.default_ai_provider,
        "available_providers": ai_config.get_available_providers()
    })

@router.get("/ai-config", response_class=HTMLResponse)
async def web_ai_config(
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """AI configuration page"""
    from ..services.config_service import get_config_service

    config_service = get_config_service()
    current_config = config_service.get_all_config()

    # "gemini" is an alias for the Google provider; the UI exposes it as "google".
    current_provider = ai_config.default_ai_provider
    if (isinstance(current_provider, str) and current_provider.strip().lower() == "gemini"):
        current_provider = "google"

    supported_providers = [
        "openai",
        "deepseek",
        "kimi",
        "minimax",
        "anthropic",
        "google",
        "ollama",
        "302ai",
    ]

    return templates.TemplateResponse("ai_config.html", {
        "request": request,
        "current_provider": current_provider,
        "available_providers": ai_config.get_available_providers(),
        "supported_providers": supported_providers,
        "provider_status": {provider: ai_config.is_provider_available(provider) for provider in supported_providers},
        "current_config": current_config,
        "user": user.to_dict()
    })


@router.get("/image-generation-test", response_class=HTMLResponse)
async def web_image_generation_test(
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """AI图片生成测试页面"""
    return templates.TemplateResponse("image_generation_test.html", {
        "request": request,
        "user": user.to_dict()
    })


def _is_truthy_config_value(value: Any) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return bool(value)
    if isinstance(value, str):
        return value.strip().lower() in {"true", "1", "yes", "on"}
    return bool(value)


def _extract_responses_output_text(response_data: Dict[str, Any]) -> str:
    output_text = response_data.get("output_text")
    if isinstance(output_text, str) and output_text:
        return output_text

    texts: List[str] = []
    for item in response_data.get("output", []) or []:
        if not isinstance(item, dict) or item.get("type") != "message":
            continue
        for content in item.get("content", []) or []:
            if isinstance(content, dict) and content.get("type") == "output_text" and content.get("text"):
                texts.append(content["text"])

    return "".join(texts)


def _extract_responses_usage(response_data: Dict[str, Any]) -> Dict[str, int]:
    usage = response_data.get("usage") or {}
    return {
        "prompt_tokens": int(usage.get("input_tokens") or 0),
        "completion_tokens": int(usage.get("output_tokens") or 0),
        "total_tokens": int(usage.get("total_tokens") or 0),
    }


@router.post("/api/ai/providers/openai/models")
async def get_openai_models(
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Proxy endpoint to get OpenAI models list, avoiding CORS issues - uses frontend provided config"""
    try:
        import aiohttp
        import json
        
        # Get configuration from frontend request
        data = await request.json()
        base_url = data.get('base_url', 'https://api.openai.com/v1')
        api_key = data.get('api_key', '')
        
        logger.info(f"Frontend requested models from: {base_url}")
        
        if not api_key:
            return {"success": False, "error": "API Key is required"}
        
        # Ensure base URL ends with /v1
        if not base_url.endswith('/v1'):
            base_url = base_url.rstrip('/') + '/v1'
        
        models_url = f"{base_url}/models"
        logger.info(f"Fetching models from: {models_url}")
        
        # Make request to OpenAI API using frontend provided credentials
        async with aiohttp.ClientSession() as session:
            headers = {
                'Authorization': f'Bearer {api_key}',
                'Content-Type': 'application/json'
            }
            
            async with session.get(models_url, headers=headers, timeout=30) as response:
                if response.status == 200:
                    data = await response.json()
                    
                    # Filter and sort models
                    models = []
                    if 'data' in data and isinstance(data['data'], list):
                        for model in data['data']:
                            if model.get('id'):
                                models.append({
                                    'id': model['id'],
                                    'created': model.get('created', 0),
                                    'owned_by': model.get('owned_by', 'unknown')
                                })
                        
                        # Sort models with GPT-4 first, then GPT-3.5, then others
                        def get_priority(model_id):
                            if 'gpt-4' in model_id:
                                return 0
                            elif 'gpt-3.5' in model_id:
                                return 1
                            else:
                                return 2
                        
                        models.sort(key=lambda x: (get_priority(x['id']), x['id']))
                    logger.info(f"Successfully fetched {len(models)} models from {base_url}")
                    return {"success": True, "models": models}
                else:
                    error_text = await response.text()
                    logger.error(f"Failed to fetch models from {base_url}: {response.status} - {error_text}")
                    return {"success": False, "error": f"API returned status {response.status}: {error_text}"}
                    
    except Exception as e:
        logger.error(f"Error fetching OpenAI models from frontend config: {e}")
        return {"success": False, "error": str(e)}

@router.post("/api/ai/providers/openai/test")
async def test_openai_provider_proxy(
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Proxy endpoint to test OpenAI provider, avoiding CORS issues - uses frontend provided config"""
    try:
        import aiohttp
        
        # Get configuration from frontend request
        data = await request.json()
        base_url = data.get('base_url', 'https://api.openai.com/v1')
        api_key = data.get('api_key', '')
        model = data.get('model', 'gpt-4o')
        use_responses_api = _is_truthy_config_value(data.get('use_responses_api', False))
        enable_reasoning = _is_truthy_config_value(data.get('enable_reasoning', False))
        reasoning_effort = str(data.get('reasoning_effort', 'medium') or 'medium').strip().lower()
        
        logger.info(
            f"Frontend requested test with: base_url={base_url}, model={model}, "
            f"use_responses_api={use_responses_api}, enable_reasoning={enable_reasoning}, "
            f"reasoning_effort={reasoning_effort}"
        )
        
        if not api_key:
            return {"success": False, "error": "API Key is required"}
        
        # Ensure base URL ends with /v1
        if not base_url.endswith('/v1'):
            base_url = base_url.rstrip('/') + '/v1'
        
        request_url = f"{base_url}/responses" if use_responses_api else f"{base_url}/chat/completions"
        logger.info(f"Testing OpenAI provider at: {request_url}")
        
        # Make test request to OpenAI API using frontend provided credentials
        async with aiohttp.ClientSession() as session:
            headers = {
                'Authorization': f'Bearer {api_key}',
                'Content-Type': 'application/json'
            }
            
            if use_responses_api:
                payload = {
                    "model": model,
                    "input": "Say 'Hello, I am working!' in exactly 5 words.",
                    "max_output_tokens": 32
                }
                if enable_reasoning:
                    payload["reasoning"] = {"effort": reasoning_effort}
            else:
                payload = {
                    "model": model,
                    "messages": [
                        {
                            "role": "user",
                            "content": "Say 'Hello, I am working!' in exactly 5 words."
                        }
                    ]
                }
                if enable_reasoning:
                    payload["reasoning_effort"] = reasoning_effort
            
            async with session.post(request_url, headers=headers, json=payload, timeout=30) as response:
                if response.status == 200:
                    data = await response.json()
                    response_preview = (
                        _extract_responses_output_text(data)
                        if use_responses_api
                        else data['choices'][0]['message']['content']
                    )
                    usage = (
                        _extract_responses_usage(data)
                        if use_responses_api
                        else data.get('usage', {
                            "prompt_tokens": 0,
                            "completion_tokens": 0,
                            "total_tokens": 0
                        })
                    )
                    
                    logger.info(f"Test successful for {base_url} with model {model}")
                    
                    # Return with consistent format that frontend expects
                    return {
                        "success": True,
                        "status": "success",  # Add status field for compatibility
                        "provider": "openai",
                        "model": model,
                        "api_mode": "responses" if use_responses_api else "chat_completions",
                        "response_preview": response_preview,
                        "usage": usage
                    }
                else:
                    error_text = await response.text()
                    try:
                        error_data = json.loads(error_text)
                        error_message = error_data.get('error', {}).get('message', f"API returned status {response.status}")
                    except:
                        error_message = f"API returned status {response.status}: {error_text}"
                    
                    logger.error(f"Test failed for {base_url}: {error_message}")
                    
                    return {
                        "success": False,
                        "status": "error",  # Add status field for compatibility
                        "error": error_message
                    }

    except Exception as e:
        logger.error(f"Error testing OpenAI provider with frontend config: {e}")
        return {
            "success": False,
            "status": "error",  # Add status field for compatibility
            "error": str(e)
        }

@router.post("/api/ai/providers/anthropic/test")
async def test_anthropic_provider_proxy(
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Proxy endpoint to test Anthropic provider, avoiding CORS issues - uses frontend provided config"""
    try:
        import aiohttp

        # Get configuration from frontend request
        data = await request.json()
        base_url = data.get('base_url', 'https://api.anthropic.com')
        api_key = data.get('api_key', '')
        model = data.get('model', 'claude-3-5-sonnet-20241022')

        logger.info(f"Frontend requested Anthropic test with: base_url={base_url}, model={model}")

        if not api_key:
            return {"success": False, "error": "API Key is required"}

        # Ensure base URL format
        base_url = base_url.rstrip('/')
        if not base_url.endswith('/v1'):
            base_url = base_url + '/v1'

        messages_url = f"{base_url}/messages"
        logger.info(f"Testing Anthropic provider at: {messages_url}")

        # Make test request to Anthropic API using frontend provided credentials
        async with aiohttp.ClientSession() as session:
            headers = {
                'x-api-key': api_key,
                'Content-Type': 'application/json',
                'anthropic-version': '2023-06-01'
            }

            payload = {
                "model": model,
                "messages": [
                    {
                        "role": "user",
                        "content": "Say 'Hello, I am working!' in exactly 5 words."
                    }
                ],
                "max_tokens": 1024,
            }

            async with session.post(messages_url, headers=headers, json=payload, timeout=30) as response:
                if response.status == 200:
                    data = await response.json()

                    logger.info(f"Anthropic test successful for {base_url} with model {model}")

                    # Anthropic response format: data.content[0].text
                    content = data.get('content', [])
                    response_text = content[0].get('text', '') if content else ''

                    # Anthropic usage format: input_tokens, output_tokens
                    # Frontend expects: prompt_tokens, completion_tokens, total_tokens
                    anthropic_usage = data.get('usage', {})
                    usage = {
                        "prompt_tokens": anthropic_usage.get('input_tokens', 0),
                        "completion_tokens": anthropic_usage.get('output_tokens', 0),
                        "total_tokens": anthropic_usage.get('input_tokens', 0) + anthropic_usage.get('output_tokens', 0)
                    }

                    # Return with consistent format that frontend expects
                    return {
                        "success": True,
                        "status": "success",
                        "provider": "anthropic",
                        "model": model,
                        "response_preview": response_text,
                        "usage": usage
                    }
                else:
                    error_text = await response.text()
                    try:
                        error_data = json.loads(error_text)
                        error_message = error_data.get('error', {}).get('message', f"API returned status {response.status}")
                    except:
                        error_message = f"API returned status {response.status}: {error_text}"

                    logger.error(f"Anthropic test failed for {base_url}: {error_message}")

                    return {
                        "success": False,
                        "status": "error",
                        "error": error_message
                    }

    except Exception as e:
        logger.error(f"Error testing Anthropic provider with frontend config: {e}")
        return {
            "success": False,
            "status": "error",
            "error": str(e)
        }

@router.get("/scenarios", response_class=HTMLResponse)
async def web_scenarios(
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Scenarios selection page"""
    scenarios = [
        {"id": "general", "name": "通用", "description": "适用于各种通用场景的PPT模板", "icon": "📋"},
        {"id": "tourism", "name": "旅游观光", "description": "旅游线路、景点介绍等旅游相关PPT", "icon": "🌍"},
        {"id": "education", "name": "儿童科普", "description": "适合儿童的科普教育PPT", "icon": "🎓"},
        {"id": "analysis", "name": "深入分析", "description": "数据分析、研究报告等深度分析PPT", "icon": "📊"},
        {"id": "history", "name": "历史文化", "description": "历史事件、文化介绍等人文类PPT", "icon": "🏛️"},
        {"id": "technology", "name": "科技技术", "description": "技术介绍、产品发布等科技类PPT", "icon": "💻"},
        {"id": "business", "name": "方案汇报", "description": "商业计划、项目汇报等商务PPT", "icon": "💼"}
    ]
    return templates.TemplateResponse("scenarios.html", {"request": request, "scenarios": scenarios})

# Legacy route removed - now using /projects/create for new project workflow

# Legacy task status route removed - now using project detail pages

# Legacy preview route removed - now using project-based preview at /projects/{project_id}/fullscreen

# Legacy tasks list route removed - now using /projects for project management

@router.post("/upload", response_class=HTMLResponse)
async def web_upload_file(
    request: Request,
    file: UploadFile = File(...),
    user: User = Depends(get_current_user_required)
):
    """Upload file via web interface"""
    try:
        # Validate file type
        allowed_types = [".docx", ".pdf", ".txt", ".md"]
        file_extension = "." + file.filename.split(".")[-1].lower()

        if file_extension not in allowed_types:
            return templates.TemplateResponse("upload_result.html", {
                "request": request,
                "success": False,
                "error": f"Unsupported file type. Allowed types: {', '.join(allowed_types)}"
            })

        # Read file content in thread pool to avoid blocking
        content = await file.read()

        # Process file in thread pool
        processed_content = await ppt_service.process_uploaded_file(
            filename=file.filename,
            content=content,
            file_type=file_extension
        )

        return templates.TemplateResponse("upload_result.html", {
            "request": request,
            "success": True,
            "filename": file.filename,
            "size": len(content),
            "type": file_extension,
            "processed_content": processed_content[:500] + "..." if len(processed_content) > 500 else processed_content
        })

    except Exception as e:
        return templates.TemplateResponse("upload_result.html", {
            "request": request,
            "success": False,
            "error": str(e)
        })

@router.get("/demo", response_class=HTMLResponse)
async def web_demo(
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Demo page with sample PPT"""
    # Create a demo PPT
    demo_request = PPTGenerationRequest(
        scenario="technology",
        topic="人工智能技术发展趋势",
        requirements="面向技术人员的深度分析",
        network_mode=False,
        language="zh"
    )

    task_id = "demo-" + str(uuid.uuid4())[:8]
    result = await ppt_service.generate_ppt(task_id, demo_request)

    return templates.TemplateResponse("demo.html", {
        "request": request,
        "task_id": task_id,
        "outline": result.get("outline"),
        "slides_html": result.get("slides_html"),
        "demo_topic": demo_request.topic
    })

@router.get("/research", response_class=HTMLResponse)
async def web_research_status(
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """DEEP Research status and management page"""
    return templates.TemplateResponse("research_status.html", {
        "request": request
    })

# New Project Management Routes

@router.get("/dashboard", response_class=HTMLResponse)
async def web_dashboard(
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Project dashboard with overview"""
    try:
        # Get project statistics
        projects_response = await ppt_service.project_manager.list_projects(page=1, page_size=100)
        projects = projects_response.projects

        total_projects = len(projects)
        completed_projects = len([p for p in projects if p.status == "completed"])
        in_progress_projects = len([p for p in projects if p.status == "in_progress"])
        draft_projects = len([p for p in projects if p.status == "draft"])

        # Get recent projects (last 5)
        recent_projects = sorted(projects, key=lambda x: x.updated_at, reverse=True)[:5]

        # Get active TODO boards
        active_todo_boards = []
        for project in projects:
            if project.status == "in_progress" and project.todo_board:
                todo_board = await ppt_service.get_project_todo_board(project.project_id)
                if todo_board:
                    active_todo_boards.append(todo_board)

        return templates.TemplateResponse("project_dashboard.html", {
            "request": request,
            "total_projects": total_projects,
            "completed_projects": completed_projects,
            "in_progress_projects": in_progress_projects,
            "draft_projects": draft_projects,
            "recent_projects": recent_projects,
            "active_todo_boards": active_todo_boards[:3]  # Show max 3 boards
        })

    except Exception as e:
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })

@router.get("/projects", response_class=HTMLResponse)
async def web_projects_list(
    request: Request,
    page: int = 1,
    status: str = None,
    user: User = Depends(get_current_user_required)
):
    """List all projects"""
    try:
        projects_response = await ppt_service.project_manager.list_projects(
            page=page, page_size=10, status=status
        )

        return templates.TemplateResponse("projects_list.html", {
            "request": request,
            "projects": projects_response.projects,
            "total": projects_response.total,
            "page": projects_response.page,
            "page_size": projects_response.page_size,
            "status_filter": status
        })

    except Exception as e:
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })

@router.get("/projects/{project_id}", response_class=HTMLResponse)
async def web_project_detail(
    request: Request,
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """Project detail page"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return templates.TemplateResponse("error.html", {
                "request": request,
                "error": "Project not found"
            })

        todo_board = await ppt_service.get_project_todo_board(project_id)
        versions = await ppt_service.project_manager.get_project_versions(project_id)

        return templates.TemplateResponse("project_detail.html", {
            "request": request,
            "project": project,
            "todo_board": todo_board,
            "versions": versions
        })

    except Exception as e:
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })

@router.get("/projects/{project_id}/todo", response_class=HTMLResponse)
async def web_project_todo_board(
    request: Request,
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """TODO board page for a project with integrated editor"""
    try:
        # Validate project_id format (should be UUID-like)
        if project_id in ["template-selection", "todo", "edit", "preview", "fullscreen"]:
            error_msg = f"无效的项目ID: {project_id}。\n\n"
            error_msg += "可能的原因：\n"
            error_msg += "1. URL格式错误，正确格式应为: /projects/[项目ID]/todo\n"
            error_msg += "2. 您可能访问了错误的链接\n\n"
            error_msg += "建议解决方案：\n"
            error_msg += "• 返回项目列表页面选择正确的项目\n"
            error_msg += "• 检查浏览器地址栏中的URL是否完整"

            return templates.TemplateResponse("error.html", {
                "request": request,
                "error": error_msg
            })

        # Check if project exists first
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return templates.TemplateResponse("error.html", {
                "request": request,
                "error": f"项目不存在 (ID: {project_id})。请检查项目ID是否正确。"
            })

        todo_board = await ppt_service.get_project_todo_board(project_id)
        if not todo_board:
            return templates.TemplateResponse("error.html", {
                "request": request,
                "error": f"项目 '{project.topic}' 的TODO看板不存在。请联系技术支持。"
            })

        # Check if we should use the integrated editor version
        project = await ppt_service.project_manager.get_project(project_id)
        use_integrated_editor = (
            project and
            project.confirmed_requirements and
            len(todo_board.stages) > 2 and
            (todo_board.stages[1].status in ['running', 'completed'] or
             todo_board.stages[2].status in ['running', 'completed'])
        )

        # Also use integrated editor if PPT creation stage is about to start or running
        if (project and project.confirmed_requirements and len(todo_board.stages) > 2 and
            todo_board.stages[1].status == 'completed'):
            use_integrated_editor = True

        template_name = "todo_board_with_editor.html" if use_integrated_editor else "todo_board.html"

        # Ensure project is not None for template
        template_context = {
            "request": request,
            "todo_board": todo_board
        }

        # Only add project if it exists
        if project:
            template_context["project"] = project

        return templates.TemplateResponse(template_name, template_context)

    except Exception as e:
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })



@router.get("/projects/{project_id}/fullscreen", response_class=HTMLResponse)
async def web_project_fullscreen(
    request: Request,
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """Fullscreen preview of project PPT with modern presentation interface"""
    try:
        # 直接从数据库获取最新的项目数据，确保数据实时性
        from ..services.db_project_manager import DatabaseProjectManager
        db_manager = DatabaseProjectManager()
        project = await db_manager.get_project(project_id)

        if not project:
            return templates.TemplateResponse("error.html", {
                "request": request,
                "error": "项目未找到"
            })

        # 检查是否有幻灯片数据
        if not project.slides_data or len(project.slides_data) == 0:
            return templates.TemplateResponse("error.html", {
                "request": request,
                "error": "PPT尚未生成或无幻灯片内容"
            })

        # 使用新的分享演示模板
        return templates.TemplateResponse("project_fullscreen_presentation.html", {
            "request": request,
            "project": project,
            "slides_count": len(project.slides_data)
        })

    except Exception as e:
        logger.error(f"Error in fullscreen presentation: {e}")
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": f"加载演示时出错: {str(e)}"
        })

@router.get("/share/{share_token}", response_class=HTMLResponse)
async def web_shared_presentation(
    request: Request,
    share_token: str,
    db: Session = Depends(get_db)
):
    """Public presentation view - no authentication required"""
    try:
        from ..services.share_service import ShareService
        share_service = ShareService(db)

        # Validate share token and get project
        project_model = share_service.validate_share_token(share_token)

        if not project_model:
            return templates.TemplateResponse("error.html", {
                "request": request,
                "error": "分享链接无效或已失效"
            })

        # Check if project has slides
        if not project_model.slides_data or len(project_model.slides_data) == 0:
            return templates.TemplateResponse("error.html", {
                "request": request,
                "error": "演示文稿尚未生成"
            })

        # Convert to PPTProject for template compatibility
        from ..api.models import PPTProject
        project = PPTProject(
            project_id=project_model.project_id,
            title=project_model.title,
            scenario=project_model.scenario,
            topic=project_model.topic,
            requirements=project_model.requirements,
            status=project_model.status,
            outline=project_model.outline,
            slides_html=project_model.slides_html,
            slides_data=project_model.slides_data,
            confirmed_requirements=project_model.confirmed_requirements,
            version=project_model.version,
            created_at=project_model.created_at,
            updated_at=project_model.updated_at
        )

        # Render presentation template
        return templates.TemplateResponse("project_fullscreen_presentation.html", {
            "request": request,
            "project": project,
            "slides_count": len(project.slides_data),
            "is_shared": True  # Flag to indicate this is a shared view
        })

    except Exception as e:
        logger.error(f"Error displaying shared presentation: {e}")
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": f"加载分享演示时出错: {str(e)}"
        })


@router.get("/api/share/{share_token}/slides-data")
async def get_shared_slides_data(
    share_token: str,
    db: Session = Depends(get_db)
):
    """Get slides data for public shared presentation - no authentication required"""
    try:
        from ..services.share_service import ShareService
        share_service = ShareService(db)

        # Validate share token and get project
        project = share_service.validate_share_token(share_token)

        if not project:
            raise HTTPException(status_code=404, detail="分享链接无效或已失效")

        if not project.slides_data or len(project.slides_data) == 0:
            return {
                "status": "no_slides",
                "message": "PPT尚未生成",
                "slides_data": [],
                "total_slides": 0
            }

        return {
            "status": "success",
            "slides_data": project.slides_data,
            "total_slides": len(project.slides_data),
            "project_title": project.title,
            "updated_at": project.updated_at
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error getting shared slides data: {e}")
        raise HTTPException(status_code=500, detail=f"获取幻灯片数据失败: {str(e)}")


@router.get("/api/projects/{project_id}/slides-data")
async def get_project_slides_data(
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """获取项目最新的幻灯片数据 - 用于分享演示实时更新"""
    try:
        # 直接从数据库获取最新数据
        from ..services.db_project_manager import DatabaseProjectManager
        db_manager = DatabaseProjectManager()
        project = await db_manager.get_project(project_id)

        if not project:
            raise HTTPException(status_code=404, detail="项目未找到")

        if not project.slides_data or len(project.slides_data) == 0:
            return {
                "status": "no_slides",
                "message": "PPT尚未生成",
                "slides_data": [],
                "total_slides": 0
            }

        return {
            "status": "success",
            "slides_data": project.slides_data,
            "total_slides": len(project.slides_data),
            "project_title": project.title,
            "updated_at": project.updated_at
        }

    except Exception as e:
        logger.error(f"Error getting slides data: {e}")
        raise HTTPException(status_code=500, detail=f"获取幻灯片数据失败: {str(e)}")


@router.post("/api/projects/{project_id}/share/generate")
async def generate_share_link(
    project_id: str,
    user: User = Depends(get_current_user_required),
    db: Session = Depends(get_db)
):
    """Generate a public share link for a project"""
    try:
        from ..services.share_service import ShareService
        share_service = ShareService(db)

        # Verify project exists and belongs to user
        from ..services.db_project_manager import DatabaseProjectManager
        db_manager = DatabaseProjectManager()
        project = await db_manager.get_project(project_id)

        if not project:
            raise HTTPException(status_code=404, detail="项目未找到")

        # Generate share token
        share_token = share_service.generate_share_token(project_id)

        if not share_token:
            raise HTTPException(status_code=500, detail="生成分享链接失败")

        # Construct full share URL
        share_url = f"/share/{share_token}"

        return {
            "success": True,
            "share_token": share_token,
            "share_url": share_url,
            "message": "分享链接已生成"
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error generating share link: {e}")
        raise HTTPException(status_code=500, detail=f"生成分享链接失败: {str(e)}")


@router.post("/api/projects/{project_id}/share/disable")
async def disable_share_link(
    project_id: str,
    user: User = Depends(get_current_user_required),
    db: Session = Depends(get_db)
):
    """Disable sharing for a project"""
    try:
        from ..services.share_service import ShareService
        share_service = ShareService(db)

        # Verify project exists
        from ..services.db_project_manager import DatabaseProjectManager
        db_manager = DatabaseProjectManager()
        project = await db_manager.get_project(project_id)

        if not project:
            raise HTTPException(status_code=404, detail="项目未找到")

        # Disable sharing
        success = share_service.disable_sharing(project_id)

        if not success:
            raise HTTPException(status_code=500, detail="禁用分享失败")

        return {
            "success": True,
            "message": "分享已禁用"
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error disabling share: {e}")
        raise HTTPException(status_code=500, detail=f"禁用分享失败: {str(e)}")


@router.get("/api/projects/{project_id}/share/info")
async def get_share_info(
    project_id: str,
    user: User = Depends(get_current_user_required),
    db: Session = Depends(get_db)
):
    """Get share information for a project"""
    try:
        from ..services.share_service import ShareService
        share_service = ShareService(db)

        # Verify project exists
        from ..services.db_project_manager import DatabaseProjectManager
        db_manager = DatabaseProjectManager()
        project = await db_manager.get_project(project_id)

        if not project:
            raise HTTPException(status_code=404, detail="项目未找到")

        # Get share info
        share_info = share_service.get_share_info(project_id)

        return {
            "success": True,
            **share_info
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error getting share info: {e}")
        raise HTTPException(status_code=500, detail=f"获取分享信息失败: {str(e)}")


@router.get("/test/slides-navigation", response_class=HTMLResponse)
async def test_slides_navigation(
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """测试幻灯片导航功能"""
    with open("test_slides_navigation.html", "r", encoding="utf-8") as f:
        content = f.read()
    return HTMLResponse(content=content)

@router.get("/temp/{file_path:path}")
async def serve_temp_file(
    file_path: str,
    user: User = Depends(get_current_user_required)
):
    """Serve temporary slide files"""
    try:
        # Construct the full path to the temp file using system temp directory
        import tempfile
        temp_dir = Path(tempfile.gettempdir()) / "landppt"
        full_path = temp_dir / file_path

        # Security check: ensure the file is within the temp directory
        if not str(full_path.resolve()).startswith(str(temp_dir.resolve())):
            raise HTTPException(status_code=403, detail="Access denied")

        # Check if file exists
        if not full_path.exists():
            raise HTTPException(status_code=404, detail="File not found")

        # Return the file
        return FileResponse(
            path=str(full_path),
            media_type="text/html; charset=utf-8",
            headers={"Cache-Control": "no-cache"}
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/projects/create", response_class=HTMLResponse)
async def web_create_project(
    request: Request,
    scenario: str = Form(...),
    topic: str = Form(...),
    requirements: str = Form(None),
    language: str = Form("zh"),
    network_mode: bool = Form(False),
    user: User = Depends(get_current_user_required)
):
    """Create new project via web interface"""
    try:
        # Create project request
        project_request = PPTGenerationRequest(
            scenario=scenario,
            topic=topic,
            requirements=requirements,
            network_mode=network_mode,
            language=language
        )

        # Create project with TODO board (without starting workflow yet)
        project = await ppt_service.project_manager.create_project(project_request)

        # Update project status to in_progress
        await ppt_service.project_manager.update_project_status(project.project_id, "in_progress")

        # Redirect directly to TODO page without showing redirect page
        from fastapi.responses import RedirectResponse
        return RedirectResponse(
            url=f"/projects/{project.project_id}/todo",
            status_code=302
        )

    except Exception as e:
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })

@router.post("/projects/{project_id}/start-workflow")
async def start_project_workflow(
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """Start the AI workflow for a project (only if requirements are confirmed)"""
    try:
        # Get project
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if requirements are confirmed
        if not project.confirmed_requirements:
            return {"status": "waiting", "message": "Waiting for requirements confirmation"}

        # Extract network_mode from project metadata
        network_mode = False
        if project.project_metadata and isinstance(project.project_metadata, dict):
            network_mode = project.project_metadata.get("network_mode", False)

        # Extract language from project metadata (set during project creation)
        language = "zh"  # Default language
        if project.project_metadata and isinstance(project.project_metadata, dict):
            language = project.project_metadata.get("language", "zh")

        # Create project request from project data
        confirmed_requirements = project.confirmed_requirements or {}
        project_request = PPTGenerationRequest(
            scenario=project.scenario,
            topic=project.topic,
            requirements=project.requirements,
            language=language,
            network_mode=network_mode,
            target_audience=confirmed_requirements.get('target_audience', '普通大众'),
            ppt_style=confirmed_requirements.get('ppt_style', 'general'),
            custom_style_prompt=confirmed_requirements.get('custom_style_prompt'),
            description=confirmed_requirements.get('description')
        )

        # Start the workflow in background
        asyncio.create_task(ppt_service._execute_project_workflow(project_id, project_request))

        return {"status": "success", "message": "Workflow started"}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/projects/{project_id}/requirements", response_class=HTMLResponse)
async def project_requirements_page(
    request: Request,
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """Show project requirements confirmation page"""
    try:
        # Get project
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # 提供默认的展示类型选项，不再调用AI生成建议
        default_type_options = [
            "技术分享",
            "产品介绍",
            "学术报告",
            "商业汇报",
            "教学课件",
            "项目展示",
            "数据分析",
            "综合介绍"
        ]

        return templates.TemplateResponse("project_requirements.html", {
            "request": request,
            "project": project,
            "ai_suggestions": {
                "type_options": default_type_options
            }
        })

    except Exception as e:
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })

# 移除AI生成需求建议的API端点，改为使用默认选项

@router.get("/projects/{project_id}/outline-stream")
async def stream_outline_generation(
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """Stream outline generation for a project"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        async def generate():
            try:
                async for chunk in ppt_service.generate_outline_streaming(project_id):
                    yield chunk
            except Exception as e:
                import json
                error_response = {'error': str(e)}
                yield f"data: {json.dumps(error_response)}\n\n"

        return StreamingResponse(generate(), media_type="text/plain")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/projects/{project_id}/generate-outline")
async def generate_outline(
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """Generate outline for a project (non-streaming)"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if project has confirmed requirements
        if not project.confirmed_requirements:
            return {
                "status": "error",
                "error": "项目需求尚未确认，请先完成需求确认步骤"
            }

        # Create PPTGenerationRequest from project data
        confirmed_requirements = project.confirmed_requirements

        # Extract network_mode and language from project metadata
        network_mode = False
        language = "zh"  # Default language
        if project.project_metadata and isinstance(project.project_metadata, dict):
            network_mode = project.project_metadata.get("network_mode", False)
            language = project.project_metadata.get("language", "zh")

        project_request = PPTGenerationRequest(
            scenario=project.scenario,
            topic=confirmed_requirements.get('topic', project.topic),
            requirements=project.requirements,
            language=language,
            network_mode=network_mode,
            target_audience=confirmed_requirements.get('target_audience', '普通大众'),
            ppt_style=confirmed_requirements.get('ppt_style', 'general'),
            custom_style_prompt=confirmed_requirements.get('custom_style_prompt'),
            description=confirmed_requirements.get('description')
        )

        # Extract page count settings from confirmed requirements
        page_count_settings = confirmed_requirements.get('page_count_settings', {})

        # Generate outline using AI with page count settings
        outline = await ppt_service.generate_outline(project_request, page_count_settings)

        # Convert outline to dict format
        outline_dict = {
            "title": outline.title,
            "slides": outline.slides,
            "metadata": outline.metadata
        }

        # Format as JSON
        import json
        formatted_json = json.dumps(outline_dict, ensure_ascii=False, indent=2)

        # Update outline generation stage
        await ppt_service._update_outline_generation_stage(project_id, outline_dict)

        return {
            "status": "success",
            "outline_content": formatted_json,
            "message": "Outline generated successfully"
        }

    except Exception as e:
        logger.error(f"Error generating outline: {e}")
        return {
            "status": "error",
            "error": str(e)
        }

@router.post("/projects/{project_id}/regenerate-outline")
async def regenerate_outline(
    project_id: str,
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Regenerate outline for a project (overwrites existing outline) with optional custom requirements"""
    try:
        # Get request body to extract custom requirements if provided
        request_data = {}
        try:
            request_data = await request.json()
        except:
            pass  # If no body or invalid JSON, use empty dict
        
        custom_requirements = (request_data.get('custom_requirements') or '').strip()
        
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if project has confirmed requirements
        if not project.confirmed_requirements:
            return {
                "status": "error",
                "error": "项目需求尚未确认，请先完成需求确认步骤"
            }

        # Create project request from confirmed requirements
        confirmed_requirements = project.confirmed_requirements
        
        # Extract language from project metadata (set during project creation)
        language = "zh"  # Default language
        if project.project_metadata and isinstance(project.project_metadata, dict):
            language = project.project_metadata.get("language", "zh")
        
        # 如果提供了自定义需求，将其追加或覆盖原有需求
        final_requirements = confirmed_requirements.get('requirements', project.requirements) or ''
        if custom_requirements:
            # 将自定义需求追加到原有需求
            if final_requirements:
                final_requirements = f"{final_requirements}\n\n【本次重新生成的额外要求】\n{custom_requirements}"
            else:
                final_requirements = custom_requirements
        
        project_request = PPTGenerationRequest(
            scenario=confirmed_requirements.get('scenario', 'general'),
            topic=confirmed_requirements.get('topic', project.topic),
            requirements=final_requirements,
            language=language,
            network_mode=confirmed_requirements.get('network_mode', False),
            target_audience=confirmed_requirements.get('target_audience', '普通大众'),
            ppt_style=confirmed_requirements.get('ppt_style', 'general'),
            custom_style_prompt=confirmed_requirements.get('custom_style_prompt'),
            description=confirmed_requirements.get('description')
        )

        # Extract page count settings from confirmed requirements
        page_count_settings = confirmed_requirements.get('page_count_settings', {})

        # Check if this is a file-based project
        is_file_project = confirmed_requirements.get('content_source') == 'file'

        if is_file_project:
            # Check if file path exists
            file_path = confirmed_requirements.get('file_path')
            if not file_path:
                return {
                    "status": "error",
                    "error": "文件路径信息丢失，请重新上传文件并确认需求"
                }

            # Use file-based outline generation
            file_request = FileOutlineGenerationRequest(
                file_path=file_path,
                filename=confirmed_requirements.get('filename', 'uploaded_file'),
                topic=project_request.topic,
                scenario=project_request.scenario,
                requirements=final_requirements,
                target_audience=confirmed_requirements.get('target_audience', '普通大众'),
                language=language,
                page_count_mode=page_count_settings.get('mode', 'ai_decide'),
                min_pages=page_count_settings.get('min_pages', 5),
                max_pages=page_count_settings.get('max_pages', 20),
                fixed_pages=page_count_settings.get('fixed_pages', 10),
                ppt_style=confirmed_requirements.get('ppt_style', 'general'),
                custom_style_prompt=confirmed_requirements.get('custom_style_prompt'),
                file_processing_mode=confirmed_requirements.get('file_processing_mode', 'markitdown'),
                content_analysis_depth=confirmed_requirements.get('content_analysis_depth', 'standard')
            )

            result = await ppt_service.generate_outline_from_file(file_request)

            if not result.success:
                return {
                    "status": "error",
                    "error": result.error or "文件大纲生成失败"
                }

            # Update outline generation stage
            await ppt_service._update_outline_generation_stage(project_id, result.outline)

            # Format outline as JSON string
            import json
            outline_content = json.dumps(result.outline, ensure_ascii=False, indent=2)

            return {
                "status": "success",
                "outline_content": outline_content,
                "message": "File-based outline regenerated successfully"
            }
        else:
            # Use standard outline generation
            outline = await ppt_service.generate_outline(project_request, page_count_settings)

            # Convert outline to dict format
            outline_dict = {
                "title": outline.title,
                "slides": outline.slides,
                "metadata": outline.metadata
            }

            # Format as JSON
            import json
            formatted_json = json.dumps(outline_dict, ensure_ascii=False, indent=2)

            # Update outline generation stage
            await ppt_service._update_outline_generation_stage(project_id, outline_dict)

            return {
                "status": "success",
                "outline_content": formatted_json,
                "message": "Outline regenerated successfully"
            }

    except Exception as e:
        logger.error(f"Error regenerating outline: {e}")
        return {
            "status": "error",
            "error": str(e)
        }

@router.post("/projects/{project_id}/generate-file-outline")
async def generate_file_outline(
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """Generate outline from uploaded file (non-streaming)"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if project has file-generated outline
        file_generated_outline = None

        # 首先检查项目的outline字段
        if project.outline and project.outline.get('slides'):
            # 检查是否是从文件生成的大纲
            metadata = project.outline.get('metadata', {})
            if metadata.get('generated_with_summeryfile') or metadata.get('generated_with_file'):
                file_generated_outline = project.outline
                logger.info(f"Project {project_id} has file-generated outline in project.outline, using it")

        # 如果项目outline中没有，再检查confirmed_requirements
        if not file_generated_outline and project.confirmed_requirements and project.confirmed_requirements.get('file_generated_outline'):
            file_generated_outline = project.confirmed_requirements['file_generated_outline']
            logger.info(f"Project {project_id} has file-generated outline in confirmed_requirements, using it")

        # If no existing outline but file upload is configured, wait a bit and check again
        if not file_generated_outline and project.confirmed_requirements and project.confirmed_requirements.get('content_source') == 'file':
            logger.info(f"Project {project_id} has file upload but no outline yet, waiting for file processing...")

            # Wait for file processing to complete (it should be done during requirements confirmation)
            import asyncio
            max_wait_time = 10  # Maximum wait time in seconds
            wait_interval = 1   # Check every 1 second

            for i in range(max_wait_time):
                await asyncio.sleep(wait_interval)

                # Refresh project data
                project = await ppt_service.project_manager.get_project(project_id)
                if project.confirmed_requirements and project.confirmed_requirements.get('file_generated_outline'):
                    file_generated_outline = project.confirmed_requirements['file_generated_outline']
                    logger.info(f"Project {project_id} file outline found after waiting {i+1} seconds")
                    break

            if not file_generated_outline:
                logger.warning(f"Project {project_id} file outline not found after waiting {max_wait_time} seconds")

        if file_generated_outline:
            # Return the existing file-generated outline
            import json
            existing_outline = {
                "title": file_generated_outline.get('title', project.topic),
                "slides": file_generated_outline.get('slides', []),
                "metadata": file_generated_outline.get('metadata', {})
            }

            # Ensure metadata includes correct identification
            if 'metadata' not in existing_outline:
                existing_outline['metadata'] = {}
            existing_outline['metadata']['generated_with_summeryfile'] = True
            existing_outline['metadata']['generated_at'] = time.time()

            formatted_json = json.dumps(existing_outline, ensure_ascii=False, indent=2)

            # Update outline generation stage
            await ppt_service._update_outline_generation_stage(project_id, existing_outline)

            return {
                "status": "success",
                "outline_content": formatted_json,
                "message": "File outline generated successfully"
            }
        else:
            # Check if there's an uploaded file that needs processing
            if (project.confirmed_requirements and
                (project.confirmed_requirements.get('uploaded_files') or
                 project.confirmed_requirements.get('content_source') == 'file')):
                logger.info(f"Project {project_id} has uploaded files, starting file outline generation")

                # Start file outline generation using summeryfile
                try:
                    # Create a request object for file outline generation
                    from ..api.models import FileOutlineGenerationRequest

                    # Get file information from confirmed requirements
                    uploaded_files = project.confirmed_requirements.get('uploaded_files', [])
                    if uploaded_files:
                        file_info = uploaded_files[0]  # Use first file
                        # 使用确认的要求或项目创建时的要求作为fallback
                        confirmed_reqs = project.confirmed_requirements.get('requirements', '')
                        project_reqs = project.requirements or ''
                        final_reqs = confirmed_reqs or project_reqs

                        # Extract language from project metadata (set during project creation)
                        language = "zh"
                        if project.project_metadata and isinstance(project.project_metadata, dict):
                            language = project.project_metadata.get("language", "zh")

                        file_request = FileOutlineGenerationRequest(
                            filename=file_info.get('filename', 'uploaded_file'),
                            file_path=file_info.get('file_path', ''),
                            topic=project.topic,
                            scenario='general',
                            requirements=final_reqs,
                            target_audience=project.confirmed_requirements.get('target_audience', '普通大众'),
                            language=language,
                            page_count_mode=project.confirmed_requirements.get('page_count_settings', {}).get('mode', 'ai_decide'),
                            min_pages=project.confirmed_requirements.get('page_count_settings', {}).get('min_pages', 8),
                            max_pages=project.confirmed_requirements.get('page_count_settings', {}).get('max_pages', 15),
                            fixed_pages=project.confirmed_requirements.get('page_count_settings', {}).get('fixed_pages', 10),
                            ppt_style=project.confirmed_requirements.get('ppt_style', 'general'),
                            custom_style_prompt=project.confirmed_requirements.get('custom_style_prompt'),
                            file_processing_mode=project.confirmed_requirements.get('file_processing_mode', 'markitdown'),
                            content_analysis_depth=project.confirmed_requirements.get('content_analysis_depth', 'standard')
                        )

                        # Generate outline from file using summeryfile
                        outline_response = await ppt_service.generate_outline_from_file(file_request)

                        if outline_response.success and outline_response.outline:
                            # Format the generated outline
                            import json
                            formatted_outline = outline_response.outline

                            # Ensure metadata includes correct identification
                            if 'metadata' not in formatted_outline:
                                formatted_outline['metadata'] = {}
                            formatted_outline['metadata']['generated_with_summeryfile'] = True
                            formatted_outline['metadata']['generated_at'] = time.time()

                            formatted_json = json.dumps(formatted_outline, ensure_ascii=False, indent=2)

                            # Update outline generation stage
                            await ppt_service._update_outline_generation_stage(project_id, formatted_outline)

                            return {
                                "status": "success",
                                "outline_content": formatted_json,
                                "message": "File outline generated successfully"
                            }
                        else:
                            error_msg = outline_response.error if hasattr(outline_response, 'error') else "Unknown error"
                            return {
                                "status": "error",
                                "error": f"Failed to generate outline from uploaded file: {error_msg}"
                            }
                    else:
                        return {
                            "status": "error",
                            "error": "No uploaded file information found in project requirements."
                        }

                except Exception as gen_error:
                    logger.error(f"Error generating outline from file: {gen_error}")
                    return {
                        "status": "error",
                        "error": f"Failed to generate outline from file: {str(gen_error)}"
                    }
            else:
                # No file outline found and no uploaded files
                return {
                    "status": "error",
                    "error": "No file outline found. Please ensure you uploaded a file during requirements confirmation."
                }

    except Exception as e:
        logger.error(f"Error generating file outline: {e}")
        return {
            "status": "error",
            "error": str(e)
        }

@router.post("/projects/{project_id}/update-outline")
async def update_project_outline(
    project_id: str,
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Update project outline content"""
    try:
        data = await request.json()
        outline_content = data.get('outline_content', '')

        success = await ppt_service.update_project_outline(project_id, outline_content)
        if success:
            return {"status": "success", "message": "Outline updated"}
        else:
            raise HTTPException(status_code=500, detail="Failed to update outline")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/projects/{project_id}/confirm-outline")
async def confirm_project_outline(
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """Confirm project outline and enable PPT generation"""
    try:
        success = await ppt_service.confirm_project_outline(project_id)
        if success:
            return {"status": "success", "message": "Outline confirmed"}
        else:
            raise HTTPException(status_code=500, detail="Failed to confirm outline")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/projects/{project_id}/todo-editor")
async def web_project_todo_editor(
    request: Request,
    project_id: str,
    auto_start: bool = False,
    user: User = Depends(get_current_user_required)
):
    """Project TODO board with editor"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return templates.TemplateResponse("error.html", {
                "request": request,
                "error": "Project not found"
            })

        return templates.TemplateResponse("todo_board_with_editor.html", {
            "request": request,
            "todo_board": project.todo_board,
            "project": project,
            "auto_start": auto_start
        })

    except Exception as e:
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })

@router.post("/projects/{project_id}/confirm-requirements")
async def confirm_project_requirements(
    request: Request,
    project_id: str,
    topic: str = Form(...),
    audience_type: str = Form(...),
    custom_audience: str = Form(None),
    page_count_mode: str = Form("ai_decide"),
    min_pages: int = Form(8),
    max_pages: int = Form(15),
    fixed_pages: int = Form(10),
    ppt_style: str = Form("general"),
    custom_style_prompt: str = Form(None),
    description: str = Form(None),
    content_source: str = Form("manual"),
    file_upload: List[UploadFile] = File(None),
    file_processing_mode: str = Form("markitdown"),
    content_analysis_depth: str = Form("standard"),
    user: User = Depends(get_current_user_required)
):
    """Confirm project requirements and generate TODO list - 支持多文件上传和联网搜索集成"""
    try:
        # Get project to access original requirements
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Extract network_mode from project metadata (set during project creation)
        network_mode = False
        if project.project_metadata and isinstance(project.project_metadata, dict):
            network_mode = project.project_metadata.get("network_mode", False)

        # Process audience information
        target_audience = audience_type
        if audience_type == "自定义" and custom_audience:
            target_audience = custom_audience

        # Extract language from project metadata (set during project creation)
        language = "zh"  # Default language
        if project.project_metadata and isinstance(project.project_metadata, dict):
            language = project.project_metadata.get("language", "zh")

        # Handle file upload if content source is file
        file_outline = None
        if content_source == "file" and file_upload:
            # Process uploaded files (support multiple files) and generate outline
            # 使用项目创建时的 network_mode 和 language 参数
            file_outline = await _process_uploaded_files_for_outline(
                file_upload, topic, target_audience, page_count_mode, min_pages, max_pages,
                fixed_pages, ppt_style, custom_style_prompt,
                file_processing_mode, content_analysis_depth, project.requirements,
                enable_web_search=network_mode,  # 使用项目的 network_mode
                scenario=project.scenario,  # 传递场景参数
                language=language  # 传递用户选择的语言参数
            )

            # Update topic if it was extracted from file
            if file_outline and file_outline.get('title') and not topic.strip():
                topic = file_outline['title']

        # Process page count settings
        page_count_settings = {
            "mode": page_count_mode,
            "min_pages": min_pages if page_count_mode == "custom_range" else None,
            "max_pages": max_pages if page_count_mode == "custom_range" else None,
            "fixed_pages": fixed_pages if page_count_mode == "fixed" else None
        }

        # Update project with confirmed requirements
        confirmed_requirements = {
            "topic": topic,
            "requirements": project.requirements,  # 使用项目创建时的具体要求
            "target_audience": target_audience,
            "audience_type": audience_type,
            "custom_audience": custom_audience if audience_type == "自定义" else None,
            "page_count_settings": page_count_settings,
            "ppt_style": ppt_style,
            "custom_style_prompt": custom_style_prompt if ppt_style == "custom" else None,
            "description": description,
            "content_source": content_source,
            "file_processing_mode": file_processing_mode if content_source == "file" else None,
            "content_analysis_depth": content_analysis_depth if content_source == "file" else None,
            "file_generated_outline": file_outline
        }

        # 如果是文件项目，保存文件信息
        if content_source == "file" and file_outline and 'file_info' in file_outline:
            file_info = file_outline['file_info']
            file_path = file_info.get('file_path') or file_info.get('merged_file_path')
            filename = file_info.get('filename') or file_info.get('merged_filename')
            uploaded_files = file_info.get('uploaded_files')

            file_metadata = {}
            if file_path:
                file_metadata["file_path"] = file_path
            if filename:
                file_metadata["filename"] = filename
            if uploaded_files:
                file_metadata["uploaded_files"] = uploaded_files

            if file_metadata:
                confirmed_requirements.update(file_metadata)

        # Store confirmed requirements in project
        # 直接确认需求并更新TODO板，无需AI生成待办清单
        success = await ppt_service.confirm_requirements_and_update_workflow(project_id, confirmed_requirements)

        if not success:
            raise Exception("需求确认失败")

        # Return JSON success response for AJAX request
        from fastapi.responses import JSONResponse
        return JSONResponse({
            "status": "success",
            "message": "需求确认完成",
            "redirect_url": f"/projects/{project_id}/todo"
        })

    except Exception as e:
        from fastapi.responses import JSONResponse
        return JSONResponse({
            "status": "error",
            "message": str(e)
        }, status_code=500)

@router.get("/projects/{project_id}/stage-stream/{stage_id}")
async def stream_stage_response(
    project_id: str,
    stage_id: str,
    user: User = Depends(get_current_user_required)
):
    """Stream AI response for a complete stage"""

    async def generate_stage_stream():
        try:
            # Get project and stage info
            project = await ppt_service.project_manager.get_project(project_id)
            if not project:
                yield f"data: {json.dumps({'error': 'Project not found'})}\n\n"
                return

            if not project.confirmed_requirements:
                yield f"data: {json.dumps({'error': 'Project requirements not confirmed'})}\n\n"
                return

            todo_board = await ppt_service.get_project_todo_board(project_id)
            if not todo_board:
                yield f"data: {json.dumps({'error': 'TODO board not found'})}\n\n"
                return

            # Find the stage
            stage = None
            for s in todo_board.stages:
                if s.id == stage_id:
                    stage = s
                    break

            if not stage:
                yield f"data: {json.dumps({'error': 'Stage not found'})}\n\n"
                return

            # Extract confirmed requirements from project
            confirmed_requirements = project.confirmed_requirements

            # Check if stage is already running or completed
            if stage.status == "running":
                yield f"data: {json.dumps({'error': 'Stage is already running'})}\n\n"
                return
            elif stage.status == "completed":
                yield f"data: {json.dumps({'error': 'Stage is already completed'})}\n\n"
                return

            # Update stage status to running
            await ppt_service.project_manager.update_stage_status(
                project_id, stage_id, "running", 0.0
            )

            # Execute the complete stage using the enhanced service
            try:
                if stage_id == "outline_generation":
                    response_content = await ppt_service._execute_outline_generation(
                        project_id, confirmed_requirements, ppt_service._load_prompts_md_system_prompt()
                    )
                elif stage_id == "ppt_creation":
                    response_content = await ppt_service._execute_ppt_creation(
                        project_id, confirmed_requirements, ppt_service._load_prompts_md_system_prompt()
                    )
                else:
                    # Fallback for other stages
                    response_content = await ppt_service._execute_general_stage(
                        project_id, stage_id, confirmed_requirements
                    )

                # Stream the response word by word for better UX
                if isinstance(response_content, dict):
                    content_text = response_content.get('message', str(response_content))
                else:
                    content_text = str(response_content)

                words = content_text.split()
                for i, word in enumerate(words):
                    yield f"data: {json.dumps({'content': word + ' ', 'done': False})}\n\n"
                    await asyncio.sleep(0.05)  # Small delay for streaming effect

            except Exception as e:
                # Fallback to basic stage execution
                prompt = f"""
作为PPT生成助手，请完成以下阶段任务：

项目主题：{project.topic}
项目场景：{project.scenario}
项目要求：{project.requirements or '无特殊要求'}

当前阶段：{stage.name}
阶段描述：{stage.description}

请根据以上信息完成当前阶段的完整任务，并提供详细的执行结果。
"""

                # Stream AI response using real streaming
                async for chunk in ppt_service.ai_provider.stream_text_completion(
                    prompt=prompt,
                    max_tokens=2000,
                    temperature=ai_config.temperature,
                    top_p=ai_config.top_p
                ):
                    if chunk:
                        yield f"data: {json.dumps({'content': chunk, 'done': False})}\n\n"

            # Update stage status to completed
            await ppt_service.project_manager.update_stage_status(
                project_id, stage_id, "completed", 100.0
            )

            # Send completion signal
            yield f"data: {json.dumps({'content': '', 'done': True})}\n\n"

        except Exception as e:
            yield f"data: {json.dumps({'error': str(e)})}\n\n"

    return StreamingResponse(
        generate_stage_stream(),
        media_type="text/plain",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"}
    )



@router.get("/projects/{project_id}/edit", response_class=HTMLResponse)
async def edit_project_ppt(
    request: Request,
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """Edit PPT slides with advanced editor"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # 允许编辑器在PPT生成过程中显示，提供更好的用户体验
        # 如果没有slides_data，创建一个空的结构供编辑器使用
        if not project.slides_data:
            project.slides_data = []

        return templates.TemplateResponse("project_slides_editor.html", {
            "request": request,
            "project": project,
            "enable_auto_layout_repair": ai_config.enable_auto_layout_repair
        })

    except Exception as e:
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })

@router.post("/api/projects/{project_id}/update-html")
async def update_project_html(
    project_id: str,
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Update project HTML content and mark all slides as user-edited"""
    try:
        data = await request.json()
        slides_html = data.get('slides_html', '')

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Update project HTML
        project.slides_html = slides_html
        project.updated_at = time.time()

        # 解析HTML内容，提取各个页面并标记为用户编辑
        if project.slides_data and slides_html:
            try:
                # 解析HTML内容，提取各个页面
                updated_slides_data = await _extract_slides_from_html(slides_html, project.slides_data)

                # 标记所有页面为用户编辑状态
                for slide_data in updated_slides_data:
                    slide_data["is_user_edited"] = True

                # 更新项目的slides_data
                project.slides_data = updated_slides_data

                logger.info(f"Marked {len(updated_slides_data)} slides as user-edited for project {project_id}")

            except Exception as parse_error:
                logger.warning(f"Failed to parse HTML content for slide extraction: {parse_error}")
                # 如果解析失败，至少标记现有的slides_data为用户编辑
                if project.slides_data:
                    for slide_data in project.slides_data:
                        slide_data["is_user_edited"] = True

        # 保存更新的HTML和slides_data到数据库
        try:
            from ..services.db_project_manager import DatabaseProjectManager
            db_manager = DatabaseProjectManager()

            # 保存幻灯片HTML和数据到数据库
            save_success = await db_manager.save_project_slides(
                project_id,
                project.slides_html,
                project.slides_data or []
            )

            if save_success:
                logger.info(f"Successfully saved updated HTML and slides data to database for project {project_id}")
            else:
                logger.error(f"Failed to save updated HTML and slides data to database for project {project_id}")

        except Exception as save_error:
            logger.error(f"Exception while saving updated HTML and slides data to database: {save_error}")
            # 继续返回成功，因为内存中的数据已经更新

        return {"status": "success", "message": "HTML updated successfully and slides marked as user-edited"}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/api/projects/{project_id}")
async def get_project_data(
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """Get project data for real-time updates"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        return {
            "project_id": project.project_id,
            "title": project.title,
            "status": project.status,
            "slides_data": project.slides_data or [],
            "slides_count": len(project.slides_data) if project.slides_data else 0,
            "updated_at": project.updated_at
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.put("/api/projects/{project_id}/slides")
async def update_project_slides(
    project_id: str,
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Update project slides data"""
    try:
        logger.info(f"🔄 开始更新项目 {project_id} 的幻灯片数据")

        data = await request.json()
        slides_data = data.get('slides_data', [])

        logger.info(f"📊 接收到 {len(slides_data)} 页幻灯片数据")

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            logger.error(f"❌ 项目 {project_id} 不存在")
            raise HTTPException(status_code=404, detail="Project not found")

        logger.info(f"📝 更新项目幻灯片数据...")

        # Update project slides data
        project.slides_data = slides_data
        project.updated_at = time.time()

        # Regenerate combined HTML
        if slides_data:
            # 安全地获取大纲标题
            outline_title = project.title
            if project.outline:
                if isinstance(project.outline, dict):
                    outline_title = project.outline.get('title', project.title)
                elif hasattr(project.outline, 'title'):
                    outline_title = project.outline.title

            project.slides_html = ppt_service._combine_slides_to_full_html(
                slides_data, outline_title
            )

        # 标记所有幻灯片为用户编辑状态
        for i, slide_data in enumerate(project.slides_data):
            slide_data["is_user_edited"] = True

        # 保存更新的幻灯片数据到数据库
        save_success = False
        save_error_message = None

        try:
            from ..services.db_project_manager import DatabaseProjectManager
            db_manager = DatabaseProjectManager()

            # 保存幻灯片数据到数据库
            save_success = await db_manager.save_project_slides(
                project_id,
                project.slides_html or "",
                project.slides_data
            )

            if save_success:
                logger.info(f"Successfully saved updated slides data to database for project {project_id}")
            else:
                logger.error(f"Failed to save updated slides data to database for project {project_id}")
                save_error_message = "Failed to save slides data to database"

        except Exception as save_error:
            logger.error(f"❌ 保存幻灯片数据到数据库时发生异常: {save_error}")
            import traceback
            traceback.print_exc()
            save_success = False
            save_error_message = str(save_error)

        # 根据保存结果返回相应的响应
        if save_success:
            return {
                "status": "success",
                "success": True,
                "message": "Slides updated and saved to database successfully"
            }
        else:
            # 即使数据库保存失败，内存中的数据已经更新，所以仍然返回成功，但包含警告信息
            return {
                "status": "success",
                "success": True,
                "message": "Slides updated in memory successfully",
                "warning": f"Database save failed: {save_error_message}",
                "database_saved": False
            }

    except Exception as e:
        logger.error(f"❌ 更新项目幻灯片数据时发生错误: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/api/projects/{project_id}/regenerate-html")
async def regenerate_project_html(project_id: str):
    """Regenerate project HTML with fixed encoding"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        if not project.slides_data:
            raise HTTPException(status_code=400, detail="No slides data found")

        # Regenerate combined HTML using the fixed method
        # 安全地获取大纲标题
        outline_title = project.title
        if project.outline:
            if isinstance(project.outline, dict):
                outline_title = project.outline.get('title', project.title)
            elif hasattr(project.outline, 'title'):
                outline_title = project.outline.title

        project.slides_html = ppt_service._combine_slides_to_full_html(
            project.slides_data, outline_title
        )

        project.updated_at = time.time()

        # 保存重新生成的HTML到数据库
        try:
            from ..services.db_project_manager import DatabaseProjectManager
            db_manager = DatabaseProjectManager()

            # 保存幻灯片数据到数据库
            save_success = await db_manager.save_project_slides(
                project_id,
                project.slides_html,
                project.slides_data
            )

            if save_success:
                logger.info(f"Successfully saved regenerated HTML to database for project {project_id}")
            else:
                logger.error(f"Failed to save regenerated HTML to database for project {project_id}")

        except Exception as save_error:
            logger.error(f"Exception while saving regenerated HTML to database: {save_error}")
            # 继续返回成功，因为内存中的数据已经更新

        return {
            "success": True,
            "message": "Project HTML regenerated successfully"
        }

    except Exception as e:
        return {"success": False, "error": str(e)}

@router.post("/api/projects/{project_id}/slides/{slide_number}/regenerate")
async def regenerate_slide(project_id: str, slide_number: int):
    """Regenerate a specific slide"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        if not project.outline:
            raise HTTPException(status_code=400, detail="Project outline not found")

        if not project.confirmed_requirements:
            raise HTTPException(status_code=400, detail="Project requirements not confirmed")

        # Handle different outline structures
        if isinstance(project.outline, dict):
            slides = project.outline.get('slides', [])
        else:
            # If outline is a PPTOutline object
            slides = project.outline.slides if hasattr(project.outline, 'slides') else []

        if slide_number < 1 or slide_number > len(slides):
            raise HTTPException(status_code=400, detail="Invalid slide number")

        # 关键修复：当 slides_data 缺页（例如只存在第2、3页）时，列表索引会错位，
        # 可能导致重新生成第2页时覆盖第1页。这里按 outline/page_number 归一化 slides_data。
        try:
            outline_total = len(slides)
            if outline_total > 0:
                if project.slides_data is None:
                    project.slides_data = []

                normalized = [None] * outline_total
                unplaced = []

                for s in (project.slides_data or []):
                    if not isinstance(s, dict):
                        continue
                    pn = s.get("page_number", None)
                    if isinstance(pn, str):
                        try:
                            pn = int(pn)
                        except Exception:
                            pn = None
                    if isinstance(pn, int) and 1 <= pn <= outline_total and normalized[pn - 1] is None:
                        normalized[pn - 1] = s
                    else:
                        unplaced.append(s)

                for s in unplaced:
                    try:
                        idx = normalized.index(None)
                    except ValueError:
                        break
                    normalized[idx] = s

                for i in range(outline_total):
                    if normalized[i] is None:
                        oslide = slides[i] if i < len(slides) else {}
                        title = oslide.get("title") if isinstance(oslide, dict) else None
                        slide_type = (oslide.get("slide_type") or oslide.get("type")) if isinstance(oslide, dict) else None
                        content_points = oslide.get("content_points") if isinstance(oslide, dict) else None
                        normalized[i] = {
                            "page_number": i + 1,
                            "title": title or f"Slide {i + 1}",
                            "html_content": "<div>Pending</div>",
                            "slide_type": slide_type or "content",
                            "content_points": content_points if isinstance(content_points, list) else [],
                            "is_user_edited": False,
                        }
                    else:
                        normalized[i]["page_number"] = i + 1

                project.slides_data = normalized
        except Exception as normalize_err:
            logger.warning(f"Slides normalization skipped for regenerate_slide {project_id}: {normalize_err}")

        slide_data = slides[slide_number - 1]

        # Load system prompt
        system_prompt = ppt_service._load_prompts_md_system_prompt()

        # Ensure project has a global template selected (use default if none selected)
        selected_template = await ppt_service._ensure_global_master_template_selected(project_id)

        # Regenerate the slide using template-based generation if template is available
        if selected_template:
            logger.info(f"Regenerating slide {slide_number} using template: {selected_template['template_name']}")
            new_html_content = await ppt_service._generate_slide_with_template(
                slide_data, selected_template, slide_number, len(slides), project.confirmed_requirements
            )
        else:
            # Fallback to original generation method if no template available
            logger.warning(f"No template available for project {project_id}, using fallback generation")
            new_html_content = await ppt_service._generate_single_slide_html_with_prompts(
                slide_data, project.confirmed_requirements, system_prompt, slide_number, len(slides),
                slides, project.slides_data, project_id=project_id
            )

        # Update the slide in project data
        if not project.slides_data:
            project.slides_data = []

        # Ensure slides_data has enough entries
        while len(project.slides_data) < slide_number:
            new_page_number = len(project.slides_data) + 1
            project.slides_data.append({
                "page_number": new_page_number,
                "title": f"第{new_page_number}页",
                "html_content": "<div>待生成</div>",
                "slide_type": "content",
                "content_points": [],
                "is_user_edited": False
            })

        # Update the specific slide - 保持与现有数据结构一致
        existing_slide = project.slides_data[slide_number - 1] if slide_number <= len(project.slides_data) else {}

        # 更新幻灯片数据，保留现有字段并确保必要字段存在
        updated_slide = {
            "page_number": slide_number,
            "title": slide_data.get('title', f'第{slide_number}页'),
            "html_content": new_html_content,
            "slide_type": slide_data.get('slide_type', existing_slide.get('slide_type', 'content')),
            "content_points": slide_data.get('content_points', existing_slide.get('content_points', [])),
            "is_user_edited": existing_slide.get('is_user_edited', False),
            # 保留其他可能存在的字段
            **{k: v for k, v in existing_slide.items() if k not in ['page_number', 'title', 'html_content', 'slide_type', 'content_points', 'is_user_edited']}
        }

        project.slides_data[slide_number - 1] = updated_slide

        # Regenerate combined HTML
        outline_title = project.title
        if isinstance(project.outline, dict):
            outline_title = project.outline.get('title', project.title)
        elif hasattr(project.outline, 'title'):
            outline_title = project.outline.title

        project.slides_html = ppt_service._combine_slides_to_full_html(
            project.slides_data, outline_title
        )

        project.updated_at = time.time()

        # 保存更新后的幻灯片数据到数据库
        try:
            from ..services.db_project_manager import DatabaseProjectManager
            db_manager = DatabaseProjectManager()

            # 只保存单个重新生成的幻灯片，而不是整个项目的幻灯片数据
            # 这样可以避免删除所有幻灯片再重新创建的问题
            save_success = await db_manager.save_single_slide(
                project_id,
                slide_number - 1,  # 转换为0基索引
                updated_slide
            )

            if save_success:
                logger.info(f"Successfully saved regenerated slide {slide_number} to database for project {project_id}")

                # 同时更新项目的slides_html字段
                await db_manager.update_project_data(project_id, {
                    "slides_html": project.slides_html,
                    "updated_at": project.updated_at
                })
            else:
                logger.error(f"Failed to save regenerated slide {slide_number} to database for project {project_id}")

        except Exception as save_error:
            logger.error(f"Exception while saving regenerated slide to database: {save_error}")
            # 继续返回成功，因为内存中的数据已经更新

        return {
            "success": True,
            "message": f"Slide {slide_number} regenerated successfully",
            "slide_data": project.slides_data[slide_number - 1]
        }

    except Exception as e:
        return {"success": False, "error": str(e)}

@router.post("/api/projects/{project_id}/slides/batch-regenerate")
async def batch_regenerate_slides(project_id: str, payload: SlideBatchRegenerateRequest):
    """Regenerate multiple slides (or all slides) in one request."""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        if not project.outline:
            raise HTTPException(status_code=400, detail="Project outline not found")

        if not project.confirmed_requirements:
            raise HTTPException(status_code=400, detail="Project requirements not confirmed")

        if isinstance(project.outline, dict):
            outline_slides = project.outline.get("slides", [])
            outline_title = project.outline.get("title", project.title)
        else:
            outline_slides = project.outline.slides if hasattr(project.outline, "slides") else []
            outline_title = project.outline.title if hasattr(project.outline, "title") else project.title

        total_slides = len(outline_slides)
        if total_slides <= 0:
            raise HTTPException(status_code=400, detail="No slides found in outline")

        # 关键修复：当 slides_data 缺页时，按 outline/page_number 归一化，避免批量重新生成错页写入。
        try:
            if project.slides_data is None:
                project.slides_data = []

            normalized = [None] * total_slides
            unplaced = []

            for s in (project.slides_data or []):
                if not isinstance(s, dict):
                    continue
                pn = s.get("page_number", None)
                if isinstance(pn, str):
                    try:
                        pn = int(pn)
                    except Exception:
                        pn = None
                if isinstance(pn, int) and 1 <= pn <= total_slides and normalized[pn - 1] is None:
                    normalized[pn - 1] = s
                else:
                    unplaced.append(s)

            for s in unplaced:
                try:
                    idx = normalized.index(None)
                except ValueError:
                    break
                normalized[idx] = s

            for i in range(total_slides):
                if normalized[i] is None:
                    oslide = outline_slides[i] if i < len(outline_slides) else {}
                    title = oslide.get("title") if isinstance(oslide, dict) else None
                    slide_type = (oslide.get("slide_type") or oslide.get("type")) if isinstance(oslide, dict) else None
                    content_points = oslide.get("content_points") if isinstance(oslide, dict) else None
                    normalized[i] = {
                        "page_number": i + 1,
                        "title": title or f"Slide {i + 1}",
                        "html_content": "<div>Pending</div>",
                        "slide_type": slide_type or "content",
                        "content_points": content_points if isinstance(content_points, list) else [],
                        "is_user_edited": False
                    }
                else:
                    normalized[i]["page_number"] = i + 1

            project.slides_data = normalized
        except Exception as normalize_err:
            logger.warning(f"Slides normalization skipped for batch_regenerate {project_id}: {normalize_err}")

        # Determine target indices (0-based).
        if payload.regenerate_all or not payload.slide_indices:
            target_indices = list(range(total_slides))
        else:
            target_indices = sorted(set(payload.slide_indices))

        invalid_indices = [i for i in target_indices if i < 0 or i >= total_slides]
        if invalid_indices:
            raise HTTPException(status_code=400, detail=f"Invalid slide indices: {invalid_indices}")

        # Prepare generation context once.
        system_prompt = ppt_service._load_prompts_md_system_prompt()
        selected_template = await ppt_service._ensure_global_master_template_selected(project_id)

        if project.slides_data is None:
            project.slides_data = []

        # Ensure slides_data has enough entries for all slides.
        while len(project.slides_data) < total_slides:
            page_number = len(project.slides_data) + 1
            project.slides_data.append({
                "page_number": page_number,
                "title": f"Slide {page_number}",
                "html_content": "<div>Pending</div>",
                "slide_type": "content",
                "content_points": [],
                "is_user_edited": False
            })

        results: List[Dict[str, Any]] = []

        for slide_index in target_indices:
            slide_number = slide_index + 1  # 1-based for prompts/templates
            slide_outline = outline_slides[slide_index]
            try:
                if selected_template:
                    new_html_content = await ppt_service._generate_slide_with_template(
                        slide_outline,
                        selected_template,
                        slide_number,
                        total_slides,
                        project.confirmed_requirements
                    )
                else:
                    new_html_content = await ppt_service._generate_single_slide_html_with_prompts(
                        slide_outline,
                        project.confirmed_requirements,
                        system_prompt,
                        slide_number,
                        total_slides,
                        outline_slides,
                        project.slides_data,
                        project_id=project_id
                    )

                existing_slide = project.slides_data[slide_index] if slide_index < len(project.slides_data) else {}
                updated_slide = {
                    "page_number": slide_number,
                    "title": slide_outline.get("title", existing_slide.get("title", f"Slide {slide_number}")),
                    "html_content": new_html_content,
                    "slide_type": slide_outline.get("slide_type", existing_slide.get("slide_type", "content")),
                    "content_points": slide_outline.get("content_points", existing_slide.get("content_points", [])),
                    "is_user_edited": existing_slide.get("is_user_edited", False),
                    **{k: v for k, v in (existing_slide or {}).items() if k not in ["page_number", "title", "html_content", "slide_type", "content_points", "is_user_edited"]}
                }

                project.slides_data[slide_index] = updated_slide

                results.append({
                    "slide_index": slide_index,
                    "slide_number": slide_number,
                    "success": True,
                    "slide_data": updated_slide
                })
            except Exception as e:
                logger.error(f"Batch regenerate failed for project {project_id} slide {slide_number}: {e}")
                results.append({
                    "slide_index": slide_index,
                    "slide_number": slide_number,
                    "success": False,
                    "error": str(e)
                })

        # Rebuild combined HTML once.
        project.slides_html = ppt_service._combine_slides_to_full_html(project.slides_data, outline_title)
        project.updated_at = time.time()

        updated_count = len([r for r in results if r.get("success")])

        # Persist: save regenerated slides and update project HTML.
        try:
            from ..services.db_project_manager import DatabaseProjectManager
            db_manager = DatabaseProjectManager()

            for r in results:
                if not r.get("success") or not r.get("slide_data"):
                    continue
                await db_manager.save_single_slide(project_id, int(r["slide_index"]), r["slide_data"])

            await db_manager.update_project_data(project_id, {
                "slides_html": project.slides_html,
                "updated_at": project.updated_at
            })
        except Exception as save_error:
            logger.error(f"Batch regenerate DB save failed for project {project_id}: {save_error}")

        return {
            "success": updated_count > 0,
            "updated_count": updated_count,
            "total_requested": len(target_indices),
            "results": results
        }

    except HTTPException:
        raise
    except Exception as e:
        return {"success": False, "error": str(e)}

@router.post("/api/projects/{project_id}/slides/{slide_index}/auto-repair-layout")
async def auto_repair_layout(
    project_id: str,
    slide_index: int,
    request: AutoLayoutRepairRequest,
    user: User = Depends(get_current_user_required)
):
    """Run multimodal layout inspection and repair workflow for a single slide."""
    try:
        if slide_index < 1:
            raise HTTPException(status_code=400, detail="Slide index must be >= 1")

        html_content = (request.html_content or "").strip()
        if not html_content:
            raise HTTPException(status_code=400, detail="HTML content is required")

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        slides_data = project.slides_data or []
        total_pages = len(slides_data)
        if total_pages == 0:
            total_pages = request.slide_data.get("total_pages") or request.slide_data.get("totalSlides") or slide_index

        slide_payload = dict(request.slide_data or {})
        slide_payload.setdefault("page_number", slide_index)
        slide_payload.setdefault("title", slide_payload.get("title", f"第{slide_index}页"))

        repaired_html = await ppt_service._apply_auto_layout_repair(
            html_content,
            slide_payload,
            slide_index,
            total_pages or slide_index
        )

        changed = repaired_html.strip() != html_content

        if project.slides_data is None:
            project.slides_data = []

        while len(project.slides_data) < slide_index:
            page_number = len(project.slides_data) + 1
            project.slides_data.append({
                "page_number": page_number,
                "title": f"第{page_number}页",
                "html_content": "",
                "slide_type": "content",
                "content_points": [],
                "is_user_edited": False
            })

        existing_slide = project.slides_data[slide_index - 1]
        updated_slide = {
            **existing_slide,
            "page_number": slide_index,
            "title": slide_payload.get("title", existing_slide.get("title", f"第{slide_index}页")),
            "html_content": repaired_html,
        }

        project.slides_data[slide_index - 1] = updated_slide

        if changed:
            outline_title = project.title
            if isinstance(project.outline, dict):
                outline_title = project.outline.get('title', project.title)
            elif hasattr(project.outline, 'title'):
                outline_title = project.outline.title

            project.slides_html = ppt_service._combine_slides_to_full_html(
                project.slides_data,
                outline_title
            )
            project.updated_at = time.time()

        try:
            from ..services.db_project_manager import DatabaseProjectManager
            db_manager = DatabaseProjectManager()
            await db_manager.save_single_slide(project_id, slide_index - 1, updated_slide)

            if changed:
                await db_manager.update_project_data(project_id, {
                    "slides_html": project.slides_html,
                    "updated_at": project.updated_at
                })

        except Exception as save_error:
            logger.error(f"Failed to persist auto layout repair result: {save_error}")

        return {
            "success": True,
            "repaired_html": repaired_html,
            "changed": changed
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Auto layout repair failed: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/api/ai/slide-edit")
async def ai_slide_edit(
    request: AISlideEditRequest,
    user: User = Depends(get_current_user_required)
):
    """AI编辑幻灯片接口"""
    try:
        # 获取AI提供者
        provider, settings = get_role_provider("editor")

        # 构建AI编辑上下文
        outline_info = ""
        if request.slideOutline:
            outline_info = f"""
当前幻灯片大纲信息：
- 幻灯片类型：{request.slideOutline.get('slide_type', '未知')}
- 描述：{request.slideOutline.get('description', '无')}
- 要点：{', '.join(request.slideOutline.get('content_points', [])) if request.slideOutline.get('content_points') else '无'}
"""

        context = f"""
你是一位专业的PPT设计师和编辑助手。用户想要对当前幻灯片进行编辑修改。

当前幻灯片信息：
- 页码：第{request.slideIndex}页
- 标题：{request.slideTitle}
- 项目主题：{request.projectInfo.get('title', '未知')}
- 项目场景：{request.projectInfo.get('scenario', '未知')}
{outline_info}
用户的编辑要求：
{request.userRequest}

当前幻灯片的HTML内容：
{request.slideContent}

请根据用户的要求和幻灯片大纲信息，提供以下内容：
1. 对用户要求的理解和分析
2. 具体的修改建议
3. 如果需要，提供修改后的完整HTML代码

注意事项：
- 确保修改后的内容符合PPT演示的专业标准和大纲要求
- 生成的HTML应该是完整的，包含必要的CSS样式
- 保持1280x720的PPT标准尺寸
- 参考大纲信息中的要点和描述来优化内容
"""

        # 构建AI消息，包含对话历史
        messages = [
            AIMessage(role=MessageRole.SYSTEM, content="你是一位专业的PPT设计师和编辑助手，擅长根据用户需求修改和优化PPT内容。")
        ]

        # 添加对话历史
        if request.chatHistory:
            logger.debug(f"AI编辑接收到对话历史，共 {len(request.chatHistory)} 条消息")
            for i, chat_msg in enumerate(request.chatHistory):
                role = MessageRole.USER if chat_msg.get('role') == 'user' else MessageRole.ASSISTANT
                content = chat_msg.get('content', '')
                logger.debug(f"对话历史 {i+1}: {role.value} - {content[:100]}...")
                messages.append(AIMessage(role=role, content=content))
        else:
            logger.debug("AI编辑未接收到对话历史")

        # 添加当前用户请求
        messages.append(AIMessage(role=MessageRole.USER, content=context))

        # 调用AI生成回复
        response = await provider.chat_completion(
            messages=messages,
            max_tokens=ai_config.max_tokens,
            temperature=ai_config.temperature,
            top_p=ai_config.top_p,
            model=settings.get('model')
        )

        ai_response = response.content

        # 检查是否包含HTML代码
        new_html_content = None
        if "```html" in ai_response:
            import re
            html_match = re.search(r'```html\s*(.*?)\s*```', ai_response, re.DOTALL)
            if html_match:
                new_html_content = html_match.group(1).strip()

        return {
            "success": True,
            "response": ai_response,
            "newHtmlContent": new_html_content
        }

    except Exception as e:
        logger.error(f"AI编辑请求失败: {e}")
        return {
            "success": False,
            "error": str(e),
            "response": "抱歉，AI编辑服务暂时不可用。请稍后重试。"
        }

@router.post("/api/ai/slide-edit/stream")
async def ai_slide_edit_stream(
    request: AISlideEditRequest,
    user: User = Depends(get_current_user_required)
):
    """AI编辑幻灯片流式接口"""
    try:
        # 获取AI提供者
        provider, settings = get_role_provider("editor")

        # 构建AI编辑上下文
        outline_info = ""
        if request.slideOutline:
            outline_info = f"""
当前幻灯片大纲信息：
{request.slideOutline}
"""

        # 构建图片信息
        images_info = ""
        if request.images and len(request.images) > 0:
            images_info = f"""

用户上传的图片信息：
"""
            for i, image in enumerate(request.images, 1):
                url = image.get("url", "") if isinstance(image, dict) else ""
                # 避免把 data URL/base64 整段塞进文本上下文（图片会以多模态内容附带）
                if isinstance(url, str) and url.startswith("data:image"):
                    url_display = "（data URL 已随消息附带）"
                else:
                    url_display = url

                images_info += f"""
- 图片{i}：{image.get('name', '未知')}
  - URL：{url_display}
  - 大小：{image.get('size', '未知')}
  - 说明：请分析这张图片的内容，理解用户的意图，并根据编辑要求进行相应的处理
"""

        # 构建视觉上下文信息
        vision_context = ""
        if request.visionEnabled and request.slideScreenshot:
            vision_context = f"""

🔍 视觉上下文：
- 当前幻灯片的视觉截图已提供
- 请结合截图中的视觉内容来理解用户的编辑需求
- 注意截图中的布局、颜色、字体、图片位置等视觉元素
- 在提供编辑建议时，请考虑当前的视觉呈现效果
"""

        context = f"""
你是一位专业的PPT设计师和编辑助手。用户想要对当前幻灯片进行编辑修改。

当前幻灯片信息：
- 页码：第{request.slideIndex}页
- 标题：{request.slideTitle}
- 项目主题：{request.projectInfo.get('title', '未知')}
- 项目场景：{request.projectInfo.get('scenario', '未知')}
{outline_info}{images_info}{vision_context}
用户的编辑要求：
{request.userRequest}

当前幻灯片的HTML内容：
{request.slideContent}

请根据用户的要求和幻灯片大纲信息，提供以下内容：
1. 对用户要求的理解和分析
2. 具体的修改建议
3. 默认提供修改后的完整HTML代码

注意事项：
- 保持原有的设计风格和布局结构
- 确保修改后的内容符合PPT演示的专业标准和大纲要求
- 如果用户要求不明确，请提供多个可选方案
- 生成的HTML应该是完整的，包含必要的CSS样式
- 保持1280x720的PPT标准尺寸
- 参考大纲信息中的要点和描述来优化内容
"""

        # 构建AI消息，包含对话历史
        messages = [
            AIMessage(role=MessageRole.SYSTEM, content="你是一位专业的PPT设计师和编辑助手，擅长根据用户需求修改和优化PPT内容。")
        ]

        # 添加对话历史
        if request.chatHistory:
            logger.info(f"AI流式编辑接收到对话历史，共 {len(request.chatHistory)} 条消息")
            for i, chat_msg in enumerate(request.chatHistory):
                role = MessageRole.USER if chat_msg.get('role') == 'user' else MessageRole.ASSISTANT
                content = chat_msg.get('content', '')
                logger.info(f"对话历史 {i+1}: {role.value} - {content[:100]}...")
                messages.append(AIMessage(role=role, content=content))
        else:
            logger.info("AI流式编辑未接收到对话历史")

        # 添加当前用户请求（视觉模式：支持截图 + 上传图片的多模态内容）
        if request.visionEnabled:
            from ..ai.base import TextContent, ImageContent

            user_content = [TextContent(text=context)]

            if request.slideScreenshot:
                user_content.append(ImageContent(image_url={"url": request.slideScreenshot}))

            if request.images:
                for img in request.images:
                    url = (img or {}).get("url")
                    if url:
                        user_content.append(ImageContent(image_url={"url": url}))

            messages.append(AIMessage(role=MessageRole.USER, content=user_content))
        else:
            messages.append(AIMessage(role=MessageRole.USER, content=context))

        async def generate_ai_stream():
            try:
                # 发送开始信号
                yield f"data: {json.dumps({'type': 'start', 'content': ''})}\n\n"

                # 流式生成AI回复
                full_response = ""
                if hasattr(provider, 'stream_chat_completion'):
                    async for chunk in provider.stream_chat_completion(
                        messages=messages,
                        max_tokens=ai_config.max_tokens,
                        temperature=ai_config.temperature,
                        top_p=ai_config.top_p,
                        model=settings.get('model')
                    ):
                        if chunk:
                            full_response += chunk
                            yield f"data: {json.dumps({'type': 'content', 'content': chunk})}\n\n"
                else:
                    response = await provider.chat_completion(
                        messages=messages,
                        max_tokens=ai_config.max_tokens,
                        temperature=ai_config.temperature,
                        top_p=ai_config.top_p,
                        model=settings.get('model')
                    )
                    if response.content:
                        full_response = response.content
                        yield f"data: {json.dumps({'type': 'content', 'content': response.content})}\n\n"

                # 检查是否包含HTML代码 - 改进版本，支持多种格式
                new_html_content = None
                import re
                
                # 尝试多种HTML代码块格式
                html_patterns = [
                    r'```html\s*(.*?)\s*```',  # 标准格式
                    r'```HTML\s*(.*?)\s*```',  # 大写
                    r'```\s*html\s*(.*?)\s*```',  # 带空格
                    r'<html[^>]*>.*?</html>',  # 完整HTML文档
                    r'<div[^>]*style[^>]*>.*?</div>',  # PPT幻灯片div
                ]
                
                for pattern in html_patterns:
                    html_match = re.search(pattern, full_response, re.DOTALL | re.IGNORECASE)
                    if html_match:
                        new_html_content = html_match.group(1).strip() if html_match.groups() else html_match.group(0).strip()
                        logger.info(f"HTML内容提取成功，使用模式: {pattern}，内容长度: {len(new_html_content)}")
                        break
                
                if not new_html_content:
                    logger.warning(f"未能从AI响应中提取HTML内容。响应长度: {len(full_response)}")
                    logger.debug(f"AI完整响应: {full_response[:500]}...")

                # 发送完成信号
                yield f"data: {json.dumps({'type': 'complete', 'content': '', 'newHtmlContent': new_html_content, 'fullResponse': full_response})}\n\n"

            except Exception as e:
                logger.error(f"AI流式编辑请求失败: {e}")
                yield f"data: {json.dumps({'type': 'error', 'content': '', 'error': str(e)})}\n\n"

        return StreamingResponse(
            generate_ai_stream(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Headers": "Cache-Control"
            }
        )

    except Exception as e:
        logger.error(f"AI流式编辑请求失败: {e}")
        return {
            "success": False,
            "error": str(e),
            "response": "抱歉，AI编辑服务暂时不可用。请稍后重试。"
        }

@router.post("/api/ai/slide-native-dialog/stream")
async def ai_slide_native_dialog_stream(
    request: AISlideNativeDialogRequest,
    user: User = Depends(get_current_user_required)
):
    """AI自由对话流式接口（不设系统提示词，仅当前页）"""
    try:
        provider, settings = get_role_provider("editor")

        images_info = ""
        if request.images:
            images_info = f"\n用户上传/粘贴的图片数量：{len(request.images)}（图片内容在消息中以多模态形式附带）\n"

        context = f"""
当前页面信息（仅此页）：
- 页码：第{request.slideIndex}页
- 标题：{request.slideTitle}

当前页面HTML内容：
{request.slideContent}

用户问题：
{request.userRequest}
{images_info}
"""

        messages: List[AIMessage] = []

        # 添加对话历史（忽略 system 角色，避免“系统提示词”进入模型）
        if request.chatHistory:
            for chat_msg in request.chatHistory:
                role_str = (chat_msg.get("role") or "").lower()
                if role_str == "system":
                    continue

                if role_str == "assistant":
                    role = MessageRole.ASSISTANT
                else:
                    role = MessageRole.USER

                messages.append(AIMessage(role=role, content=chat_msg.get("content", "")))

        # 添加当前用户请求（支持粘贴/上传图片的多模态内容）
        if request.images and len(request.images) > 0:
            from ..ai.base import TextContent, ImageContent

            user_content = [TextContent(text=context)]
            for img in request.images:
                url = img.get("url")
                if url:
                    user_content.append(ImageContent(image_url={"url": url}))

            messages.append(AIMessage(role=MessageRole.USER, content=user_content))
        else:
            messages.append(AIMessage(role=MessageRole.USER, content=context))

        async def generate_ai_stream():
            try:
                yield f"data: {json.dumps({'type': 'start', 'content': ''})}\n\n"

                full_response = ""
                if hasattr(provider, "stream_chat_completion"):
                    async for chunk in provider.stream_chat_completion(
                        messages=messages,
                        max_tokens=ai_config.max_tokens,
                        temperature=ai_config.temperature,
                        top_p=ai_config.top_p,
                        model=settings.get("model"),
                    ):
                        if chunk:
                            full_response += chunk
                            yield f"data: {json.dumps({'type': 'content', 'content': chunk})}\n\n"
                else:
                    response = await provider.chat_completion(
                        messages=messages,
                        max_tokens=ai_config.max_tokens,
                        temperature=ai_config.temperature,
                        top_p=ai_config.top_p,
                        model=settings.get("model"),
                    )
                    if response.content:
                        full_response = response.content
                        yield f"data: {json.dumps({'type': 'content', 'content': response.content})}\n\n"

                yield f"data: {json.dumps({'type': 'complete', 'content': '', 'fullResponse': full_response})}\n\n"

            except Exception as e:
                logger.error(f"AI自由对话流式请求失败: {e}")
                yield f"data: {json.dumps({'type': 'error', 'content': '', 'error': str(e)})}\n\n"

        return StreamingResponse(
            generate_ai_stream(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Headers": "Cache-Control",
            },
        )

    except Exception as e:
        logger.error(f"AI自由对话请求失败: {e}")
        return StreamingResponse(
            iter([f"data: {json.dumps({'type': 'error', 'content': '', 'error': str(e)})}\n\n"]),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Headers": "Cache-Control",
            },
        )

# 大纲AI优化请求数据模型
class OutlineAIOptimizeRequest(BaseModel):
    outline_content: str  # JSON格式的大纲内容
    user_request: str  # 用户的优化需求
    project_info: Dict[str, Any]  # 项目信息
    optimization_type: str = "full"  # full=全大纲优化, single=单页优化
    slide_index: Optional[int] = None  # 当optimization_type=single时使用
    language: Optional[str] = None  # 目标语言（如 zh/en/ja...），优先级高于大纲metadata.language

@router.post("/api/ai/optimize-outline")
async def ai_optimize_outline(
    request: OutlineAIOptimizeRequest,
    user: User = Depends(get_current_user_required)
):
    """AI优化大纲接口 - 支持全大纲优化和单页优化"""
    try:
        # 获取AI提供者
        provider, settings = get_role_provider("editor")
        
        # 解析大纲JSON
        try:
            outline_data = json.loads(request.outline_content)
        except json.JSONDecodeError as e:
            return {
                "success": False,
                "error": f"大纲JSON格式错误: {str(e)}"
            }

        def _normalize_language_code(value: Any) -> Optional[str]:
            if not isinstance(value, str):
                return None
            code = value.strip().lower()
            if not code:
                return None
            # Normalize common variants (zh-cn -> zh, en-us -> en, etc.)
            if code.startswith("zh"):
                return "zh"
            if code.startswith("en"):
                return "en"
            if code.startswith("ja"):
                return "ja"
            if code.startswith("ko"):
                return "ko"
            if code.startswith("fr"):
                return "fr"
            if code.startswith("de"):
                return "de"
            if code.startswith("es"):
                return "es"
            return code

        outline_language = None
        if isinstance(outline_data, dict):
            metadata = outline_data.get("metadata")
            if isinstance(metadata, dict):
                outline_language = metadata.get("language") or outline_data.get("language")
            else:
                outline_language = outline_data.get("language")

        target_language = (
            _normalize_language_code(request.language)
            or _normalize_language_code(outline_language)
            or "zh"
        )
        
        # 根据优化类型构建不同的提示词
        if request.optimization_type == "single" and request.slide_index is not None:
            # 单页优化
            if request.slide_index < 0 or request.slide_index >= len(outline_data.get('slides', [])):
                return {
                    "success": False,
                    "error": "无效的幻灯片索引"
                }
            
            slide = outline_data['slides'][request.slide_index]
            
            context = f"""Output language: {target_language}
你是一位专业的PPT大纲设计专家。用户想要优化PPT大纲中的第{request.slide_index + 1}页内容。

项目信息：
- 主题：{request.project_info.get('topic', '未知')}
- 场景：{request.project_info.get('scenario', '通用')}
- 目标受众：{request.project_info.get('target_audience', '普通大众')}

当前页面信息：
- 页码：第{slide.get('page_number', request.slide_index + 1)}页
- 标题：{slide.get('title', '未命名')}
- 类型：{slide.get('slide_type', 'content')}
- 内容要点：{json.dumps(slide.get('content_points', []), ensure_ascii=False, indent=2)}

用户的优化需求：
{request.user_request}

请根据用户需求优化这一页的内容。

【重要】直接返回优化后的JSON数据，不要包含任何解释性文字或markdown标记（如```json）。

返回格式示例：
{{
  "page_number": {slide.get('page_number', request.slide_index + 1)},
  "title": "优化后的标题",
  "subtitle": "副标题（可选）",
  "content_points": ["要点1", "要点2", "要点3"],
  "slide_type": "content",
  "description": "页面描述（可选）"
}}

优化要求：
1. 保持与整体大纲的连贯性和逻辑性
2. 确保内容要点清晰、具体、有价值
3. 标题要简洁有力，能够准确概括页面内容
4. content_points数组中的字符串可以包含代码示例（用```标记），这是合法的JSON字符串内容
5. 【关键】只返回纯JSON对象，不要用```json包裹整个JSON，不要添加任何其他解释文字
"""
        else:
            # 全大纲优化
            context = f"""Output language: {target_language}
你是一位专业的PPT大纲设计专家。用户想要优化整个PPT大纲。

项目信息：
- 主题：{request.project_info.get('topic', '未知')}
- 场景：{request.project_info.get('scenario', '通用')}
- 目标受众：{request.project_info.get('target_audience', '普通大众')}
- 当前页数：{len(outline_data.get('slides', []))}页

当前大纲：
{json.dumps(outline_data, ensure_ascii=False, indent=2)}

用户的优化需求：
{request.user_request}

请根据用户需求优化整个大纲。

【重要】直接返回完整的优化后的JSON数据，不要包含任何解释性文字、markdown标记或注释。

返回格式示例：
{{
  "title": "优化后的PPT标题",
  "slides": [
    {{
      "page_number": 1,
      "title": "页面标题",
      "subtitle": "副标题（可选）",
      "content_points": ["要点1", "要点2"],
      "slide_type": "title",
      "description": "页面描述（可选）"
    }}
  ],
  "metadata": {{
    "scenario": "{request.project_info.get('scenario', '通用')}",
    "language": "{target_language}",
    "target_audience": "{request.project_info.get('target_audience', '普通大众')}",
    "optimized": true
  }}
}}

优化要求：
1. 保持大纲的整体逻辑性和连贯性
2. 确保每页内容要点清晰、具体、有价值
3. 可以调整页面顺序、合并或拆分页面，但要保持总体结构合理
4. 标题要简洁有力
5. 【关键】只返回纯JSON格式，不要添加任何解释、注释或markdown标记
"""
        
        # 构建AI消息
        messages = [
            AIMessage(role=MessageRole.SYSTEM, content="你是一位专业的PPT大纲设计专家，擅长优化和改进PPT大纲结构和内容。你的回复必须是纯JSON格式，不要包含任何解释性文字、markdown标记或注释。"),
            AIMessage(role=MessageRole.SYSTEM, content=f"Output language: {target_language}. Return pure JSON only."),
            AIMessage(role=MessageRole.USER, content=context)
        ]
        
        # 调用AI生成回复
        response = await provider.chat_completion(
            messages=messages,
            temperature=ai_config.temperature,
            top_p=ai_config.top_p,
            model=settings.get('model')
        )
        
        ai_response = response.content
        
        # 智能提取JSON内容
        import re
        
        def extract_json_from_response(text: str) -> str:
            """从AI响应中提取JSON内容，支持多种格式"""
            
            # 优先方法: 查找第一个{到最后一个}之间的内容
            # 这样可以避免错误提取content_points字段内的代码块
            first_brace = text.find('{')
            last_brace = text.rfind('}')
            if first_brace != -1 and last_brace != -1 and first_brace < last_brace:
                potential_json = text[first_brace:last_brace + 1]
                # 尝试解析，如果成功则返回
                try:
                    json.loads(potential_json)
                    return potential_json.strip()
                except json.JSONDecodeError:
                    # 如果解析失败，尝试清理注释后再试
                    cleaned_json = re.sub(r'//[^\n]*', '', potential_json)  # 单行注释
                    cleaned_json = re.sub(r'/\*.*?\*/', '', cleaned_json, flags=re.DOTALL)  # 多行注释
                    try:
                        json.loads(cleaned_json)
                        return cleaned_json.strip()
                    except json.JSONDecodeError:
                        pass  # 继续尝试其他方法
            
            # 备用方法: 提取markdown代码块中的JSON（仅当标记为json时）
            # 使用更严格的匹配，确保是JSON代码块而不是其他代码块
            json_match = re.search(r'```json\s*(.*?)\s*```', text, re.DOTALL | re.IGNORECASE)
            if json_match:
                extracted = json_match.group(1).strip()
                # 验证提取的内容是否是有效JSON
                try:
                    json.loads(extracted)
                    return extracted
                except json.JSONDecodeError:
                    pass  # 继续尝试其他方法
            
            # 最后尝试: 直接返回清理后的文本
            cleaned = text.strip()
            if cleaned.startswith('{'):
                return cleaned
            
            # 尝试找到JSON开始的位置
            for line in cleaned.split('\n'):
                line = line.strip()
                if line.startswith('{'):
                    start_idx = cleaned.find(line)
                    return cleaned[start_idx:].strip()
            
            return cleaned
        
        optimized_json = extract_json_from_response(ai_response)
        
        # 验证JSON格式
        try:
            optimized_data = json.loads(optimized_json)
        except json.JSONDecodeError as e:
            # 提供更详细的错误信息，帮助调试
            return {
                "success": False,
                "error": f"AI返回的内容不是有效的JSON格式: {str(e)}",
                "raw_response": ai_response,
                "extracted_json": optimized_json[:500] if len(optimized_json) > 500 else optimized_json
            }
        
        return {
            "success": True,
            "optimized_content": json.dumps(optimized_data, ensure_ascii=False, indent=2),
            "optimization_type": request.optimization_type,
            "raw_response": ai_response
        }
        
    except Exception as e:
        logger.error(f"AI优化大纲请求失败: {e}")
        return {
            "success": False,
            "error": str(e)
        }

@router.post("/api/ai/regenerate-image")
async def ai_regenerate_image(
    request: AIImageRegenerateRequest,
    user: User = Depends(get_current_user_required)
):
    """AI重新生成图像接口 - 完全遵循enhanced_ppt_service.py的标准流程"""
    try:
        # 获取图像服务和AI提供者
        from ..services.image.image_service import get_image_service

        image_service = get_image_service()
        if not image_service:
            return {
                "success": False,
                "message": "图像服务不可用"
            }

        provider, settings = get_role_provider("editor")
        # Ensure we have a general AI provider instance as well (some image processors expect ai_provider)
        ai_provider = get_ai_provider()
        if not ai_provider:
            return {
                "success": False,
                "message": "AI提供者不可用"
            }

        # 获取图像配置
        from ..services.config_service import config_service
        image_config = config_service.get_config_by_category('image_service')

        # 检查是否启用图片生成服务
        enable_image_service = image_config.get('enable_image_service', False)
        if not enable_image_service:
            return {
                "success": False,
                "message": "图片生成服务未启用，请在配置中启用"
            }

        # 第一步：检查启用的图片来源（完全遵循PPTImageProcessor的逻辑）
        from ..services.models.slide_image_info import ImageSource

        enabled_sources = []
        if image_config.get('enable_local_images', True):
            enabled_sources.append(ImageSource.LOCAL)
        if image_config.get('enable_network_search', False):
            enabled_sources.append(ImageSource.NETWORK)
        if image_config.get('enable_ai_generation', False):
            enabled_sources.append(ImageSource.AI_GENERATED)

        if not enabled_sources:
            return {
                "success": False,
                "message": "没有启用任何图片来源，请在配置中启用至少一种图片来源"
            }

        # 初始化PPT图像处理器
        from ..services.ppt_image_processor import PPTImageProcessor

        image_processor = PPTImageProcessor(
            image_service=image_service,
            ai_provider=ai_provider
        )

        # 提取图像信息和幻灯片内容
        image_info = request.image_info
        slide_content = request.slide_content

        # 构建幻灯片数据结构（遵循PPTImageProcessor期望的格式）
        slide_data = {
            'title': slide_content.get('title', ''),
            'content_points': [slide_content.get('title', '')],  # 简化的内容点
        }

        # 构建确认需求结构
        confirmed_requirements = {
            'project_topic': request.project_topic,
            'project_scenario': request.project_scenario
        }

        # 第二步：直接创建图像重新生成需求（跳过AI配图适用性判断）
        logger.info(f"开始图片重新生成，启用的来源: {[source.value for source in enabled_sources]}")

        # 分析原图像的用途和上下文
        image_context = await analyze_image_context(
            image_info, slide_content, request.project_topic, request.project_scenario
        )

        # 根据启用的来源和配置，智能选择最佳的图片来源
        selected_source = select_best_image_source(enabled_sources, image_config, image_context)

        # 创建图像需求对象（直接生成，不需要AI判断是否适合配图）
        from ..services.models.slide_image_info import ImageRequirement, ImagePurpose

        # 将字符串用途转换为ImagePurpose枚举
        purpose_str = image_context.get('image_purpose', 'illustration')
        purpose_mapping = {
            'background': ImagePurpose.BACKGROUND,
            'icon': ImagePurpose.ICON,
            'chart_support': ImagePurpose.CHART_SUPPORT,
            'decoration': ImagePurpose.DECORATION,
            'illustration': ImagePurpose.ILLUSTRATION
        }
        purpose = purpose_mapping.get(purpose_str, ImagePurpose.ILLUSTRATION)

        requirement = ImageRequirement(
            source=selected_source,
            count=1,
            purpose=purpose,
            description=f"重新生成图像: {image_info.get('alt', '')} - {request.project_topic}",
            priority=5  # 高优先级，因为是用户明确请求的重新生成
        )

        logger.info(f"选择图片来源: {selected_source.value}, 用途: {purpose.value}")

        # 第三步：直接处理图片生成（单个需求）
        from ..services.models.slide_image_info import SlideImagesCollection

        images_collection = SlideImagesCollection(page_number=request.slide_index + 1, images=[])

        # 根据选择的来源处理图片生成
        if requirement.source == ImageSource.LOCAL and ImageSource.LOCAL in enabled_sources:
            local_images = await image_processor._process_local_images(
                requirement, request.project_topic, request.project_scenario,
                slide_content.get('title', ''), slide_content.get('title', '')
            )
            images_collection.images.extend(local_images)

        elif requirement.source == ImageSource.NETWORK and ImageSource.NETWORK in enabled_sources:
            network_images = await image_processor._process_network_images(
                requirement, request.project_topic, request.project_scenario,
                slide_content.get('title', ''), slide_content.get('title', ''), image_config
            )
            images_collection.images.extend(network_images)

        elif requirement.source == ImageSource.AI_GENERATED and ImageSource.AI_GENERATED in enabled_sources:
            ai_images = await image_processor._process_ai_generated_images(
                requirement=requirement,
                project_topic=request.project_topic,
                project_scenario=request.project_scenario,
                slide_title=slide_content.get('title', ''),
                slide_content=slide_content.get('title', ''),
                image_config=image_config,
                page_number=request.slide_index + 1,
                total_pages=1,
                template_html=slide_content.get('html_content', '')
            )
            images_collection.images.extend(ai_images)

        # 重新计算统计信息
        images_collection.__post_init__()

        if images_collection.total_count == 0:
            return {
                "success": False,
                "message": "未能生成任何图片，请检查配置和网络连接"
            }

        # 获取第一张生成的图像（用于替换）
        new_image = images_collection.images[0]
        new_image_url = new_image.absolute_url

        # 替换HTML中的图像
        updated_html = replace_image_in_html(
            slide_content.get('html_content', ''),
            image_info,
            new_image_url
        )

        logger.info(f"图片重新生成成功: {new_image.source.value}来源, URL: {new_image_url}")

        return {
            "success": True,
            "message": f"图像重新生成成功（来源：{new_image.source.value}）",
            "new_image_url": new_image_url,
            "new_image_id": new_image.image_id,
            "updated_html_content": updated_html,
            "generation_prompt": getattr(new_image, 'generation_prompt', ''),
            "image_source": new_image.source.value,
            "ai_analysis": {
                "total_images_analyzed": 1,
                "reasoning": f"用户请求重新生成{image_context.get('image_purpose', '图像')}，选择{selected_source.value}来源",
                "enabled_sources": [source.value for source in enabled_sources],
                "selected_source": selected_source.value
            },
            "image_info": {
                "width": new_image.width,
                "height": new_image.height,
                "format": getattr(new_image, 'format', 'unknown'),
                "alt_text": new_image.alt_text,
                "title": new_image.title,
                "source": new_image.source.value,
                "purpose": new_image.purpose.value
            }
        }

    except Exception as e:
        logger.error(f"AI图像重新生成失败: {e}")
        import traceback
        traceback.print_exc()
        return {
            "success": False,
            "message": f"图像重新生成失败: {str(e)}"
        }

@router.post("/api/ai/auto-generate-slide-images")
async def ai_auto_generate_slide_images(
    request: AIAutoImageGenerateRequest,
    user: User = Depends(get_current_user_required)
):
    """AI一键配图接口 - 自动分析幻灯片内容并生成相关配图"""
    try:
        # 获取图像服务和AI提供者
        from ..services.image.image_service import get_image_service

        image_service = get_image_service()
        if not image_service:
            return {
                "success": False,
                "message": "图像服务不可用"
            }

        ai_provider = get_ai_provider()
        if not ai_provider:
            return {
                "success": False,
                "message": "AI提供者不可用"
            }

        # 获取图像处理器
        from ..services.ppt_image_processor import PPTImageProcessor
        image_processor = PPTImageProcessor(image_service, ai_provider)

        slide_content = request.slide_content
        slide_title = slide_content.get('title', f'第{request.slide_index + 1}页')
        slide_html = slide_content.get('html_content', '')

        logger.info(f"开始为第{request.slide_index + 1}页进行一键配图")

        # 第一步：AI分析幻灯片内容，确定是否需要配图以及配图需求
        analysis_prompt = f"""作为专业的PPT设计师，请分析以下幻灯片内容，判断是否需要配图以及配图需求。

项目主题：{request.project_topic}
项目场景：{request.project_scenario}
幻灯片标题：{slide_title}
幻灯片HTML内容：{slide_html[:1000]}...

请分析：
1. 这个幻灯片是否需要配图？
2. 如果需要，应该配几张图？
3. 每张图的用途和描述是什么？
4. 图片应该插入到什么位置？

请以JSON格式回复：
{{
    "needs_images": true/false,
    "image_count": 数量,
    "images": [
        {{
            "purpose": "图片用途（如：主要插图、装饰图、背景图等）",
            "description": "图片内容描述",
            "keywords": "搜索关键词",
            "position": "插入位置（如：标题下方、内容中间、页面右侧等）"
        }}
    ],
    "reasoning": "分析理由"
}}"""

        analysis_response = await ai_provider.text_completion(
            prompt=analysis_prompt,
            temperature=0.3
        )

        # 解析AI分析结果
        import json
        try:
            analysis_result = json.loads(analysis_response.content.strip())
        except json.JSONDecodeError:
            # 如果JSON解析失败，使用默认配置
            analysis_result = {
                "needs_images": True,
                "image_count": 1,
                "images": [{
                    "purpose": "主要插图",
                    "description": f"与{slide_title}相关的配图",
                    "keywords": f"{request.project_topic} {slide_title}",
                    "position": "内容中间"
                }],
                "reasoning": "默认为幻灯片添加一张主要配图"
            }

        if not analysis_result.get("needs_images", False):
            return {
                "success": True,
                "message": "AI分析认为此幻灯片不需要配图",
                "updated_html_content": slide_html,
                "generated_images_count": 0,
                "ai_analysis": analysis_result
            }

        # 第二步：根据分析结果生成图片需求
        from ..services.models.slide_image_info import ImageRequirement, ImagePurpose, ImageSource, SlideImagesCollection

        images_collection = SlideImagesCollection(page_number=request.slide_index + 1, images=[])

        # 获取图像配置（使用与重新生成图片相同的配置键）
        from ..services.config_service import config_service
        image_config = config_service.get_config_by_category('image_service')

        # 检查是否启用图片生成服务
        enable_image_service = image_config.get('enable_image_service', False)
        if not enable_image_service:
            return {
                "success": False,
                "message": "图片生成服务未启用，请在配置中启用"
            }

        # 获取启用的图像来源（使用与重新生成图片相同的逻辑）
        from ..services.models.slide_image_info import ImageSource

        enabled_sources = []
        if image_config.get('enable_local_images', True):
            enabled_sources.append(ImageSource.LOCAL)
        if image_config.get('enable_network_search', False):
            enabled_sources.append(ImageSource.NETWORK)
        if image_config.get('enable_ai_generation', False):
            enabled_sources.append(ImageSource.AI_GENERATED)

        if not enabled_sources:
            return {
                "success": False,
                "message": "没有启用的图像来源，请在设置中配置图像获取方式"
            }

        # 使用与重新生成图片完全相同的图片来源选择逻辑
        image_context = {
            'image_purpose': 'illustration',  # 一键配图默认为说明性图片
            'slide_title': slide_title,
            'slide_content': slide_html
        }

        selected_source = select_best_image_source(enabled_sources, image_config, image_context)

        # 为每个图片需求生成图片
        for i, image_info in enumerate(analysis_result.get("images", [])[:3]):  # 最多3张图
            # 创建图片需求
            requirement = ImageRequirement(
                purpose=ImagePurpose.ILLUSTRATION,
                description=image_info.get("description", "相关配图"),
                priority=1,
                source=selected_source,
                count=1
            )

            # 根据选择的来源处理图片生成
            if requirement.source == ImageSource.AI_GENERATED and ImageSource.AI_GENERATED in enabled_sources:
                ai_images = await image_processor._process_ai_generated_images(
                    requirement=requirement,
                    project_topic=request.project_topic,
                    project_scenario=request.project_scenario,
                    slide_title=slide_title,
                    slide_content=slide_title,
                    image_config=image_config,
                    page_number=request.slide_index + 1,
                    total_pages=1,
                    template_html=slide_html
                )
                images_collection.images.extend(ai_images)

            elif requirement.source == ImageSource.NETWORK and ImageSource.NETWORK in enabled_sources:
                network_images = await image_processor._process_network_images(
                    requirement=requirement,
                    project_topic=request.project_topic,
                    project_scenario=request.project_scenario,
                    slide_title=slide_title,
                    slide_content=slide_title,
                    image_config=image_config
                )
                images_collection.images.extend(network_images)

            elif requirement.source == ImageSource.LOCAL and ImageSource.LOCAL in enabled_sources:
                local_images = await image_processor._process_local_images(
                    requirement=requirement,
                    project_topic=request.project_topic,
                    project_scenario=request.project_scenario,
                    slide_title=slide_title,
                    slide_content=slide_title
                )
                images_collection.images.extend(local_images)

        if not images_collection.images:
            return {
                "success": False,
                "message": "未能生成任何配图，请检查图像服务配置"
            }

        # 第三步：将生成的图片插入到幻灯片中
        updated_html = await image_processor._insert_images_into_slide(
            slide_html, images_collection, slide_title
        )

        logger.info(f"一键配图完成: 生成{len(images_collection.images)}张图片")

        return {
            "success": True,
            "message": f"一键配图完成，已生成{len(images_collection.images)}张图片",
            "updated_html_content": updated_html,
            "generated_images_count": len(images_collection.images),
            "generated_images": [
                {
                    "image_id": img.image_id,
                    "url": img.absolute_url,
                    "description": img.content_description,
                    "source": img.source.value
                } for img in images_collection.images
            ],
            "ai_analysis": analysis_result
        }

    except Exception as e:
        logger.error(f"AI一键配图失败: {e}")
        return {
            "success": False,
            "message": f"一键配图失败: {str(e)}"
        }

@router.post("/api/ai/enhance-bullet-point")
async def ai_enhance_bullet_point(
    request: AIBulletPointEnhanceRequest,
    user: User = Depends(get_current_user_required)
):
    """AI增强要点接口"""
    try:
        # 获取AI提供者
        provider, settings = get_role_provider("outline")

        # 构建上下文信息
        context_info = ""
        if request.contextInfo:
            original_point = request.contextInfo.get('originalBulletPoint', '')
            other_points = request.contextInfo.get('otherBulletPoints', [])
            point_index = request.contextInfo.get('pointIndex', 0)

            context_info = f"""
当前要点上下文信息：
- 要点位置：第{point_index + 1}个要点
- 原始要点内容：{original_point}
- 同页面其他要点：{', '.join(other_points) if other_points else '无'}
"""

        # 构建大纲信息
        outline_info = ""
        if request.slideOutline:
            outline_info = f"""
当前幻灯片大纲信息：
- 幻灯片类型：{request.slideOutline.get('slide_type', '未知')}
- 描述：{request.slideOutline.get('description', '无')}
- 所有要点：{', '.join(request.slideOutline.get('content_points', [])) if request.slideOutline.get('content_points') else '无'}
"""

        # 构建AI增强提示词
        context = f"""
你是一位专业的PPT内容编辑专家。用户需要你增强和优化一个PPT要点的内容。

项目信息：
- 项目标题：{request.projectInfo.get('title', '未知')}
- 项目主题：{request.projectInfo.get('topic', '未知')}
- 应用场景：{request.projectInfo.get('scenario', '未知')}

幻灯片信息：
- 幻灯片标题：{request.slideTitle}
- 幻灯片位置：第{request.slideIndex}页

{outline_info}

{context_info}

用户请求：{request.userRequest}

请根据以上信息，对要点进行增强和优化。要求：

1. **保持核心意思不变**：不要改变要点的基本含义和方向
2. **增加具体细节**：添加更多具体的描述、数据、例子或说明
3. **提升表达质量**：使用更专业、更有吸引力的表达方式
4. **保持简洁性**：虽然要增强内容，但仍要保持要点的简洁特性，不要过于冗长
5. **与其他要点协调**：确保增强后的要点与同页面其他要点在风格和层次上保持一致
6. **符合场景需求**：根据应用场景调整语言风格和专业程度

请直接返回增强后的要点内容，不需要额外的解释或格式化。
"""

        # 调用AI生成增强内容
        response = await provider.text_completion(
            prompt=context,
            max_tokens=ai_config.max_tokens // 2,  # 使用较小的token限制
            temperature=ai_config.temperature,
            top_p=ai_config.top_p,
            model=settings.get('model')
        )

        enhanced_text = response.content.strip()

        # 简单的内容验证
        if not enhanced_text or len(enhanced_text) < 5:
            raise ValueError("AI生成的增强内容过短或为空")

        return {
            "success": True,
            "enhancedText": enhanced_text,
            "originalText": request.contextInfo.get('originalBulletPoint', '') if request.contextInfo else ""
        }

    except Exception as e:
        logger.error(f"AI要点增强请求失败: {e}")
        return {
            "success": False,
            "error": str(e),
            "message": "抱歉，AI要点增强服务暂时不可用。请稍后重试。"
        }

@router.post("/api/ai/enhance-all-bullet-points")
async def ai_enhance_all_bullet_points(
    request: AIBulletPointEnhanceRequest,
    user: User = Depends(get_current_user_required)
):
    """AI增强所有要点接口"""
    try:
        # 获取AI提供者
        provider, settings = get_role_provider("outline")

        # 构建上下文信息
        context_info = ""
        all_points = []
        if request.contextInfo:
            all_points = request.contextInfo.get('allBulletPoints', [])
            total_points = request.contextInfo.get('totalPoints', 0)

            context_info = f"""
当前要点上下文信息：
- 要点总数：{total_points}个
- 所有要点内容：
"""
            for i, point in enumerate(all_points, 1):
                context_info += f"  {i}. {point}\n"

        # 构建大纲信息
        outline_info = ""
        if request.slideOutline:
            outline_info = f"""
当前幻灯片大纲信息：
- 幻灯片类型：{request.slideOutline.get('slide_type', '未知')}
- 描述：{request.slideOutline.get('description', '无')}
"""

        # 构建AI增强提示词
        context = f"""
请对以下PPT要点进行增强和优化。

项目背景：
- 项目：{request.projectInfo.get('title', '未知')}
- 主题：{request.projectInfo.get('topic', '未知')}
- 场景：{request.projectInfo.get('scenario', '未知')}
- 幻灯片：{request.slideTitle}（第{request.slideIndex}页）

{outline_info}

{context_info}

增强要求：
1. 保持每个要点的核心意思不变
2. 添加具体细节、数据或例子
3. 使用更专业、准确的表达
4. 保持简洁，避免冗长
5. 确保要点间逻辑连贯、风格统一
6. 符合{request.projectInfo.get('scenario', '商务')}场景的专业要求

重要：请直接返回增强后的要点列表，每行一个要点，不要包含任何解释、开场白或格式说明。不要添加编号、符号或其他标记。

示例格式：
第一个增强后的要点内容
第二个增强后的要点内容
第三个增强后的要点内容
"""

        # 调用AI生成增强内容
        response = await provider.text_completion(
            prompt=context,
            max_tokens=ai_config.max_tokens,  # 使用完整的token限制，因为要处理多个要点
            temperature=ai_config.temperature,
            top_p=ai_config.top_p,
            model=settings.get('model')
        )

        enhanced_content = response.content.strip()

        # 解析增强后的要点 - 改进的过滤逻辑
        enhanced_points = []
        if enhanced_content:
            # 按行分割，过滤空行
            lines = [line.strip() for line in enhanced_content.split('\n') if line.strip()]

            # 过滤掉常见的无关内容
            filtered_lines = []
            skip_patterns = [
                '好的', '作为', '我将', '我会', '以下是', '根据', '请注意', '需要说明',
                '增强后的要点', '优化后的', '改进后的', '以上', '总结', '希望',
                '如有', '如果', '建议', '推荐', '注意', '提醒', '说明',
                '要点1', '要点2', '要点3', '要点4', '要点5',
                '第一', '第二', '第三', '第四', '第五', '第六', '第七', '第八', '第九', '第十',
                '1.', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '9.', '10.',
                '•', '·', '-', '*', '→', '▪', '▫'
            ]

            for line in lines:
                # 跳过过短的行（可能是格式标记）
                if len(line) < 5:
                    continue

                # 跳过包含常见开场白模式的行
                should_skip = False
                for pattern in skip_patterns:
                    if line.startswith(pattern) or (pattern in ['好的', '作为', '我将', '我会'] and pattern in line[:10]):
                        should_skip = True
                        break

                # 跳过纯数字或符号开头的行（可能是编号）
                if line[0].isdigit() or line[0] in ['•', '·', '-', '*', '→', '▪', '▫']:
                    # 但保留去掉编号后的内容
                    cleaned_line = line
                    # 移除开头的编号和符号
                    import re
                    cleaned_line = re.sub(r'^[\d\s\.\-\*\•\·\→\▪\▫]+', '', cleaned_line).strip()
                    if len(cleaned_line) >= 5:
                        filtered_lines.append(cleaned_line)
                    continue

                if not should_skip:
                    filtered_lines.append(line)

            enhanced_points = filtered_lines

        # 简单的内容验证
        if not enhanced_points or len(enhanced_points) == 0:
            raise ValueError("AI生成的增强内容为空或被过滤")

        return {
            "success": True,
            "enhancedPoints": enhanced_points,
            "originalPoints": all_points,
            "totalEnhanced": len(enhanced_points)
        }

    except Exception as e:
        logger.error(f"AI增强所有要点请求失败: {e}")
        return {
            "success": False,
            "error": str(e),
            "message": "抱歉，AI要点增强服务暂时不可用。请稍后重试。"
        }

@router.post("/api/projects/{project_id}/speech-script/generate")
async def generate_speech_script(
    project_id: str,
    request: SpeechScriptGenerationRequest,
    user: User = Depends(get_current_user_required)
):
    """Generate speech scripts for presentation slides"""
    try:
        import uuid
        import asyncio

        # Generate task ID for progress tracking
        task_id = str(uuid.uuid4())

        # Get project
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return {
                "success": False,
                "error": "Project not found"
            }

        # Check if slides data exists
        if not project.slides_data or len(project.slides_data) == 0:
            return {
                "success": False,
                "error": "No slides data available"
            }

        # Import speech script service
        from ..services.speech_script_service import SpeechScriptService, SpeechScriptCustomization
        from ..services.speech_script_service import SpeechTone, TargetAudience, LanguageComplexity

        # Initialize service
        speech_service = SpeechScriptService()

        # Parse customization options
        customization_data = request.customization
        customization = SpeechScriptCustomization(
            tone=SpeechTone(customization_data.get('tone', 'conversational')),
            target_audience=TargetAudience(customization_data.get('target_audience', 'general_public')),
            language_complexity=LanguageComplexity(customization_data.get('language_complexity', 'moderate')),
            custom_style_prompt=customization_data.get('custom_style_prompt'),
            include_transitions=customization_data.get('include_transitions', True),
            include_timing_notes=customization_data.get('include_timing_notes', False),
            speaking_pace=customization_data.get('speaking_pace', 'normal')
        )

        # Validate request parameters
        if request.generation_type == "single":
            if not request.slide_indices or len(request.slide_indices) != 1:
                return {
                    "success": False,
                    "error": "Single generation requires exactly one slide index"
                }
        elif request.generation_type == "multi":
            if not request.slide_indices:
                return {
                    "success": False,
                    "error": "Multi generation requires slide indices"
                }
        elif request.generation_type != "full":
            return {
                "success": False,
                "error": "Invalid generation type"
            }

        # Start async generation task
        async def generate_async():
            try:
                logger.info(f"Starting async generation for task {task_id}")

                # Generate scripts based on type
                if request.generation_type == "single":
                    # Use multi_slide_scripts_with_retry for single slide to get progress tracking
                    result = await speech_service.generate_multi_slide_scripts_with_retry(
                        project, request.slide_indices, customization, task_id=task_id
                    )
                elif request.generation_type == "multi":
                    result = await speech_service.generate_multi_slide_scripts_with_retry(
                        project, request.slide_indices, customization, task_id=task_id
                    )
                elif request.generation_type == "full":
                    result = await speech_service.generate_full_presentation_scripts(
                        project, customization, progress_callback=None, task_id=task_id
                    )

                # Save scripts to database if successful
                if result.success:
                    logger.info(f"Generation successful for task {task_id}, saving to database")
                    from ..services.speech_script_repository import SpeechScriptRepository
                    repo = SpeechScriptRepository()

                    generation_params = {
                        'generation_type': request.generation_type,
                        'tone': customization.tone.value,
                        'target_audience': customization.target_audience.value,
                        'language_complexity': customization.language_complexity.value,
                        'custom_audience': request.customization.get('custom_audience'),
                        'custom_style_prompt': customization.custom_style_prompt,
                        'include_transitions': customization.include_transitions,
                        'include_timing_notes': customization.include_timing_notes,
                        'speaking_pace': customization.speaking_pace
                    }

                    saved_count = 0
                    for script in result.scripts:
                        await repo.save_speech_script(
                            project_id=project_id,
                            slide_index=script.slide_index,
                            slide_title=script.slide_title,
                            script_content=script.script_content,
                            generation_params=generation_params,
                            estimated_duration=script.estimated_duration
                        )
                        saved_count += 1
                        logger.debug(f"Saved script {saved_count}/{len(result.scripts)} for slide {script.slide_index}")

                    # Ensure all changes are committed before closing
                    repo.db.commit()
                    repo.close()
                    logger.info(f"All {saved_count} scripts saved and committed to database for task {task_id}")

                    # NOW mark the task as completed after database save
                    from ..services.progress_tracker import progress_tracker
                    progress_tracker.complete_task(
                        task_id,
                        f"生成完成！成功 {saved_count} 页"
                    )
                    logger.info(f"Task {task_id} marked as completed")
                else:
                    logger.error(f"Generation failed for task {task_id}: {result.error_message}")

            except Exception as e:
                logger.error(f"Async speech script generation failed for task {task_id}: {e}")
                from ..services.progress_tracker import progress_tracker
                progress_tracker.fail_task(task_id, str(e))

        # Start the async task
        asyncio.create_task(generate_async())

        # Return immediately with task_id
        return {
            "success": True,
            "task_id": task_id,
            "message": "演讲稿生成已开始，请查看进度"
        }

    except Exception as e:
        logger.error(f"Speech script generation failed: {e}")
        return {
            "success": False,
            "error": str(e)
        }

@router.post("/api/projects/{project_id}/speech-script/export")
async def export_speech_script(
    project_id: str,
    request: SpeechScriptExportRequest,
    user: User = Depends(get_current_user_required)
):
    """Export speech scripts to document format"""
    try:
        # Get project for title
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return {
                "success": False,
                "error": "Project not found"
            }

        # Import exporter
        from ..services.speech_script_exporter import get_speech_script_exporter
        from ..services.speech_script_service import SlideScriptData

        exporter = get_speech_script_exporter()

        # Validate scripts data
        if not request.scripts_data or len(request.scripts_data) == 0:
            return {
                "success": False,
                "error": "No speech scripts data provided"
            }

        # Convert scripts data to SlideScriptData objects
        scripts = []
        for script_data in request.scripts_data:
            # Validate required fields
            if not script_data.get('script_content'):
                continue  # Skip empty scripts

            script = SlideScriptData(
                slide_index=script_data.get('slide_index', 0),
                slide_title=script_data.get('slide_title', ''),
                script_content=script_data.get('script_content', ''),
                estimated_duration=script_data.get('estimated_duration'),
                speaker_notes=script_data.get('speaker_notes')
            )
            scripts.append(script)

        # Check if we have any valid scripts after filtering
        if not scripts:
            return {
                "success": False,
                "error": "No valid speech scripts found"
            }

        # Prepare metadata
        metadata = {}
        if request.include_metadata:
            # Calculate total estimated duration from all scripts
            total_duration = None
            if scripts:
                total_minutes = 0
                for script in scripts:
                    if script.estimated_duration and '分钟' in script.estimated_duration:
                        try:
                            minutes = float(script.estimated_duration.replace('分钟', ''))
                            total_minutes += minutes
                        except ValueError:
                            pass
                if total_minutes > 0:
                    total_duration = f"{total_minutes:.1f}分钟"

            metadata = {
                'generation_time': time.time(),
                'total_estimated_duration': total_duration,
                'customization': {}
            }

        # Export based on format
        if request.export_format == "docx":
            if not exporter.is_docx_available():
                return {
                    "success": False,
                    "error": "DOCX export not available. Please install python-docx."
                }

            docx_content = await exporter.export_to_docx(
                scripts, project.topic, metadata
            )

            # Return file response
            import urllib.parse
            filename = f"{project.topic}_演讲稿.docx"
            safe_filename = urllib.parse.quote(filename, safe='')

            from fastapi.responses import Response
            return Response(
                content=docx_content,
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                headers={
                    "Content-Disposition": f"attachment; filename*=UTF-8''{safe_filename}"
                }
            )

        elif request.export_format == "markdown":
            markdown_content = await exporter.export_to_markdown(
                scripts, project.topic, metadata
            )

            # Return file response
            import urllib.parse
            filename = f"{project.topic}_演讲稿.md"
            safe_filename = urllib.parse.quote(filename, safe='')

            from fastapi.responses import Response
            return Response(
                content=markdown_content.encode('utf-8'),
                media_type="text/markdown",
                headers={
                    "Content-Disposition": f"attachment; filename*=UTF-8''{safe_filename}"
                }
            )

        else:
            return {
                "success": False,
                "error": "Unsupported export format"
            }

    except Exception as e:
        logger.error(f"Speech script export failed: {e}")
        return {
            "success": False,
            "error": str(e)
        }


@router.get("/api/projects/{project_id}/speech-scripts")
async def get_current_speech_scripts(
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """获取项目的当前演讲稿"""
    try:
        from ..services.speech_script_repository import SpeechScriptRepository

        # 检查项目是否存在
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return {
                "success": False,
                "error": "Project not found"
            }

        repo = SpeechScriptRepository()

        # Expire all objects to ensure fresh data from database
        repo.db.expire_all()

        # 获取项目的当前演讲稿
        scripts = await repo.get_current_speech_scripts_by_project(project_id)
        logger.info(f"Found {len(scripts)} speech scripts for project {project_id}")

        # 转换为JSON格式
        scripts_data = []
        for script in scripts:
            scripts_data.append({
                "id": script.id,
                "slide_index": script.slide_index,
                "slide_title": script.slide_title,
                "script_content": script.script_content,
                "estimated_duration": script.estimated_duration,
                "speaker_notes": script.speaker_notes,
                "generation_type": script.generation_type,
                "tone": script.tone,
                "target_audience": script.target_audience,
                "custom_audience": script.custom_audience,
                "language_complexity": script.language_complexity,
                "speaking_pace": script.speaking_pace,
                "custom_style_prompt": script.custom_style_prompt,
                "include_transitions": script.include_transitions,
                "include_timing_notes": script.include_timing_notes,
                "created_at": script.created_at,
                "updated_at": script.updated_at
            })

        repo.close()

        return {
            "success": True,
            "scripts": scripts_data
        }

    except Exception as e:
        logger.error(f"Get current speech scripts failed: {e}")
        return {
            "success": False,
            "error": str(e)
        }



@router.delete("/api/projects/{project_id}/speech-scripts/slide/{slide_index}")
async def delete_speech_script_by_slide(
    project_id: str,
    slide_index: int,
    user: User = Depends(get_current_user_required)
):
    """删除指定幻灯片的演讲稿"""
    try:
        from ..services.speech_script_repository import SpeechScriptRepository

        # 检查项目是否存在
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return {
                "success": False,
                "error": "Project not found"
            }

        repo = SpeechScriptRepository()

        # 获取并删除指定幻灯片的演讲稿
        script = await repo.get_speech_script_by_slide(project_id, slide_index)
        if not script:
            return {
                "success": False,
                "error": "Speech script not found"
            }

        success = await repo.delete_speech_script(script.id)

        return {
            "success": success,
            "message": f"第{slide_index + 1}页演讲稿已删除" if success else "删除演讲稿失败"
        }

    except Exception as e:
        logger.error(f"Delete speech script failed: {e}")
        return {
            "success": False,
            "error": str(e)
        }


@router.get("/api/projects/{project_id}/speech-scripts/result/{task_id}")
async def get_speech_script_result(
    project_id: str,
    task_id: str,
    user: User = Depends(get_current_user_required)
):
    """获取演讲稿生成结果"""
    try:
        from ..services.progress_tracker import progress_tracker
        from ..services.speech_script_repository import SpeechScriptRepository

        # 检查项目是否存在
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return {
                "success": False,
                "error": "Project not found"
            }

        # 获取进度信息
        progress_info = progress_tracker.get_progress(task_id)

        if not progress_info:
            return {
                "success": False,
                "error": "Task not found"
            }

        # 验证任务是否属于该项目
        if progress_info.project_id != project_id:
            return {
                "success": False,
                "error": "Access denied"
            }

        # 如果任务还未完成，返回进度信息
        if progress_info.status != "completed":
            return {
                "success": False,
                "error": "Task not completed yet",
                "status": progress_info.status,
                "progress": progress_info.to_dict()
            }

        # 获取生成的演讲稿
        repo = SpeechScriptRepository()
        scripts = await repo.get_current_speech_scripts_by_project(project_id)

        # 转换为API格式
        scripts_data = []
        total_duration_seconds = 0

        for script in scripts:
            script_data = {
                "slide_index": script.slide_index,
                "slide_title": script.slide_title,
                "script_content": script.script_content,
                "estimated_duration": script.estimated_duration,
                "speaker_notes": getattr(script, 'speaker_notes', None)
            }
            scripts_data.append(script_data)

            # 计算总时长
            if script.estimated_duration:
                try:
                    if '分钟' in script.estimated_duration:
                        minutes = float(script.estimated_duration.replace('分钟', ''))
                        total_duration_seconds += minutes * 60
                    elif '秒' in script.estimated_duration:
                        seconds = float(script.estimated_duration.replace('秒', ''))
                        total_duration_seconds += seconds
                except:
                    pass

        # 格式化总时长
        if total_duration_seconds < 60:
            total_duration = f"{int(total_duration_seconds)}秒"
        else:
            minutes = total_duration_seconds / 60
            total_duration = f"{minutes:.1f}分钟"

        repo.close()

        return {
            "success": True,
            "scripts": scripts_data,
            "total_estimated_duration": total_duration,
            "generation_metadata": {
                "task_id": task_id,
                "completed_at": progress_info.last_update,
                "total_slides": progress_info.total_slides,
                "completed_slides": progress_info.completed_slides,
                "failed_slides": progress_info.failed_slides,
                "skipped_slides": progress_info.skipped_slides
            }
        }

    except Exception as e:
        logger.error(f"Get speech script result failed: {e}")
        return {
            "success": False,
            "error": str(e)
        }


@router.get("/api/projects/{project_id}/speech-scripts/progress/{task_id}")
async def get_speech_script_progress(
    project_id: str,
    task_id: str,
    user: User = Depends(get_current_user_required)
):
    """获取演讲稿生成进度"""
    try:
        from ..services.progress_tracker import progress_tracker

        # 检查项目是否存在
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return {
                "success": False,
                "error": "Project not found"
            }

        # 获取进度信息
        progress_info = progress_tracker.get_progress(task_id)

        if not progress_info:
            return {
                "success": False,
                "error": "Task not found"
            }

        # 验证任务是否属于该项目
        if progress_info.project_id != project_id:
            return {
                "success": False,
                "error": "Access denied"
            }

        return {
            "success": True,
            "progress": progress_info.to_dict()
        }

    except Exception as e:
        logger.error(f"Get speech script progress failed: {e}")
        return {
            "success": False,
            "error": str(e)
        }


@router.put("/api/projects/{project_id}/speech-scripts/slide/{slide_index}")
async def update_speech_script_content(
    project_id: str,
    slide_index: int,
    request: dict,
    user: User = Depends(get_current_user_required)
):
    """更新演讲稿内容"""
    try:
        from ..services.speech_script_repository import SpeechScriptRepository

        # 检查项目是否存在
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return {
                "success": False,
                "error": "Project not found"
            }

        # 获取请求数据
        script_content = request.get('script_content', '').strip()
        slide_title = request.get('slide_title', f'第{slide_index + 1}页')
        estimated_duration = request.get('estimated_duration')
        speaker_notes = request.get('speaker_notes')

        if not script_content:
            return {
                "success": False,
                "error": "演讲稿内容不能为空"
            }

        repo = SpeechScriptRepository()

        # 获取现有演讲稿
        existing_script = await repo.get_speech_script_by_slide(project_id, slide_index)
        if not existing_script:
            return {
                "success": False,
                "error": "演讲稿不存在"
            }

        # 更新内容
        existing_script.script_content = script_content
        existing_script.slide_title = slide_title
        if estimated_duration:
            existing_script.estimated_duration = estimated_duration
        if speaker_notes is not None:
            existing_script.speaker_notes = speaker_notes
        existing_script.updated_at = time.time()

        repo.db.commit()
        repo.db.refresh(existing_script)
        repo.close()

        return {
            "success": True,
            "message": "演讲稿已更新",
            "script": {
                "id": existing_script.id,
                "slide_index": existing_script.slide_index,
                "slide_title": existing_script.slide_title,
                "script_content": existing_script.script_content,
                "estimated_duration": existing_script.estimated_duration,
                "speaker_notes": existing_script.speaker_notes,
                "updated_at": existing_script.updated_at
            }
        }

    except Exception as e:
        logger.error(f"Update speech script content failed: {e}")
        return {
            "success": False,
            "error": str(e)
        }


@router.get("/api/projects/{project_id}/selected-global-template")
async def get_selected_global_template(
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """获取项目选择的全局母版模板"""
    try:
        # 检查项目是否真正选择了模板
        selected_template = await ppt_service.get_selected_global_template(project_id)
        if selected_template:
            logger.info(f"Project {project_id} has selected template: {selected_template.get('template_name', 'Unknown')}")
            return {
                "status": "success",
                "template": selected_template,
                "is_user_selected": True
            }
        else:
            # 如果没有选择的模板，尝试获取默认模板
            default_template = await ppt_service.global_template_service.get_default_template()
            if default_template:
                logger.info(f"Project {project_id} using default template: {default_template.get('template_name', 'Unknown')}")
                return {
                    "status": "success",
                    "template": default_template,
                    "is_user_selected": False
                }
            else:
                logger.warning(f"No template available for project {project_id}")
                return {
                    "status": "success",
                    "template": None,
                    "is_user_selected": False
                }
    except Exception as e:
        logger.error(f"Error getting selected global template for project {project_id}: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/api/projects/{project_id}/free-template")
async def get_project_free_template(
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """Get project's free-template status and current generated template (if any)."""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        metadata = project.project_metadata or {}
        if metadata.get("template_mode") != "free":
            return {
                "success": True,
                "enabled": False,
                "message": "Project is not using free template mode",
                "status": None,
                "confirmed": False,
                "template": None
            }

        html = metadata.get("free_template_html")
        name = metadata.get("free_template_name") or "自由模板（AI决定）"

        template = None
        if isinstance(html, str) and html.strip():
            template = {
                "template_name": name,
                "description": "AI 根据大纲自动生成的项目专属模板",
                "html_template": html,
                "tags": ["自由模板", "AI生成", "项目专属"],
                "created_by": "ai_free"
            }

        return {
            "success": True,
            "enabled": True,
            "status": metadata.get("free_template_status"),
            "confirmed": bool(metadata.get("free_template_confirmed")),
            "saved_template_id": metadata.get("saved_global_template_id"),
            "template": template
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error getting free template for project {project_id}: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/api/projects/{project_id}/free-template/generate")
async def generate_project_free_template(
    project_id: str,
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Generate (or regenerate) a project's free-template via AI."""
    try:
        payload = {}
        try:
            payload = await request.json()
        except Exception:
            payload = {}

        force = bool(payload.get("force", False))

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        metadata = project.project_metadata or {}
        if metadata.get("template_mode") != "free":
            raise HTTPException(status_code=400, detail="Project is not using free template mode")

        if force:
            metadata.pop("free_template_html", None)
            metadata.pop("free_template_name", None)
            metadata.pop("free_template_generated_at", None)
            metadata.pop("free_template_prompt", None)
            metadata["free_template_status"] = "pending"
            metadata["free_template_confirmed"] = False
            metadata.pop("free_template_confirmed_at", None)
            await ppt_service.project_manager.update_project_metadata(project_id, metadata)
            ppt_service.clear_cached_style_genes(project_id)

        template = await ppt_service.get_selected_global_template(project_id)
        if not template:
            raise HTTPException(status_code=500, detail="Failed to generate free template")

        # Mark status ready (generation is synchronous here)
        project = await ppt_service.project_manager.get_project(project_id)
        if project and project.project_metadata:
            metadata = project.project_metadata
            metadata["free_template_status"] = "ready"
            await ppt_service.project_manager.update_project_metadata(project_id, metadata)

        return {
            "success": True,
            "template": template
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error generating free template for project {project_id}: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/api/projects/{project_id}/free-template/confirm")
async def confirm_project_free_template(
    project_id: str,
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Confirm using the generated free-template; optionally save it into global template list."""
    try:
        data = await request.json()
        save_to_library = bool(data.get("save_to_library", False))
        requested_name = (data.get("template_name") or "").strip()
        requested_description = (data.get("description") or "").strip()
        requested_tags = data.get("tags") or []

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        metadata = project.project_metadata or {}
        if metadata.get("template_mode") != "free":
            raise HTTPException(status_code=400, detail="Project is not using free template mode")

        html = metadata.get("free_template_html")
        if not (isinstance(html, str) and html.strip()):
            raise HTTPException(status_code=400, detail="Free template is not generated yet")

        metadata["free_template_confirmed"] = True
        metadata["free_template_confirmed_at"] = time.time()
        metadata["free_template_status"] = "ready"

        saved_template = None
        if save_to_library:
            base_name = requested_name or f"自由模板-{(project.topic or 'PPT')[:20]}-{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            description = requested_description or "由自由模板功能生成并确认的模板"
            tags: List[str] = []
            if isinstance(requested_tags, list):
                tags = [str(t).strip() for t in requested_tags if str(t).strip()]
            tags = tags or ["自由模板", "AI生成"]

            # Ensure unique name
            final_name = base_name
            for i in range(1, 6):
                try:
                    saved_template = await ppt_service.global_template_service.create_template({
                        "template_name": final_name,
                        "description": description,
                        "html_template": html,
                        "tags": tags,
                        "is_default": False,
                        "is_active": True,
                        "created_by": f"free_template:{project_id}"
                    })
                    break
                except ValueError:
                    final_name = f"{base_name}-{i}"

            if not saved_template:
                raise HTTPException(status_code=409, detail="Failed to save template to library (name conflict)")

            metadata["saved_global_template_id"] = saved_template.get("id")
            metadata["saved_global_template_name"] = saved_template.get("template_name")

        await ppt_service.project_manager.update_project_metadata(project_id, metadata)
        ppt_service.clear_cached_style_genes(project_id)

        return {
            "success": True,
            "saved_template": saved_template
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error confirming free template for project {project_id}: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/api/projects/{project_id}/free-template/adjust")
async def adjust_project_free_template(
    project_id: str,
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Adjust the generated free-template based on user feedback."""
    try:
        data = await request.json()
        adjustment_request = (data.get("adjustment_request") or "").strip()
        
        if not adjustment_request:
            raise HTTPException(status_code=400, detail="Adjustment request is required")
        
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")
        
        metadata = project.project_metadata or {}
        if metadata.get("template_mode") != "free":
            raise HTTPException(status_code=400, detail="Project is not using free template mode")
        
        current_html = metadata.get("free_template_html")
        if not (isinstance(current_html, str) and current_html.strip()):
            raise HTTPException(status_code=400, detail="Free template is not generated yet")
        
        template_name = metadata.get("free_template_name") or "自由模板"
        
        # Use the global template service to adjust the template
        adjusted_html = None
        async for chunk in ppt_service.global_template_service.adjust_template_with_ai_stream(
            current_html=current_html,
            adjustment_request=adjustment_request,
            template_name=template_name
        ):
            if chunk.get('type') == 'complete':
                adjusted_html = chunk.get('html_template')
                break
            elif chunk.get('type') == 'error':
                raise HTTPException(status_code=500, detail=chunk.get('message', 'Template adjustment failed'))
        
        if not adjusted_html:
            raise HTTPException(status_code=500, detail="Failed to adjust template")
        
        # Update project metadata with adjusted template
        metadata["free_template_html"] = adjusted_html
        metadata["free_template_adjusted_at"] = time.time()
        metadata["free_template_adjustment_request"] = adjustment_request
        metadata["free_template_confirmed"] = False  # Reset confirmation after adjustment
        
        await ppt_service.project_manager.update_project_metadata(project_id, metadata)
        ppt_service.clear_cached_style_genes(project_id)
        
        return {
            "success": True,
            "template": {
                "template_name": template_name,
                "html_template": adjusted_html,
                "description": "AI 根据用户建议调整后的模板"
            }
        }
    except HTTPException:
        raise
async def generate_project_free_template(
    project_id: str,
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Generate (or regenerate) a project's free-template via AI."""
    try:
        payload = {}
        try:
            payload = await request.json()
        except Exception:
            payload = {}

        force = bool(payload.get("force", False))

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        metadata = project.project_metadata or {}
        if metadata.get("template_mode") != "free":
            raise HTTPException(status_code=400, detail="Project is not using free template mode")

        if force:
            metadata.pop("free_template_html", None)
            metadata.pop("free_template_name", None)
            metadata.pop("free_template_generated_at", None)
            metadata.pop("free_template_prompt", None)
            metadata["free_template_status"] = "pending"
            metadata["free_template_confirmed"] = False
            metadata.pop("free_template_confirmed_at", None)
            await ppt_service.project_manager.update_project_metadata(project_id, metadata)
            ppt_service.clear_cached_style_genes(project_id)

        template = await ppt_service.get_selected_global_template(project_id)
        if not template:
            raise HTTPException(status_code=500, detail="Failed to generate free template")

        # Mark status ready (generation is synchronous here)
        project = await ppt_service.project_manager.get_project(project_id)
        if project and project.project_metadata:
            metadata = project.project_metadata
            metadata["free_template_status"] = "ready"
            await ppt_service.project_manager.update_project_metadata(project_id, metadata)

        return {
            "success": True,
            "template": template
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error generating free template for project {project_id}: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/api/projects/{project_id}/free-template/confirm")
async def confirm_project_free_template(
    project_id: str,
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Confirm using the generated free-template; optionally save it into global template list."""
    try:
        data = await request.json()
        save_to_library = bool(data.get("save_to_library", False))
        requested_name = (data.get("template_name") or "").strip()
        requested_description = (data.get("description") or "").strip()
        requested_tags = data.get("tags") or []

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        metadata = project.project_metadata or {}
        if metadata.get("template_mode") != "free":
            raise HTTPException(status_code=400, detail="Project is not using free template mode")

        html = metadata.get("free_template_html")
        if not (isinstance(html, str) and html.strip()):
            raise HTTPException(status_code=400, detail="Free template is not generated yet")

        metadata["free_template_confirmed"] = True
        metadata["free_template_confirmed_at"] = time.time()
        metadata["free_template_status"] = "ready"

        saved_template = None
        if save_to_library:
            base_name = requested_name or f"自由模板-{(project.topic or 'PPT')[:20]}-{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            description = requested_description or "由自由模板功能生成并确认的模板"
            tags: List[str] = []
            if isinstance(requested_tags, list):
                tags = [str(t).strip() for t in requested_tags if str(t).strip()]
            tags = tags or ["自由模板", "AI生成"]

            # Ensure unique name
            final_name = base_name
            for i in range(1, 6):
                try:
                    saved_template = await ppt_service.global_template_service.create_template({
                        "template_name": final_name,
                        "description": description,
                        "html_template": html,
                        "tags": tags,
                        "is_default": False,
                        "is_active": True,
                        "created_by": f"free_template:{project_id}"
                    })
                    break
                except ValueError:
                    final_name = f"{base_name}-{i}"

            if not saved_template:
                raise HTTPException(status_code=409, detail="Failed to save template to library (name conflict)")

            metadata["saved_global_template_id"] = saved_template.get("id")
            metadata["saved_global_template_name"] = saved_template.get("template_name")

        await ppt_service.project_manager.update_project_metadata(project_id, metadata)
        ppt_service.clear_cached_style_genes(project_id)

        return {
            "success": True,
            "saved_template": saved_template
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error confirming free template for project {project_id}: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/api/projects/{project_id}/free-template/adjust")
async def adjust_project_free_template(
    project_id: str,
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Adjust the generated free-template based on user feedback."""
    try:
        data = await request.json()
        adjustment_request = (data.get("adjustment_request") or "").strip()
        
        if not adjustment_request:
            raise HTTPException(status_code=400, detail="Adjustment request is required")
        
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")
        
        metadata = project.project_metadata or {}
        if metadata.get("template_mode") != "free":
            raise HTTPException(status_code=400, detail="Project is not using free template mode")
        
        current_html = metadata.get("free_template_html")
        if not (isinstance(current_html, str) and current_html.strip()):
            raise HTTPException(status_code=400, detail="Free template is not generated yet")
        
        template_name = metadata.get("free_template_name") or "自由模板"
        
        # Use the global template service to adjust the template
        adjusted_html = None
        async for chunk in ppt_service.global_template_service.adjust_template_with_ai_stream(
            current_html=current_html,
            adjustment_request=adjustment_request,
            template_name=template_name
        ):
            if chunk.get('type') == 'complete':
                adjusted_html = chunk.get('html_template')
                break
            elif chunk.get('type') == 'error':
                raise HTTPException(status_code=500, detail=chunk.get('message', 'Template adjustment failed'))
        
        if not adjusted_html:
            raise HTTPException(status_code=500, detail="Failed to adjust template")
        
        # Update project metadata with adjusted template
        metadata["free_template_html"] = adjusted_html
        metadata["free_template_adjusted_at"] = time.time()
        metadata["free_template_adjustment_request"] = adjustment_request
        metadata["free_template_confirmed"] = False  # Reset confirmation after adjustment
        
        await ppt_service.project_manager.update_project_metadata(project_id, metadata)
        ppt_service.clear_cached_style_genes(project_id)
        
        return {
            "success": True,
            "template": {
                "template_name": template_name,
                "html_template": adjusted_html,
                "description": "AI 根据用户建议调整后的模板"
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error adjusting free template for project {project_id}: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/api/projects/{project_id}/slides/{slide_index}/save")
async def save_single_slide_content(
    project_id: str,
    slide_index: int,
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """保存单个幻灯片内容到数据库
    
    重要：此函数只保存被编辑的单个幻灯片，不会触碰其他幻灯片数据，
    以避免与正在进行的PPT生成过程产生冲突。
    """
    try:
        logger.info(f"🔄 开始保存项目 {project_id} 的第 {slide_index + 1} 页 (索引: {slide_index})")

        data = await request.json()
        html_content = data.get('html_content', '')
        raw_is_user_edited = data.get('is_user_edited', None)

        logger.info(f"📄 接收到HTML内容，长度: {len(html_content)} 字符")

        if not html_content:
            logger.error("❌ HTML内容为空")
            raise HTTPException(status_code=400, detail="HTML content is required")

        if slide_index < 0:
            logger.error(f"❌ 幻灯片索引不能为负数: {slide_index}")
            raise HTTPException(status_code=400, detail=f"Slide index cannot be negative: {slide_index}")

        # 直接从数据库获取该幻灯片的当前数据
        from ..services.db_project_manager import DatabaseProjectManager
        db_manager = DatabaseProjectManager()
        
        # 获取项目基本信息确认项目存在
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            logger.error(f"❌ 项目 {project_id} 不存在")
            raise HTTPException(status_code=404, detail="Project not found")

        # 获取当前幻灯片数据（如果存在）
        existing_slide = await db_manager.get_single_slide(project_id, slide_index)
        
        def _parse_bool(value):
            if isinstance(value, bool):
                return value
            if isinstance(value, str):
                normalized = value.strip().lower()
                if normalized in ("true", "1", "yes", "y", "on"):
                    return True
                if normalized in ("false", "0", "no", "n", "off"):
                    return False
            return None

        client_is_user_edited = _parse_bool(raw_is_user_edited)
        html_changed = True
        if existing_slide:
            html_changed = (existing_slide.get("html_content") != html_content)

        # is_user_edited 优先使用前端传入值；否则：保留既有标记，若HTML有变化则标记为True
        if client_is_user_edited is None:
            is_user_edited = (existing_slide.get("is_user_edited", False) if existing_slide else False) or html_changed
        else:
            is_user_edited = client_is_user_edited

        # 构建要保存的幻灯片数据
        # 保留现有数据的其他字段，只更新html_content和is_user_edited
        if existing_slide:
            slide_data = existing_slide.copy()
            slide_data['html_content'] = html_content
            slide_data['is_user_edited'] = is_user_edited
            # 兼容字段：DatabaseService.save_single_slide 使用 content_type
            if 'content_type' not in slide_data:
                slide_data['content_type'] = slide_data.get('slide_type', 'content')
        else:
            # 如果幻灯片不存在（理论上不应该发生，但做防御处理）
            slide_data = {
                "page_number": slide_index + 1,
                "title": f"Slide {slide_index + 1}",
                "html_content": html_content,
                "content_type": "content",
                "is_user_edited": is_user_edited
            }

        logger.debug(f"📝 更新第 {slide_index + 1} 页的内容")
        logger.debug(f"📊 幻灯片数据: 标题='{slide_data.get('title', '无标题')}', 用户编辑={is_user_edited}, 索引={slide_index}")

        # 只保存这一个幻灯片到数据库，不影响其他幻灯片
        save_success = await db_manager.save_single_slide(project_id, slide_index, slide_data)

        if save_success:
            logger.debug(f"✅ 第 {slide_index + 1} 页已成功保存到数据库")

            return {
                "success": True,
                "message": f"Slide {slide_index + 1} saved successfully to database",
                "slide_data": slide_data,
                "database_saved": True
            }
        else:
            logger.error(f"❌ 保存第 {slide_index + 1} 页到数据库失败")
            return {
                "success": False,
                "error": "Failed to save slide to database",
                "database_saved": False
            }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ 保存单个幻灯片时发生错误: {e}")
        import traceback
        traceback.print_exc()
        return {
            "success": False,
            "error": str(e),
            "database_saved": False
        }

@router.get("/api/projects/{project_id}/slides/stream")
async def stream_slides_generation(project_id: str):
    """Stream slides generation process"""
    try:
        # Guard: free-template must be confirmed before starting generation
        try:
            project = await ppt_service.project_manager.get_project(project_id)
            if project and project.project_metadata:
                metadata = project.project_metadata or {}
                if metadata.get("template_mode") == "free" and not metadata.get("free_template_confirmed"):
                    async def blocked_stream():
                        yield f"data: {json.dumps({'type': 'error', 'message': '自由模板尚未确认，请先在预览中确认/保存模板后再开始生成PPT。'})}\n\n"
                    return StreamingResponse(
                        blocked_stream(),
                        media_type="text/event-stream",
                        headers={
                            "Cache-Control": "no-cache",
                            "Connection": "keep-alive",
                            "Access-Control-Allow-Origin": "*",
                            "Access-Control-Allow-Headers": "Cache-Control"
                        }
                    )
        except Exception:
            # If guard fails, do not block generation
            pass

        async def generate_slides_stream():
            async for chunk in ppt_service.generate_slides_streaming(project_id):
                yield chunk

        return StreamingResponse(
            generate_slides_stream(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Headers": "Cache-Control"
            }
        )

    except Exception as e:
        return {"error": str(e)}


@router.post("/api/projects/{project_id}/slides/cleanup")
async def cleanup_excess_slides(
    project_id: str,
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """清理项目中多余的幻灯片"""
    try:
        logger.info(f"🧹 开始清理项目 {project_id} 的多余幻灯片")

        data = await request.json()
        current_slide_count = data.get('current_slide_count', 0)

        if current_slide_count <= 0:
            logger.error("❌ 无效的幻灯片数量")
            raise HTTPException(status_code=400, detail="Invalid slide count")

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            logger.error(f"❌ 项目 {project_id} 不存在")
            raise HTTPException(status_code=404, detail="Project not found")

        # 清理数据库中多余的幻灯片
        from ..services.db_project_manager import DatabaseProjectManager
        db_manager = DatabaseProjectManager()
        deleted_count = await db_manager.cleanup_excess_slides(project_id, current_slide_count)

        logger.info(f"✅ 项目 {project_id} 清理完成，删除了 {deleted_count} 张多余的幻灯片")

        return {
            "success": True,
            "message": f"Successfully cleaned up {deleted_count} excess slides",
            "deleted_count": deleted_count
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ 清理幻灯片失败: {e}")
        return {"success": False, "error": str(e)}


@router.post("/api/projects/{project_id}/slides/batch-save")
async def batch_save_slides(
    project_id: str,
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """批量保存所有幻灯片 - 高效版本"""
    try:
        logger.debug(f"🔄 开始批量保存项目 {project_id} 的所有幻灯片")

        data = await request.json()
        slides_data = data.get('slides_data', [])

        if not slides_data:
            logger.error("❌ 幻灯片数据为空")
            raise HTTPException(status_code=400, detail="Slides data is required")

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            logger.error(f"❌ 项目 {project_id} 不存在")
            raise HTTPException(status_code=404, detail="Project not found")

        # 更新项目内存中的数据
        project.slides_data = slides_data
        project.updated_at = time.time()

        # 重新生成完整HTML
        outline_title = project.title
        if hasattr(project, 'outline') and project.outline:
            outline_title = project.outline.get('title', project.title)

        project.slides_html = ppt_service._combine_slides_to_full_html(
            project.slides_data, outline_title
        )

        # 使用批量保存到数据库
        from ..services.db_project_manager import DatabaseProjectManager
        db_manager = DatabaseProjectManager()

        # 批量保存幻灯片
        batch_success = await db_manager.batch_save_slides(project_id, slides_data)

        # 更新项目信息
        if batch_success:
            await db_manager.update_project_data(project_id, {
                "slides_html": project.slides_html,
                "slides_data": project.slides_data,
                "updated_at": project.updated_at
            })

        logger.debug(f"✅ 项目 {project_id} 批量保存完成，共 {len(slides_data)} 张幻灯片")

        return {
            "success": batch_success,
            "message": f"Successfully batch saved {len(slides_data)} slides" if batch_success else "Batch save failed",
            "slides_count": len(slides_data)
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"❌ 批量保存幻灯片失败: {e}")
        return {"success": False, "error": str(e)}


@router.get("/api/projects/{project_id}/export/pdf")
async def export_project_pdf(project_id: str, individual: bool = False):
    """Export project as PDF using Pyppeteer"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if we have slides data
        if not project.slides_data or len(project.slides_data) == 0:
            raise HTTPException(status_code=400, detail="PPT not generated yet")

        # Check if Pyppeteer is available
        pdf_converter = get_pdf_converter()
        if not pdf_converter.is_available():
            raise HTTPException(
                status_code=503,
                detail="PDF generation service unavailable. Please ensure Pyppeteer is installed: pip install pyppeteer"
            )

        # Create temp file in thread pool to avoid blocking
        temp_pdf_path = await run_blocking_io(
            lambda: tempfile.NamedTemporaryFile(suffix='.pdf', delete=False).name
        )

        logging.info("Generating PDF with Pyppeteer")
        success = await _generate_pdf_with_pyppeteer(project, temp_pdf_path, individual)

        if not success:
            # Clean up temp file and raise error
            await run_blocking_io(lambda: os.unlink(temp_pdf_path) if os.path.exists(temp_pdf_path) else None)
            raise HTTPException(status_code=500, detail="PDF generation failed")

        # Return PDF file
        logging.info("PDF generated successfully using Pyppeteer")
        safe_filename = urllib.parse.quote(f"{project.topic}_PPT.pdf", safe='')

        # 使用BackgroundTask来清理临时文件
        from starlette.background import BackgroundTask

        def cleanup_temp_file():
            try:
                os.unlink(temp_pdf_path)
            except:
                pass

        return FileResponse(
            temp_pdf_path,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f"attachment; filename*=UTF-8''{safe_filename}",
                "X-PDF-Generator": "Pyppeteer"
            },
            background=BackgroundTask(cleanup_temp_file)
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/api/projects/{project_id}/slides/{slide_index}/user-edited")
async def set_slide_user_edited_status(
    project_id: str,
    slide_index: int,
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """更新单个幻灯片的 is_user_edited 状态（用于修复误标记导致生成跳过保存的问题）"""
    try:
        data = await request.json()

        def _parse_bool(value):
            if isinstance(value, bool):
                return value
            if isinstance(value, str):
                normalized = value.strip().lower()
                if normalized in ("true", "1", "yes", "y", "on"):
                    return True
                if normalized in ("false", "0", "no", "n", "off"):
                    return False
            return None

        is_user_edited = _parse_bool(data.get("is_user_edited", None))
        if is_user_edited is None:
            raise HTTPException(status_code=400, detail="is_user_edited must be a boolean")

        # Ensure project exists
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        if slide_index < 0:
            raise HTTPException(status_code=400, detail=f"Slide index cannot be negative: {slide_index}")

        from ..database.database import AsyncSessionLocal
        from ..database.repositories import SlideDataRepository, ProjectRepository

        session = AsyncSessionLocal()
        try:
            slide_repo = SlideDataRepository(session)
            updated = await slide_repo.update_slide_user_edited_status(project_id, slide_index, is_user_edited=is_user_edited)

            # Fallback: if slide_data table没有该页，尝试更新projects.slides_data字段
            if not updated:
                project_repo = ProjectRepository(session)
                db_project = await project_repo.get_by_id(project_id)
                if db_project and db_project.slides_data and slide_index < len(db_project.slides_data):
                    db_project.slides_data[slide_index]["is_user_edited"] = is_user_edited
                    await project_repo.update(project_id, {"slides_data": db_project.slides_data})
                    updated = True

            return {"success": bool(updated), "is_user_edited": is_user_edited}
        finally:
            await session.close()

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error setting slide user edited status for project {project_id} slide {slide_index}: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/api/projects/{project_id}/export/pdf/individual")
async def export_project_pdf_individual(project_id: str):
    """Export project as individual PDF files for each slide"""
    return await export_project_pdf(project_id, individual=True)

@router.get("/api/projects/{project_id}/export/pptx")
async def export_project_pptx(project_id: str):
    """Export project as PPTX by first generating PDF then converting to PowerPoint"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if we have slides data
        if not project.slides_data or len(project.slides_data) == 0:
            raise HTTPException(status_code=400, detail="PPT not generated yet")

        # Get PDF to PPTX converter
        converter = get_pdf_to_pptx_converter()
        if not converter.is_available():
            raise HTTPException(
                status_code=503,
                detail="PPTX conversion service unavailable. Please ensure Apryse SDK is installed and licensed."
            )

        # Check if Pyppeteer is available for PDF generation
        pdf_converter = get_pdf_converter()
        if not pdf_converter.is_available():
            raise HTTPException(
                status_code=503,
                detail="PDF generation service unavailable. Please ensure Pyppeteer is installed: pip install pyppeteer"
            )

        # Step 1: Generate PDF using existing PDF export functionality
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf_file:
            temp_pdf_path = temp_pdf_file.name

        logging.info("Step 1: Generating PDF for PPTX conversion")
        pdf_success = await _generate_pdf_with_pyppeteer(project, temp_pdf_path, individual=False)

        if not pdf_success:
            # Clean up temp file and raise error
            try:
                os.unlink(temp_pdf_path)
            except:
                pass
            raise HTTPException(status_code=500, detail="PDF generation failed")

        # Step 2: 启动PDF转PPTX后台任务
        logging.info("Step 2: Starting PDF to PPTX conversion task")

        from ..services.background_tasks import get_task_manager, TaskStatus

        task_manager = get_task_manager()

        # 创建临时PPTX文件路径
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_pptx_file:
            temp_pptx_path = temp_pptx_file.name

        # 定义转换任务函数
        async def pdf_to_pptx_task():
            """PDF to PPTX conversion task (runs in subprocess)."""
            try:
                success, result = await converter.convert_pdf_to_pptx_async(
                    temp_pdf_path,
                    temp_pptx_path
                )
                if success:
                    # 转换成功后，添加演讲稿到备注
                    try:
                        from pptx import Presentation
                        from ..services.speech_script_repository import SpeechScriptRepository

                        # 获取演讲稿数据
                        repo = SpeechScriptRepository()
                        scripts_list = await repo.get_current_speech_scripts_by_project(project_id)
                        speech_scripts = {script.slide_index: script.script_content for script in scripts_list}
                        repo.close()

                        if len(speech_scripts) > 0:
                            # 打开生成的PPTX文件
                            prs = Presentation(temp_pptx_path)

                            # 为每张幻灯片添加演讲稿备注
                            for i, slide in enumerate(prs.slides):
                                if i in speech_scripts:
                                    notes_slide = slide.notes_slide
                                    text_frame = notes_slide.notes_text_frame
                                    text_frame.text = speech_scripts[i]
                                    logging.info(f"Added speech script to slide {i+1} notes")

                            # 保存修改后的PPTX
                            prs.save(temp_pptx_path)
                            logging.info(f"Added {len(speech_scripts)} speech scripts to PPTX notes")
                    except Exception as e:
                        logging.warning(f"Failed to add speech scripts to PPTX: {e}")
                        # 继续执行，即使添加演讲稿失败也返回PPTX

                    return {
                        "success": True,
                        "pptx_path": temp_pptx_path,
                        "pdf_path": temp_pdf_path
                    }
                else:
                    return {
                        "success": False,
                        "error": result
                    }
            except Exception as e:
                return {
                    "success": False,
                    "error": str(e)
                }


        # 提交后台任务
        task_id = task_manager.submit_task(
            task_type="pdf_to_pptx_conversion",
            func=pdf_to_pptx_task,
            metadata={
                "project_id": project_id,
                "project_topic": project.topic,
                "pdf_path": temp_pdf_path,
                "pptx_path": temp_pptx_path
            }
        )

        # 立即返回任务ID，不等待任务完成
        return JSONResponse({
            "status": "processing",
            "task_id": task_id,
            "message": "PPTX conversion started in background",
            "polling_endpoint": f"/api/landppt/tasks/{task_id}"
        })

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/api/projects/{project_id}/export/pptx-images")
async def export_project_pptx_from_images(project_id: str, request: ImagePPTXExportRequest):
    """Export project as PPTX using high-quality Playwright screenshots"""
    try:
        from io import BytesIO
        from pptx import Presentation
        from pptx.util import Inches

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # 验证是否有幻灯片数据
        slides = getattr(request, 'slides', None)
        if not slides or len(slides) == 0:
            raise HTTPException(status_code=400, detail="No slides provided")

        # 检查Playwright是否可用
        pdf_converter = get_pdf_converter()
        if not pdf_converter.is_available():
            raise HTTPException(
                status_code=503,
                detail="Screenshot service unavailable. Please ensure Playwright is installed."
            )

        # 创建后台任务
        from ..services.background_tasks import get_task_manager
        task_manager = get_task_manager()

        # 创建临时目录和PPTX文件路径
        temp_dir = tempfile.mkdtemp()
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_pptx_file:
            temp_pptx_path = temp_pptx_file.name

        # 定义HTML到图片到PPTX的任务函数
        async def html_to_pptx_task():
            """使用Playwright截图并生成PPTX"""
            screenshot_paths = []
            try:
                logging.info(f"Starting screenshot-based PPTX export for {len(slides)} slides")

                # 第1步：获取演讲稿数据
                speech_scripts = {}
                try:
                    from ..services.speech_script_repository import SpeechScriptRepository
                    repo = SpeechScriptRepository()
                    scripts_list = await repo.get_current_speech_scripts_by_project(project_id)
                    # 构建幻灯片索引到演讲稿的映射
                    for script in scripts_list:
                        speech_scripts[script.slide_index] = script.script_content
                    repo.close()
                    logging.info(f"Loaded {len(speech_scripts)} speech scripts for slides")
                except Exception as e:
                    logging.warning(f"Failed to load speech scripts: {e}")
                    # 继续执行，即使没有演讲稿也可以生成PPTX

                # 第2步：为每张幻灯片创建临时HTML文件
                html_files = []
                for i, slide in enumerate(slides):
                    html_file = os.path.join(temp_dir, f"slide_{i}.html")
                    with open(html_file, 'w', encoding='utf-8') as f:
                        f.write(slide['html_content'])
                    html_files.append(html_file)

                # 第3步：使用Playwright对每张幻灯片进行截图
                for i, html_file in enumerate(html_files):
                    screenshot_path = os.path.join(temp_dir, f"slide_{i}.png")

                    # 使用PDF converter的截图功能
                    success = await pdf_converter.screenshot_html(
                        html_file,
                        screenshot_path,
                        width=1280,
                        height=720
                    )

                    if success:
                        screenshot_paths.append(screenshot_path)
                        logging.info(f"Screenshot {i+1}/{len(html_files)} completed")
                    else:
                        logging.warning(f"Screenshot {i+1} failed, skipping")

                if len(screenshot_paths) == 0:
                    raise Exception("No screenshots were generated")

                # 第4步：将截图转换为PPTX
                logging.info("Creating PPTX from screenshots...")
                prs = Presentation()

                # 设置幻灯片尺寸为16:9
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(5.625)

                for i, screenshot_path in enumerate(screenshot_paths):
                    # 添加空白幻灯片
                    blank_slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_slide_layout)

                    # 添加截图，填充整个幻灯片
                    left = Inches(0)
                    top = Inches(0)
                    width = prs.slide_width
                    height = prs.slide_height

                    slide.shapes.add_picture(screenshot_path, left, top, width=width, height=height)

                    # 如果该幻灯片有演讲稿，添加到备注中
                    if i in speech_scripts:
                        notes_slide = slide.notes_slide
                        text_frame = notes_slide.notes_text_frame
                        text_frame.text = speech_scripts[i]
                        logging.info(f"Added speech script to slide {i+1} notes")

                # 保存PPTX文件
                prs.save(temp_pptx_path)
                logging.info(f"PPTX saved to {temp_pptx_path}")

                return {
                    "success": True,
                    "pptx_path": temp_pptx_path
                }

            except Exception as e:
                logging.error(f"HTML to PPTX conversion failed: {e}")
                import traceback
                traceback.print_exc()
                return {
                    "success": False,
                    "error": str(e)
                }
            finally:
                # 清理临时HTML和截图文件
                try:
                    import shutil
                    if os.path.exists(temp_dir):
                        shutil.rmtree(temp_dir)
                        logging.info(f"Cleaned up temp directory: {temp_dir}")
                except Exception as cleanup_error:
                    logging.warning(f"Failed to cleanup temp directory: {cleanup_error}")

        # 提交后台任务
        task_id = task_manager.submit_task(
            task_type="html_to_pptx_screenshot",
            func=html_to_pptx_task,
            metadata={
                "project_id": project_id,
                "project_topic": project.topic,
                "slide_count": len(slides),
                "pptx_path": temp_pptx_path
            }
        )

        # 立即返回任务ID
        return JSONResponse({
            "status": "processing",
            "task_id": task_id,
            "message": "PPTX generation with screenshots started in background",
            "polling_endpoint": f"/api/landppt/tasks/{task_id}"
        })

    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"PPTX screenshot export error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


# 后台任务查询端点
@router.get("/api/landppt/tasks/{task_id}")
async def get_task_status(task_id: str):
    """查询后台任务状态"""
    from ..services.background_tasks import get_task_manager

    task_manager = get_task_manager()
    task = task_manager.get_task(task_id)

    if not task:
        raise HTTPException(status_code=404, detail="Task not found")

    response = {
        "task_id": task.task_id,
        "task_type": task.task_type,
        "status": task.status.value,
        "progress": task.progress,
        "created_at": task.created_at.isoformat(),
        "updated_at": task.updated_at.isoformat(),
        "metadata": task.metadata
    }

    # 如果任务完成，添加结果信息
    if task.status.value == "completed" and task.result:
        response["result"] = task.result
        # 如果是PDF转PPTX任务，提供下载链接
        if task.task_type == "pdf_to_pptx_conversion" and task.result.get("success"):
            response["download_url"] = f"/api/landppt/tasks/{task_id}/download"

    # 如果任务失败，添加错误信息
    if task.status.value == "failed":
        response["error"] = task.error

    return JSONResponse(response)


@router.get("/api/landppt/tasks/{task_id}/download")
async def download_task_result(task_id: str):
    """下载任务结果文件"""
    from ..services.background_tasks import get_task_manager, TaskStatus
    from starlette.background import BackgroundTask

    task_manager = get_task_manager()
    task = task_manager.get_task(task_id)

    if not task:
        raise HTTPException(status_code=404, detail="Task not found")

    if task.status != TaskStatus.COMPLETED:
        raise HTTPException(status_code=400, detail=f"Task not completed yet (status: {task.status.value})")

    if not task.result or not task.result.get("success"):
        raise HTTPException(status_code=400, detail="Task failed or no result available")

    pptx_path = task.result.get("pptx_path")
    pdf_path = task.result.get("pdf_path")

    if not pptx_path or not os.path.exists(pptx_path):
        raise HTTPException(status_code=404, detail="Result file not found")

    # 获取项目主题作为文件名
    project_topic = task.metadata.get("project_topic", "PPT")
    safe_filename = urllib.parse.quote(f"{project_topic}_PPT.pptx", safe='')

    # 清理临时文件的后台任务
    def cleanup_temp_files():
        try:
            if pdf_path and os.path.exists(pdf_path):
                os.unlink(pdf_path)
        except:
            pass
        try:
            if pptx_path and os.path.exists(pptx_path):
                os.unlink(pptx_path)
        except:
            pass

    return FileResponse(
        pptx_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f"attachment; filename*=UTF-8''{safe_filename}",
            "X-Conversion-Method": "PDF-to-PPTX-Background"
        },
        background=BackgroundTask(cleanup_temp_files)
    )


@router.get("/api/projects/{project_id}/export/html")
async def export_project_html(project_id: str):
    """Export project as HTML ZIP package with slideshow index"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if we have slides data
        if not project.slides_data or len(project.slides_data) == 0:
            raise HTTPException(status_code=400, detail="PPT not generated yet")

        # Create temporary directory and generate files in thread pool
        zip_content = await run_blocking_io(_generate_html_export_sync, project)

        # URL encode the filename to handle Chinese characters
        zip_filename = f"{project.topic}_PPT.zip"
        safe_filename = urllib.parse.quote(zip_filename, safe='')

        from fastapi.responses import Response
        return Response(
            content=zip_content,
            media_type="application/zip",
            headers={
                "Content-Disposition": f"attachment; filename*=UTF-8''{safe_filename}"
            }
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def _generate_html_export_sync(project) -> bytes:
    """同步生成HTML导出文件（在线程池中运行）"""
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)

        # Generate individual HTML files for each slide
        slide_files = []
        for i, slide in enumerate(project.slides_data):
            slide_filename = f"slide_{i+1}.html"
            slide_files.append(slide_filename)

            # Create complete HTML document for each slide
            slide_html = _generate_individual_slide_html_sync(slide, i+1, len(project.slides_data), project.topic)

            slide_path = temp_path / slide_filename
            with open(slide_path, 'w', encoding='utf-8') as f:
                f.write(slide_html)

        # Generate index.html slideshow page
        index_html = _generate_slideshow_index_sync(project, slide_files)
        index_path = temp_path / "index.html"
        with open(index_path, 'w', encoding='utf-8') as f:
            f.write(index_html)

        # Create ZIP file
        zip_filename = f"{project.topic}_PPT.zip"
        zip_path = temp_path / zip_filename

        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            # Add index.html
            zipf.write(index_path, "index.html")

            # Add all slide files
            for slide_file in slide_files:
                slide_path = temp_path / slide_file
                zipf.write(slide_path, slide_file)

        # Read ZIP file content
        with open(zip_path, 'rb') as f:
            return f.read()


def _generate_individual_slide_html_sync(slide, slide_number: int, total_slides: int, topic: str) -> str:
    """同步生成单个幻灯片HTML（在线程池中运行）"""
    slide_html = slide.get('html_content', '')
    slide_title = slide.get('title', f'第{slide_number}页')

    # Check if it's already a complete HTML document
    import re
    if slide_html.strip().lower().startswith('<!doctype') or slide_html.strip().lower().startswith('<html'):
        # It's a complete HTML document, enhance it with navigation
        return _enhance_complete_html_with_navigation(slide_html, slide_number, total_slides, topic, slide_title)
    else:
        # It's just content, wrap it in a complete structure
        slide_content = slide_html

    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{topic} - {slide_title}</title>
    <style>
        body {{
            margin: 0;
            padding: 0;
            font-family: 'Microsoft YaHei', 'PingFang SC', sans-serif;
            background: #f5f5f5;
            display: flex;
            align-items: center;
            justify-content: center;
            min-height: 100vh;
        }}
        .slide-container {{
            width: 90vw;
            height: 90vh;
            background: white;
            border-radius: 10px;
            box-shadow: 0 4px 20px rgba(0,0,0,0.1);
            overflow: hidden;
            position: relative;
        }}
        .slide-content {{
            width: 100%;
            height: 100%;
            padding: 20px;
            box-sizing: border-box;
        }}
        .slide-number {{
            position: absolute;
            bottom: 20px;
            right: 20px;
            background: rgba(0,0,0,0.7);
            color: white;
            padding: 5px 10px;
            border-radius: 5px;
            font-size: 14px;
        }}
    </style>
</head>
<body>
    <div class="slide-container">
        <div class="slide-content">
            {slide_content}
        </div>
        <div class="slide-number">{slide_number} / {total_slides}</div>
    </div>
</body>
</html>"""


def _generate_slideshow_index_sync(project, slide_files: list) -> str:
    """同步生成幻灯片索引页面（在线程池中运行）"""
    slides_list = ""
    for i, slide_file in enumerate(slide_files):
        slide = project.slides_data[i]
        slide_title = slide.get('title', f'第{i+1}页')
        slides_list += f"""
        <div class="slide-item" onclick="openSlide('{slide_file}')">
            <div class="slide-preview">
                <div class="slide-number">{i+1}</div>
                <div class="slide-title">{slide_title}</div>
            </div>
        </div>"""

    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{project.topic} - PPT放映</title>
    <style>
        body {{
            margin: 0;
            padding: 0;
            font-family: 'Microsoft YaHei', 'PingFang SC', sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
        }}
        .header {{
            text-align: center;
            padding: 40px 20px;
            color: white;
        }}
        .header h1 {{
            margin: 0;
            font-size: 2.5em;
            font-weight: 300;
        }}
        .slides-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fill, minmax(300px, 1fr));
            gap: 20px;
            padding: 20px;
        }}
        .slide-item {{
            background: white;
            border-radius: 10px;
            padding: 20px;
            cursor: pointer;
            transition: all 0.3s ease;
            box-shadow: 0 4px 15px rgba(0,0,0,0.1);
        }}
        .slide-item:hover {{
            transform: translateY(-5px);
            box-shadow: 0 8px 25px rgba(0,0,0,0.2);
        }}
        .slide-number {{
            background: #007bff;
            color: white;
            width: 40px;
            height: 40px;
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            margin: 0 auto 15px auto;
            font-weight: bold;
        }}
        .slide-title {{
            font-size: 1.1em;
            color: #333;
            margin: 0;
            text-align: center;
        }}
    </style>
</head>
<body>
    <div class="header">
        <h1>{project.topic}</h1>
        <p>PPT演示文稿 - 共{len(slide_files)}页</p>
    </div>
    <div class="slides-grid">
        {slides_list}
    </div>
    <script>
        function openSlide(slideFile) {{
            window.open(slideFile, '_blank');
        }}
    </script>
</body>
</html>"""


async def _generate_combined_html_for_export(project, export_type: str) -> str:
    """Generate combined HTML for export (PDF or HTML)"""
    try:
        if not project.slides_data:
            raise ValueError("No slides data available")

        # Create a combined HTML document with all slides
        html_parts = []

        # HTML document header
        html_parts.append(f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{project.topic} - PPT导出</title>
    <style>
        body {{
            margin: 0;
            padding: 0;
            font-family: 'Microsoft YaHei', 'PingFang SC', sans-serif;
            background: #f5f5f5;
        }}
        .slide-container {{
            width: 100vw;
            height: 100vh;
            display: flex;
            align-items: center;
            justify-content: center;
            page-break-after: always;
            background: white;
            position: relative;
        }}
        .slide-container:last-child {{
            page-break-after: avoid;
        }}
        .slide-frame {{
            width: 90vw;
            height: 90vh;
            border: none;
            border-radius: 10px;
            box-shadow: 0 4px 20px rgba(0,0,0,0.1);
        }}
        .slide-number {{
            position: absolute;
            bottom: 20px;
            right: 20px;
            background: rgba(0,0,0,0.7);
            color: white;
            padding: 5px 10px;
            border-radius: 5px;
            font-size: 14px;
        }}
        @media print {{
            .slide-container {{
                page-break-after: always;
                width: 100%;
                height: 100vh;
            }}
        }}
    </style>
</head>
<body>""")

        # Add each slide preserving original styles
        for i, slide in enumerate(project.slides_data):
            slide_html = slide.get('html_content', '')
            if slide_html:
                # Preserve complete HTML structure
                if slide_html.strip().lower().startswith('<!doctype') or slide_html.strip().lower().startswith('<html'):
                    # Extract styles from head and content from body
                    import re

                    # Extract CSS styles from head
                    style_matches = re.findall(r'<style[^>]*>(.*?)</style>', slide_html, re.DOTALL | re.IGNORECASE)
                    slide_styles = '\n'.join(style_matches)

                    # Extract body content
                    body_match = re.search(r'<body[^>]*>(.*?)</body>', slide_html, re.DOTALL | re.IGNORECASE)
                    if body_match:
                        slide_content = body_match.group(1)
                    else:
                        slide_content = slide_html
                else:
                    slide_styles = ""
                    slide_content = slide_html

                html_parts.append(f"""
    <div class="slide-container">
        <style>
            {slide_styles}
        </style>
        <div class="slide-frame">
            {slide_content}
        </div>
        <div class="slide-number">{i + 1} / {len(project.slides_data)}</div>
    </div>""")

        # Close HTML document
        html_parts.append("""
</body>
</html>""")

        return ''.join(html_parts)

    except Exception as e:
        # Fallback: return a simple error page
        return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <title>导出错误</title>
</head>
<body>
    <h1>导出失败</h1>
    <p>错误信息: {str(e)}</p>
    <p>请确保PPT已经生成完成后再尝试导出。</p>
</body>
</html>"""



# Legacy Node.js Puppeteer check function - no longer needed with Pyppeteer
# def _check_puppeteer_available() -> bool:
#     """Check if Node.js and Puppeteer are available"""
#     # This function is deprecated - we now use Pyppeteer (Python) instead
#     return False

async def _generate_individual_slide_html(slide, slide_number: int, total_slides: int, topic: str) -> str:
    """Generate complete HTML document for individual slide preserving original styles"""
    slide_html = slide.get('html_content', '')
    slide_title = slide.get('title', f'第{slide_number}页')

    # Check if it's already a complete HTML document
    import re
    if slide_html.strip().lower().startswith('<!doctype') or slide_html.strip().lower().startswith('<html'):
        # It's a complete HTML document, enhance it with navigation
        return _enhance_complete_html_with_navigation(slide_html, slide_number, total_slides, topic, slide_title)
    else:
        # It's just content, wrap it in a complete structure
        slide_content = slide_html

    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{topic} - {slide_title}</title>
    <style>
        body {{
            margin: 0;
            padding: 0;
            font-family: 'Microsoft YaHei', 'PingFang SC', sans-serif;
            background: #f5f5f5;
            display: flex;
            align-items: center;
            justify-content: center;
            min-height: 100vh;
        }}
        .slide-container {{
            width: 90vw;
            height: 90vh;
            background: white;
            border-radius: 10px;
            box-shadow: 0 4px 20px rgba(0,0,0,0.1);
            overflow: hidden;
            position: relative;
        }}
        .slide-content {{
            width: 100%;
            height: 100%;
            padding: 20px;
            box-sizing: border-box;
        }}
        .slide-number {{
            position: absolute;
            bottom: 20px;
            right: 20px;
            background: rgba(0,0,0,0.7);
            color: white;
            padding: 5px 10px;
            border-radius: 5px;
            font-size: 14px;
        }}
        .navigation {{
            position: fixed;
            bottom: 20px;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 10px;
            z-index: 1000;
        }}
        .nav-btn {{
            background: #007bff;
            color: white;
            border: none;
            padding: 10px 15px;
            border-radius: 5px;
            cursor: pointer;
            text-decoration: none;
            display: inline-block;
        }}
        .nav-btn:hover {{
            background: #0056b3;
        }}
        .nav-btn:disabled {{
            background: #ccc;
            cursor: not-allowed;
        }}
        .fullscreen-btn {{
            position: fixed;
            top: 20px;
            right: 20px;
            background: #28a745;
            color: white;
            border: none;
            padding: 10px;
            border-radius: 5px;
            cursor: pointer;
            z-index: 1000;
        }}
        .fullscreen-btn:hover {{
            background: #1e7e34;
        }}
    </style>
</head>
<body>
    <div class="slide-container">
        <div class="slide-content">
            {slide_content}
        </div>
        <div class="slide-number">{slide_number} / {total_slides}</div>
    </div>

    <div class="navigation">
        <a href="index.html" class="nav-btn">🏠 返回目录</a>
        {"" if slide_number <= 1 else f'<a href="slide_{slide_number-1}.html" class="nav-btn">‹ 上一页</a>'}
        {"" if slide_number >= total_slides else f'<a href="slide_{slide_number+1}.html" class="nav-btn">下一页 ›</a>'}
    </div>

    <button class="fullscreen-btn" onclick="toggleFullscreen()" title="全屏显示">
        📺
    </button>

    <script>
        function toggleFullscreen() {{
            if (!document.fullscreenElement) {{
                document.documentElement.requestFullscreen();
            }} else {{
                if (document.exitFullscreen) {{
                    document.exitFullscreen();
                }}
            }}
        }}

        // Keyboard navigation
        document.addEventListener('keydown', function(e) {{
            if (e.key === 'ArrowLeft' && {slide_number} > 1) {{
                window.location.href = 'slide_{slide_number-1}.html';
            }} else if (e.key === 'ArrowRight' && {slide_number} < {total_slides}) {{
                window.location.href = 'slide_{slide_number+1}.html';
            }} else if (e.key === 'Escape') {{
                window.location.href = 'index.html';
            }}
        }});
    </script>
</body>
</html>"""

def _enhance_complete_html_with_navigation(original_html: str, slide_number: int, total_slides: int, topic: str, slide_title: str) -> str:
    """Enhance complete HTML document with navigation controls"""
    import re

    # Add navigation CSS and JavaScript to the head section
    navigation_css = """
    <style>
        .slide-navigation {
            position: fixed;
            bottom: 20px;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 10px;
            z-index: 10000;
            background: rgba(0,0,0,0.8);
            padding: 10px;
            border-radius: 25px;
        }
        .nav-btn {
            background: #007bff;
            color: white;
            border: none;
            padding: 10px 15px;
            border-radius: 5px;
            cursor: pointer;
            text-decoration: none;
            display: inline-block;
            font-size: 14px;
        }
        .nav-btn:hover {
            background: #0056b3;
        }
        .fullscreen-btn {
            position: fixed;
            top: 20px;
            right: 20px;
            background: #28a745;
            color: white;
            border: none;
            padding: 10px;
            border-radius: 5px;
            cursor: pointer;
            z-index: 10000;
            font-size: 16px;
        }
        .fullscreen-btn:hover {
            background: #1e7e34;
        }
    </style>"""

    navigation_js = f"""
    <script>
        function toggleFullscreen() {{
            if (!document.fullscreenElement) {{
                document.documentElement.requestFullscreen();
            }} else {{
                if (document.exitFullscreen) {{
                    document.exitFullscreen();
                }}
            }}
        }}

        // Keyboard navigation
        document.addEventListener('keydown', function(e) {{
            if (e.key === 'ArrowLeft' && {slide_number} > 1) {{
                window.location.href = 'slide_{slide_number-1}.html';
            }} else if (e.key === 'ArrowRight' && {slide_number} < {total_slides}) {{
                window.location.href = 'slide_{slide_number+1}.html';
            }} else if (e.key === 'Escape') {{
                window.location.href = 'index.html';
            }}
        }});
    </script>"""

    navigation_html = f"""
    <div class="slide-navigation">
        <a href="index.html" class="nav-btn">🏠 返回目录</a>
        {"" if slide_number <= 1 else f'<a href="slide_{slide_number-1}.html" class="nav-btn">‹ 上一页</a>'}
        {"" if slide_number >= total_slides else f'<a href="slide_{slide_number+1}.html" class="nav-btn">下一页 ›</a>'}
    </div>

    <button class="fullscreen-btn" onclick="toggleFullscreen()" title="全屏显示">
        📺
    </button>"""

    # Insert navigation CSS into head
    head_pattern = r'</head>'
    enhanced_html = re.sub(head_pattern, navigation_css + '\n</head>', original_html, flags=re.IGNORECASE)

    # Insert navigation HTML and JS before closing body tag
    body_pattern = r'</body>'
    enhanced_html = re.sub(body_pattern, navigation_html + '\n' + navigation_js + '\n</body>', enhanced_html, flags=re.IGNORECASE)

    return enhanced_html

async def _generate_pdf_slide_html(slide, slide_number: int, total_slides: int, topic: str) -> str:
    """Generate PDF-optimized HTML for individual slide without navigation elements"""
    slide_html = slide.get('html_content', '')
    slide_title = slide.get('title', f'第{slide_number}页')

    # Check if it's already a complete HTML document
    import re
    if slide_html.strip().lower().startswith('<!doctype') or slide_html.strip().lower().startswith('<html'):
        # It's a complete HTML document, clean it for PDF
        return _clean_html_for_pdf(slide_html, slide_number, total_slides)
    else:
        # It's just content, wrap it in a PDF-optimized structure
        slide_content = slide_html

    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{topic} - {slide_title}</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}

        html, body {{
            width: 100%;
            height: 100vh;
            margin: 0;
            padding: 0;
            font-family: 'Microsoft YaHei', 'PingFang SC', sans-serif;
            overflow: hidden;
        }}

        .slide-container {{
            width: 100vw;
            height: 100vh;
            display: flex;
            align-items: center;
            justify-content: center;
            position: relative;
        }}

        .slide-content {{
            width: 100%;
            height: 100%;
            display: flex;
            align-items: center;
            justify-content: center;
            position: relative;
        }}

        /* Ensure all backgrounds and colors are preserved for PDF */
        * {{
            -webkit-print-color-adjust: exact !important;
            print-color-adjust: exact !important;
        }}
    </style>
</head>
<body>
    <div class="slide-container">
        <div class="slide-content">
            {slide_content}
        </div>
    </div>
</body>
</html>"""

def _clean_html_for_pdf(original_html: str, slide_number: int, total_slides: int) -> str:
    """Clean complete HTML document for PDF generation by removing navigation elements"""
    import re

    # Remove navigation elements that might interfere with PDF generation
    cleaned_html = original_html

    # Remove navigation divs and buttons
    cleaned_html = re.sub(r'<div[^>]*class="[^"]*navigation[^"]*"[^>]*>.*?</div>', '', cleaned_html, flags=re.DOTALL | re.IGNORECASE)
    cleaned_html = re.sub(r'<button[^>]*class="[^"]*nav[^"]*"[^>]*>.*?</button>', '', cleaned_html, flags=re.DOTALL | re.IGNORECASE)
    cleaned_html = re.sub(r'<a[^>]*class="[^"]*nav[^"]*"[^>]*>.*?</a>', '', cleaned_html, flags=re.DOTALL | re.IGNORECASE)

    # Remove fullscreen buttons
    cleaned_html = re.sub(r'<button[^>]*fullscreen[^>]*>.*?</button>', '', cleaned_html, flags=re.DOTALL | re.IGNORECASE)

    # Add PDF-specific styles
    pdf_styles = """
    <style>
        /* PDF optimization styles */
        * {
            -webkit-print-color-adjust: exact !important;
            print-color-adjust: exact !important;
        }

        html, body {
            width: 100% !important;
            height: 100vh !important;
            margin: 0 !important;
            padding: 0 !important;
            overflow: hidden !important;
        }

        /* Hide any remaining navigation elements */
        .navigation, .nav-btn, .fullscreen-btn, .slide-navigation {
            display: none !important;
        }
    </style>
    """

    # Insert PDF styles before closing head tag
    head_pattern = r'</head>'
    cleaned_html = re.sub(head_pattern, pdf_styles + '\n</head>', cleaned_html, flags=re.IGNORECASE)

    return cleaned_html

async def _generate_pdf_with_pyppeteer(project, output_path: str, individual: bool = False) -> bool:
    """Generate PDF using Pyppeteer (Python)"""
    try:
        pdf_converter = get_pdf_converter()

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)

            # Always generate individual HTML files for each slide for better page separation
            # This ensures each slide becomes a separate PDF page
            html_files = []
            for i, slide in enumerate(project.slides_data):
                # Use a specialized PDF-optimized HTML generator without navigation
                slide_html = await _generate_pdf_slide_html(
                    slide, i+1, len(project.slides_data), project.topic
                )

                html_file = temp_path / f"slide_{i+1}.html"
                # Write HTML file in thread pool to avoid blocking
                def write_html_file(content, path):
                    with open(path, 'w', encoding='utf-8') as f:
                        f.write(content)

                await run_blocking_io(write_html_file, slide_html, str(html_file))
                html_files.append(str(html_file))

            # Use Pyppeteer to convert multiple files and merge them
            pdf_dir = temp_path / "pdfs"
            await run_blocking_io(pdf_dir.mkdir)

            logging.info(f"Starting PDF generation for {len(html_files)} files")

            # Convert HTML files to PDFs and merge them
            pdf_files = await pdf_converter.convert_multiple_html_to_pdf(
                html_files, str(pdf_dir), output_path
            )

            if pdf_files and os.path.exists(output_path):
                logging.info("Pyppeteer PDF generation successful")
                return True
            else:
                logging.error("Pyppeteer PDF generation failed: No output file created")
                return False

    except Exception as e:
        logging.error(f"Pyppeteer PDF generation failed: {e}")
        return False



async def _generate_combined_html_for_pdf(project) -> str:
    """Generate combined HTML for PDF export with all slides preserving original styles"""
    slides_html = ""
    global_styles = ""

    for i, slide in enumerate(project.slides_data):
        slide_html = slide.get('html_content', '')
        slide_title = slide.get('title', f'第{i+1}页')

        # Enhanced style extraction to preserve all styling
        if slide_html.strip().lower().startswith('<!doctype') or slide_html.strip().lower().startswith('<html'):
            import re

            # Extract all CSS styles from head (including link tags and style tags)
            style_matches = re.findall(r'<style[^>]*>(.*?)</style>', slide_html, re.DOTALL | re.IGNORECASE)
            link_matches = re.findall(r'<link[^>]*rel=["\']stylesheet["\'][^>]*>', slide_html, re.IGNORECASE)

            slide_styles = '\n'.join(style_matches)
            slide_links = '\n'.join(link_matches)

            # Extract body content with preserved attributes
            body_match = re.search(r'<body([^>]*)>(.*?)</body>', slide_html, re.DOTALL | re.IGNORECASE)
            if body_match:
                body_attrs = body_match.group(1)
                slide_content = body_match.group(2)
                # Preserve body styles if any
                if 'style=' in body_attrs:
                    body_style_match = re.search(r'style=["\']([^"\']*)["\']', body_attrs)
                    if body_style_match:
                        slide_styles += f"\n.slide-content {{ {body_style_match.group(1)} }}"
            else:
                slide_content = slide_html
                slide_links = ""

            # Add to global styles to avoid duplication
            if slide_links and slide_links not in global_styles:
                global_styles += slide_links + "\n"
        else:
            slide_styles = ""
            slide_content = slide_html
            slide_links = ""

        # Create a separate page for each slide with proper page break
        slides_html += f"""
        <div class="slide-page" data-slide="{i+1}" style="page-break-before: always; page-break-after: always; page-break-inside: avoid;">
            <style>
                /* Slide {i+1} specific styles */
                .slide-page[data-slide="{i+1}"] .slide-content {{
                    /* Preserve original styling */
                }}
                {slide_styles}
            </style>
            <div class="slide-content">
                {slide_content}
            </div>
            <div class="slide-footer">
                <span class="slide-number">{i+1} / {len(project.slides_data)}</span>
                <span class="slide-title">{slide_title}</span>
            </div>
        </div>"""

    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{project.topic} - PDF导出</title>
    {global_styles}
    <style>
        /* Reset and base styles */
        * {{
            box-sizing: border-box;
        }}

        body {{
            margin: 0;
            padding: 0;
            font-family: 'Microsoft YaHei', 'PingFang SC', sans-serif;
            /* Don't force background color - let slides define their own */
        }}

        .slide-page {{
            width: 297mm;
            height: 167mm;
            margin: 0;
            padding: 0;
            page-break-before: always;
            page-break-after: always;
            page-break-inside: avoid;
            position: relative;
            aspect-ratio: 16/9;
            /* Don't force background - preserve original slide backgrounds */
            overflow: hidden;
            display: block;
            box-sizing: border-box;
        }}

        .slide-page:first-child {{
            page-break-before: avoid;
        }}

        .slide-page:last-child {{
            page-break-after: avoid;
        }}

        .slide-content {{
            width: 100%;
            height: calc(100% - 30px);
            position: relative;
            /* Preserve original content styling */
        }}

        .slide-footer {{
            position: absolute;
            bottom: 5mm;
            right: 10mm;
            font-size: 10px;
            color: rgba(255, 255, 255, 0.7);
            background: rgba(0, 0, 0, 0.3);
            padding: 2px 8px;
            border-radius: 3px;
            z-index: 1000;
        }}

        .slide-number {{
            font-weight: bold;
        }}

        .slide-title {{
            margin-left: 8px;
            opacity: 0.8;
        }}

        /* Print-specific styles */
        @media print {{
            @page {{
                size: 297mm 167mm;
                margin: 0;
            }}

            body {{
                -webkit-print-color-adjust: exact;
                print-color-adjust: exact;
                margin: 0;
                padding: 0;
            }}

            .slide-page {{
                page-break-before: always;
                page-break-after: always;
                page-break-inside: avoid;
                -webkit-print-color-adjust: exact;
                print-color-adjust: exact;
                width: 297mm;
                height: 167mm;
                margin: 0;
                padding: 0;
                display: block;
            }}

            .slide-page:first-child {{
                page-break-before: avoid;
            }}

            .slide-page:last-child {{
                page-break-after: avoid;
            }}
        }}

        /* Ensure all backgrounds and colors are preserved */
        * {{
            -webkit-print-color-adjust: exact !important;
            print-color-adjust: exact !important;
        }}
    </style>
</head>
<body>
    {slides_html}
</body>
</html>"""




async def _generate_slideshow_index(project, slide_files: list) -> str:
    """Generate slideshow index page"""
    slides_list = ""
    for i, slide_file in enumerate(slide_files):
        slide = project.slides_data[i]
        slide_title = slide.get('title', f'第{i+1}页')
        slides_list += f"""
        <div class="slide-item" onclick="openSlide('{slide_file}')">
            <div class="slide-preview">
                <div class="slide-number">{i+1}</div>
                <div class="slide-title">{slide_title}</div>
            </div>
        </div>"""

    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{project.topic} - PPT放映</title>
    <style>
        body {{
            margin: 0;
            padding: 0;
            font-family: 'Microsoft YaHei', 'PingFang SC', sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
        }}
        .header {{
            text-align: center;
            padding: 40px 20px;
            color: white;
        }}
        .header h1 {{
            margin: 0;
            font-size: 2.5em;
            font-weight: 300;
        }}
        .header p {{
            margin: 10px 0 0 0;
            font-size: 1.2em;
            opacity: 0.9;
        }}
        .container {{
            max-width: 1200px;
            margin: 0 auto;
            padding: 0 20px;
        }}
        .slides-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fill, minmax(300px, 1fr));
            gap: 20px;
            padding: 20px 0;
        }}
        .slide-item {{
            background: white;
            border-radius: 10px;
            padding: 20px;
            cursor: pointer;
            transition: all 0.3s ease;
            box-shadow: 0 4px 15px rgba(0,0,0,0.1);
        }}
        .slide-item:hover {{
            transform: translateY(-5px);
            box-shadow: 0 8px 25px rgba(0,0,0,0.2);
        }}
        .slide-preview {{
            text-align: center;
        }}
        .slide-number {{
            background: #007bff;
            color: white;
            width: 40px;
            height: 40px;
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            margin: 0 auto 15px auto;
            font-weight: bold;
        }}
        .slide-title {{
            font-size: 1.1em;
            color: #333;
            margin: 0;
        }}
        .controls {{
            text-align: center;
            padding: 40px 20px;
        }}
        .btn {{
            background: #28a745;
            color: white;
            border: none;
            padding: 15px 30px;
            border-radius: 25px;
            font-size: 1.1em;
            cursor: pointer;
            margin: 0 10px;
            transition: all 0.3s ease;
            text-decoration: none;
            display: inline-block;
        }}
        .btn:hover {{
            background: #1e7e34;
            transform: translateY(-2px);
        }}
        .btn-secondary {{
            background: #6c757d;
        }}
        .btn-secondary:hover {{
            background: #545b62;
        }}
        @media (max-width: 768px) {{
            .slides-grid {{
                grid-template-columns: repeat(auto-fill, minmax(250px, 1fr));
                gap: 15px;
            }}
            .header h1 {{
                font-size: 2em;
            }}
        }}
    </style>
</head>
<body>
    <div class="header">
        <h1>{project.topic}</h1>
        <p>PPT演示文稿 - 共{len(slide_files)}页</p>
    </div>

    <div class="container">
        <div class="controls">
            <button class="btn" onclick="startSlideshow()">🎬 开始放映</button>
            <button class="btn btn-secondary" onclick="downloadAll()">📦 下载所有文件</button>
        </div>

        <div class="slides-grid">
            {slides_list}
        </div>
    </div>

    <script>
        function openSlide(slideFile) {{
            window.open(slideFile, '_blank');
        }}

        function startSlideshow() {{
            window.open('slide_1.html', '_blank');
        }}

        function downloadAll() {{
            alert('所有文件已包含在此ZIP包中');
        }}

        // Keyboard shortcuts
        document.addEventListener('keydown', function(e) {{
            if (e.key === 'Enter' || e.key === ' ') {{
                startSlideshow();
            }}
        }});
    </script>
</body>
</html>"""

@router.get("/upload", response_class=HTMLResponse)
async def web_upload_page(
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """File upload page"""
    return templates.TemplateResponse("upload.html", {
        "request": request
    })

async def _process_uploaded_files_for_outline(
    file_uploads: List[UploadFile],
    topic: str,
    target_audience: str,
    page_count_mode: str,
    min_pages: int,
    max_pages: int,
    fixed_pages: int,
    ppt_style: str,
    custom_style_prompt: str,
    file_processing_mode: str,
    content_analysis_depth: str,
    requirements: str = None,
    enable_web_search: bool = False,  # 新增参数
    scenario: str = "general",  # 新增参数
    language: str = "zh"  # 新增参数
) -> Optional[Dict[str, Any]]:
    """处理上传的多个文件并生成PPT大纲，支持联网搜索集成"""
    try:
        from ..services.file_processor import FileProcessor
        file_processor = FileProcessor()

        # 过滤掉None值（如果没有文件上传）
        files = [f for f in file_uploads if f is not None]
        if not files:
            logger.error("No files provided")
            return None

        saved_file_paths = []
        all_processed_content = []

        try:
            # 处理每个文件
            for file_upload in files:
                # 验证文件
                is_valid, message = file_processor.validate_file(file_upload.filename, file_upload.size)
                if not is_valid:
                    logger.error(f"File validation failed for {file_upload.filename}: {message}")
                    continue

                # 读取文件内容并保存到项目文件目录
                content = await file_upload.read()
                # logger.info(f"文件内容: {content}")
                project_file_path = await run_blocking_io(
                    _save_project_file_sync, content, file_upload.filename
                )
                saved_file_paths.append(project_file_path)

                # 处理单个文件内容
                file_result = await file_processor.process_file(
                    project_file_path,
                    file_upload.filename,
                    file_processing_mode=file_processing_mode,
                )
                all_processed_content.append({
                    "filename": file_upload.filename,
                    "content": file_result.processed_content
                })
                logger.debug(f"文件处理内容: {file_result.processed_content}")
            if not all_processed_content:
                logger.error("No files were successfully processed")
                return None

            # 决定是否使用联网搜索并整合
            merged_file_path = None
            merged_filename = None

            if enable_web_search and topic and topic.strip():
                # 使用联网搜索并整合本地文件
                logger.info(f"启用联网搜索模式，主题: {topic}")

                # 构建上下文信息
                context = {
                    'scenario': scenario,
                    'target_audience': target_audience or '普通大众',
                    'requirements': requirements or '',
                    'ppt_style': ppt_style,
                    'description': f'文件数量: {len(files)}'
                }

                # 进行联网搜索并与文件整合
                merged_file_path = await ppt_service.conduct_research_and_merge_with_files(
                    topic=topic,
                    language=language,
                    file_paths=saved_file_paths,
                    context=context,
                    file_processing_mode=file_processing_mode,
                )

                merged_filename = f"merged_with_search_{len(files)}_files.md"
                logger.info(f"✅ 联网搜索和文件整合完成: {merged_file_path}")
            else:
                # 不使用联网搜索，仅合并所有文件内容
                merged_content = file_processor.merge_multiple_files_to_markdown(all_processed_content)

                # 创建临时合并文件
                import tempfile
                import os
                with tempfile.NamedTemporaryFile(mode='w', delete=False, suffix='.md', encoding='utf-8') as merged_file:
                    merged_file.write(merged_content)
                    merged_file_path = merged_file.name

                merged_filename = f"merged_content_{len(files)}_files.md"

            saved_file_paths.append(merged_file_path)

            # 创建文件大纲生成请求
            from ..api.models import FileOutlineGenerationRequest
            filenames_str = ", ".join([f.filename for f in files])
            merged_filename = f"merged_content_{len(files)}_files.md"
            outline_request = FileOutlineGenerationRequest(
                file_path=merged_file_path,
                filename=merged_filename,
                topic=topic if topic.strip() else None,
                scenario=scenario,
                requirements=requirements,
                target_audience=target_audience,
                language=language,
                page_count_mode=page_count_mode,
                min_pages=min_pages,
                max_pages=max_pages,
                fixed_pages=fixed_pages,
                ppt_style=ppt_style,
                custom_style_prompt=custom_style_prompt,
                file_processing_mode=file_processing_mode,
                content_analysis_depth=content_analysis_depth
            )

            # 使用enhanced_ppt_service生成大纲
            result = await ppt_service.generate_outline_from_file(outline_request)

            if result.success:
                logger.info(f"Successfully generated outline from {len(files)} files: {filenames_str}")
                # 在大纲中添加文件信息，用于重新生成
                outline_with_file_info = result.outline.copy()
                original_filenames = [f.filename for f in files]
                file_paths_without_merge = saved_file_paths[:-1]  # 排除临时合并文件
                uploaded_files_info = [
                    {'filename': name, 'file_path': path}
                    for name, path in zip(original_filenames, file_paths_without_merge)
                ]
                outline_with_file_info['file_info'] = {
                    'file_paths': file_paths_without_merge,
                    'merged_file_path': merged_file_path,
                    'merged_filename': merged_filename,
                    'filenames': original_filenames,
                    'files_count': len(files),
                    'processing_mode': file_processing_mode,
                    'analysis_depth': content_analysis_depth,
                    'file_path': merged_file_path,
                    'filename': merged_filename,
                    'uploaded_files': uploaded_files_info
                }
                return outline_with_file_info
            else:
                logger.error(f"Failed to generate outline from files: {result.error}")
                # 如果生成失败，清理文件
                for file_path in saved_file_paths:
                    await run_blocking_io(_cleanup_project_file_sync, file_path)
                return None

        except Exception as e:
            # 清理所有已保存的文件
            for file_path in saved_file_paths:
                try:
                    await run_blocking_io(_cleanup_project_file_sync, file_path)
                except:
                    pass
            raise e

    except Exception as e:
        logger.error(f"Error processing uploaded files for outline: {e}")
        return None


async def _process_uploaded_file_for_outline(
    file_upload: UploadFile,
    topic: str,
    target_audience: str,
    page_count_mode: str,
    min_pages: int,
    max_pages: int,
    fixed_pages: int,
    ppt_style: str,
    custom_style_prompt: str,
    file_processing_mode: str,
    content_analysis_depth: str,
    requirements: str = None
) -> Optional[Dict[str, Any]]:
    """处理上传的单个文件并生成PPT大纲（向后兼容）"""
    return await _process_uploaded_files_for_outline(
        [file_upload], topic, target_audience, page_count_mode, min_pages, max_pages,
        fixed_pages, ppt_style, custom_style_prompt, file_processing_mode,
        content_analysis_depth, requirements
    )


def _save_temp_file_sync(content: bytes, filename: str) -> str:
    """同步保存临时文件（在线程池中运行）"""
    import tempfile
    import os

    with tempfile.NamedTemporaryFile(
        delete=False,
        suffix=os.path.splitext(filename)[1]
    ) as temp_file:
        temp_file.write(content)
        return temp_file.name


def _save_project_file_sync(content: bytes, filename: str) -> str:
    """同步保存项目文件到永久位置（在线程池中运行）"""
    import os
    import time
    from pathlib import Path

    # 创建项目文件目录
    project_files_dir = Path("temp/project_files")
    project_files_dir.mkdir(parents=True, exist_ok=True)

    # 生成唯一文件名
    timestamp = int(time.time())
    file_ext = os.path.splitext(filename)[1]
    safe_filename = f"{timestamp}_{filename}"
    file_path = project_files_dir / safe_filename

    # 保存文件
    with open(file_path, 'wb') as f:
        f.write(content)

    return str(file_path)


def _cleanup_temp_file_sync(temp_file_path: str):
    """同步清理临时文件（在线程池中运行）"""
    import os
    if os.path.exists(temp_file_path):
        os.unlink(temp_file_path)


def _cleanup_project_file_sync(project_file_path: str):
    """同步清理项目文件（在线程池中运行）"""
    import os
    if os.path.exists(project_file_path):
        os.unlink(project_file_path)


@router.get("/global-master-templates", response_class=HTMLResponse)
async def global_master_templates_page(
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """Global master templates management page"""
    try:
        return templates.TemplateResponse("global_master_templates.html", {
            "request": request
        })
    except Exception as e:
        logger.error(f"Error loading global master templates page: {e}")
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })


@router.get("/image-gallery", response_class=HTMLResponse)
async def image_gallery_page(
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """本地图床管理页面"""
    try:
        return templates.TemplateResponse("image_gallery.html", {
            "request": request,
            "user": user
        })
    except Exception as e:
        logger.error(f"Error rendering image gallery page: {e}")
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })


@router.get("/image-generation-test", response_class=HTMLResponse)
async def image_generation_test_page(
    request: Request,
    user: User = Depends(get_current_user_required)
):
    """AI图片生成测试页面"""
    try:
        return templates.TemplateResponse("image_generation_test.html", {
            "request": request,
            "user": user
        })
    except Exception as e:
        logger.error(f"Error rendering image generation test page: {e}")
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })


@router.get("/projects/{project_id}/template-selection", response_class=HTMLResponse)
async def template_selection_page(
    request: Request,
    project_id: str,
    user: User = Depends(get_current_user_required)
):
    """Template selection page for PPT generation"""
    try:
        # Get project info
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        return templates.TemplateResponse("template_selection.html", {
            "request": request,
            "project_id": project_id,
            "project_topic": project.topic
        })
    except Exception as e:
        logger.error(f"Error loading template selection page: {e}")
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })


# 图像重新生成相关辅助函数
async def analyze_image_context(image_info: Dict[str, Any], slide_content: Dict[str, Any],
                               project_topic: str, project_scenario: str) -> Dict[str, Any]:
    """分析图像在幻灯片中的上下文"""
    return {
        "slide_title": slide_content.get("title", ""),
        "slide_content": slide_content.get("html_content", ""),
        "image_alt": image_info.get("alt", ""),
        "image_title": image_info.get("title", ""),
        "image_size": f"{image_info.get('width', 0)}x{image_info.get('height', 0)}",
        "image_position": image_info.get("position", {}),
        "project_topic": project_topic,
        "project_scenario": project_scenario,
        "image_purpose": determine_image_purpose(image_info, slide_content)
    }

def determine_image_purpose(image_info: Dict[str, Any], slide_content: Dict[str, Any]) -> str:
    """确定图像在幻灯片中的用途"""
    # 简单的启发式规则来确定图像用途
    width = image_info.get('width', 0)
    height = image_info.get('height', 0)
    alt_text = image_info.get('alt', '').lower()

    if width > 800 or height > 600:
        return "background"  # 大图像可能是背景
    elif 'icon' in alt_text or 'logo' in alt_text:
        return "icon"
    elif 'chart' in alt_text or 'graph' in alt_text:
        return "chart_support"
    elif width < 200 and height < 200:
        return "decoration"
    else:
        return "illustration"

# 图像重新生成相关辅助函数

def select_best_image_source(enabled_sources: List, image_config: Dict[str, Any], image_context: Dict[str, Any]):
    """智能选择最佳的图片来源"""
    from ..services.models.slide_image_info import ImageSource

    # 如果只有一个启用的来源，直接使用
    if len(enabled_sources) == 1:
        return enabled_sources[0]

    # 根据图像用途和配置智能选择
    image_purpose = image_context.get('image_purpose', 'illustration')

    # 优先级规则
    if image_purpose == 'background':
        # 背景图优先使用AI生成，其次网络搜索
        if ImageSource.AI_GENERATED in enabled_sources:
            return ImageSource.AI_GENERATED
        elif ImageSource.NETWORK in enabled_sources:
            return ImageSource.NETWORK
        elif ImageSource.LOCAL in enabled_sources:
            return ImageSource.LOCAL

    elif image_purpose == 'icon':
        # 图标优先使用本地，其次AI生成
        if ImageSource.LOCAL in enabled_sources:
            return ImageSource.LOCAL
        elif ImageSource.AI_GENERATED in enabled_sources:
            return ImageSource.AI_GENERATED
        elif ImageSource.NETWORK in enabled_sources:
            return ImageSource.NETWORK

    elif image_purpose in ['illustration', 'chart_support', 'decoration']:
        # 说明性图片优先使用网络搜索，其次AI生成
        if ImageSource.NETWORK in enabled_sources:
            return ImageSource.NETWORK
        elif ImageSource.AI_GENERATED in enabled_sources:
            return ImageSource.AI_GENERATED
        elif ImageSource.LOCAL in enabled_sources:
            return ImageSource.LOCAL

    # 默认优先级：AI生成 > 网络搜索 > 本地
    for source in [ImageSource.AI_GENERATED, ImageSource.NETWORK, ImageSource.LOCAL]:
        if source in enabled_sources:
            return source

    # 如果都没有，返回第一个可用的
    return enabled_sources[0] if enabled_sources else ImageSource.AI_GENERATED

# 注意：generate_image_prompt_for_replacement 函数已被PPTImageProcessor的标准流程替代
# 现在使用 PPTImageProcessor._ai_generate_image_prompt 方法来生成提示词

def replace_image_in_html(html_content: str, image_info: Dict[str, Any], new_image_url: str) -> str:
    """在HTML内容中替换指定的图像，支持img标签、背景图像和SVG，保持布局和样式"""
    try:
        from bs4 import BeautifulSoup
        import re

        soup = BeautifulSoup(html_content, 'html.parser')

        old_src = image_info.get('src', '')
        image_type = image_info.get('type', 'img')

        if not old_src:
            logger.warning("图像信息中没有src属性，无法替换")
            return html_content

        replacement_success = False

        if image_type == 'img':
            # 处理 <img> 标签
            replacement_success = replace_img_tag(soup, image_info, new_image_url, old_src)

        elif image_type == 'background':
            # 处理背景图像
            replacement_success = replace_background_image(soup, image_info, new_image_url, old_src)

        elif image_type == 'svg':
            # 处理SVG图像
            replacement_success = replace_svg_image(soup, image_info, new_image_url, old_src)

        if replacement_success:
            logger.info(f"成功替换{image_type}图像: {old_src} -> {new_image_url}")
            return str(soup)
        else:
            logger.warning(f"未找到匹配的{image_type}图像进行替换")
            return fallback_string_replacement(html_content, old_src, new_image_url)

    except Exception as e:
        logger.error(f"替换HTML中的图像失败: {e}")
        return fallback_string_replacement(html_content, image_info.get('src', ''), new_image_url)

def replace_img_tag(soup, image_info: Dict[str, Any], new_image_url: str, old_src: str) -> bool:
    """替换img标签"""
    img_elements = soup.find_all('img')

    for img in img_elements:
        img_src = img.get('src', '')

        # 比较图像源URL（处理相对路径和绝对路径）
        if (img_src == old_src or
            img_src.endswith(old_src.split('/')[-1]) or
            old_src.endswith(img_src.split('/')[-1])):

            # 替换图像URL
            img['src'] = new_image_url

            # 保持原有的重要属性
            preserved_attributes = ['class', 'style', 'width', 'height', 'id']
            for attr in preserved_attributes:
                if attr in image_info and image_info[attr]:
                    img[attr] = image_info[attr]

            # 更新或保持alt和title
            if image_info.get('alt'):
                img['alt'] = image_info['alt']
            if image_info.get('title'):
                img['title'] = image_info['title']

            # 确保图像加载错误时有后备处理
            if not img.get('onerror'):
                img['onerror'] = "this.style.display='none'"

            return True

    return False

def replace_background_image(soup, image_info: Dict[str, Any], new_image_url: str, old_src: str) -> bool:
    """替换CSS背景图像"""
    # 查找所有元素
    all_elements = soup.find_all()

    for element in all_elements:
        # 检查内联样式中的背景图像
        style = element.get('style', '')
        if 'background-image' in style and old_src in style:
            # 替换内联样式中的背景图像URL
            new_style = style.replace(old_src, new_image_url)
            element['style'] = new_style
            return True

        # 检查class属性，可能对应CSS规则中的背景图像
        class_names = element.get('class', [])
        if class_names and image_info.get('className'):
            # 如果class匹配，我们假设这是目标元素
            if any(cls in image_info.get('className', '') for cls in class_names):
                # 为元素添加内联背景图像样式
                current_style = element.get('style', '')
                if current_style and not current_style.endswith(';'):
                    current_style += ';'
                new_style = f"{current_style}background-image: url('{new_image_url}');"
                element['style'] = new_style
                return True

    return False

def replace_svg_image(soup, image_info: Dict[str, Any], new_image_url: str, old_src: str) -> bool:
    """替换SVG图像"""
    # 查找SVG元素
    svg_elements = soup.find_all('svg')

    for svg in svg_elements:
        # 如果SVG有src属性（虽然不常见）
        if svg.get('src') == old_src:
            svg['src'] = new_image_url
            return True

        # 检查SVG的内容或其他标识
        if image_info.get('outerHTML') and svg.get_text() in image_info.get('outerHTML', ''):
            # 对于内联SVG，我们可能需要替换整个元素
            # 这里简化处理，添加一个data属性来标记已替换
            svg['data-replaced-image'] = new_image_url
            return True

    return False

def fallback_string_replacement(html_content: str, old_src: str, new_image_url: str) -> str:
    """后备的字符串替换方案"""
    try:
        import re

        if old_src and old_src in html_content:
            # 尝试多种替换模式
            patterns = [
                # img标签的src属性
                (rf'(<img[^>]*src=")[^"]*({re.escape(old_src)}[^"]*")([^>]*>)', rf'\1{new_image_url}\3'),
                # CSS背景图像
                (rf'(background-image:\s*url\([\'"]?)[^\'")]*({re.escape(old_src)}[^\'")]*)', rf'\1{new_image_url}'),
                # 直接字符串替换
                (re.escape(old_src), new_image_url)
            ]

            for pattern, replacement in patterns:
                updated_html = re.sub(pattern, replacement, html_content, flags=re.IGNORECASE)
                if updated_html != html_content:
                    logger.info(f"使用后备方案成功替换图像: {old_src} -> {new_image_url}")
                    return updated_html

        return html_content

    except Exception as e:
        logger.error(f"后备替换方案也失败: {e}")
        return html_content
